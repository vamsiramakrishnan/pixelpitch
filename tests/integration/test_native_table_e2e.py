"""End-to-end: render an HTML deck with a real <table>, convert it, and
verify the produced PPTX contains a native table primitive (editable cells,
not floating textboxes or a raster crop)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from slidify.api import ConversionConfig, convert

FIXTURE = Path(__file__).resolve().parents[1] / "fixtures" / "table_deck.html"


@pytest.mark.asyncio
async def test_convert_html_table_to_native_pptx_table(tmp_path: Path):
    pptx = tmp_path / "table_deck.pptx"
    cfg = ConversionConfig(run_oracle=False, run_tier3=False)
    result = await convert(FIXTURE, pptx, cfg)
    assert pptx.exists() and pptx.stat().st_size > 0
    assert result.n_slides == 1

    prs = Presentation(str(pptx))
    s = prs.slides[0]
    tables = [shp for shp in s.shapes if shp.has_table]
    assert len(tables) == 1, (
        f"expected exactly one native table, got {len(tables)} "
        f"(shapes: {[shp.shape_type for shp in s.shapes]})"
    )
    t = tables[0].table
    # 1 header row + 4 body rows; 3 cols.
    assert len(t.rows) == 5
    assert len(t.columns) == 3
    # Header text round-trips.
    assert t.cell(0, 0).text_frame.text == "Quarter"
    assert t.cell(0, 2).text_frame.text == "YoY"
    # A body cell round-trips.
    assert t.cell(2, 1).text_frame.text == "$3.9M"
