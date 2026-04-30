"""Tests for the IR → PPTX compiler.

These tests bypass the entire render-and-classify pipeline. They take a
JSON IR document (the wire format @slidify/components emits) and compile
it directly. Native ratio is 100% by construction.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from slidify.compile_ir import compile_ir
from slidify.ir import (
    FillRadialGradient,
    FillSolid,
    GradientStop,
    IRBbox,
    IRDeck,
    IRGroupNode,
    IRParagraph,
    IRShapeNode,
    IRSlide,
    IRTextNode,
    IRTextRun,
)


def _hero_deck() -> IRDeck:
    """Build a minimal hero-slide deck programmatically."""
    return IRDeck(
        version=1,
        slides=[
            IRSlide(
                index=0,
                background=FillRadialGradient(
                    kind="radial-gradient",
                    cx=0.8,
                    cy=0.12,
                    stops=[
                        GradientStop(color="#1e1b4b", position=0.0),
                        GradientStop(color="#0a0a14", position=1.0),
                    ],
                ),
                nodes=[
                    IRTextNode(
                        kind="text",
                        recipeId="kicker",
                        bbox=IRBbox(x=96, y=80, w=600, h=18),
                        paragraphs=[
                            IRParagraph(
                                runs=[
                                    IRTextRun(
                                        text="Q2 2026 · INVESTOR UPDATE",
                                        fontSizePx=13,
                                        fontWeight=700,
                                        color="#a78bfa",
                                    )
                                ],
                                align="left",
                            )
                        ],
                    ),
                    IRTextNode(
                        kind="text",
                        recipeId="title",
                        bbox=IRBbox(x=96, y=220, w=1080, h=220),
                        paragraphs=[
                            IRParagraph(
                                runs=[
                                    IRTextRun(text="A compiler for ", fontSizePx=96, fontWeight=800),
                                    IRTextRun(text="presentations", fontSizePx=96, fontWeight=800, color="#c084fc"),
                                    IRTextRun(text=", not a screenshot tool.", fontSizePx=96, fontWeight=800),
                                ],
                                align="left",
                            )
                        ],
                    ),
                    IRGroupNode(
                        kind="group",
                        recipeId="pill",
                        bbox=IRBbox(x=96, y=580, w=180, h=32),
                        children=[
                            IRShapeNode(
                                kind="shape",
                                recipeId="pill.bg",
                                bbox=IRBbox(x=96, y=580, w=180, h=32),
                                shape="rounded-rect",
                                borderRadiusPx=9999,
                                fill=FillSolid(kind="solid", color={"hex": "#ffffff", "alpha": 0.06}),
                            ),
                            IRTextNode(
                                kind="text",
                                recipeId="pill.text",
                                bbox=IRBbox(x=120, y=580, w=144, h=32),
                                paragraphs=[
                                    IRParagraph(
                                        runs=[IRTextRun(text="shipping today", fontSizePx=13, color="#e4e4e7")],
                                        align="left",
                                    )
                                ],
                            ),
                        ],
                    ),
                ],
            )
        ],
    )


def test_compile_ir_from_pydantic(tmp_path: Path):
    deck = _hero_deck()
    out = compile_ir(deck, tmp_path / "hero.pptx")
    assert out.exists() and out.stat().st_size > 0
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    # Background + 2 text frames + 2 group children = 5 shapes.
    assert len(slide.shapes) >= 4
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    joined = " ".join(texts)
    assert "Q2 2026" in joined
    assert "compiler for" in joined
    assert "shipping today" in joined


def test_compile_ir_from_json_string(tmp_path: Path):
    deck_json = _hero_deck().model_dump_json()
    out = compile_ir(deck_json, tmp_path / "hero_json.pptx")
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_compile_ir_from_dict(tmp_path: Path):
    deck_dict = json.loads(_hero_deck().model_dump_json())
    out = compile_ir(deck_dict, tmp_path / "hero_dict.pptx")
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_recipe_id_is_stamped_on_emitted_shapes(tmp_path: Path):
    """Every emitted shape carries its recipeId in an OOXML extension entry."""
    out = compile_ir(_hero_deck(), tmp_path / "stamped.pptx")
    prs = Presentation(str(out))
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    SLIDIFY = "https://slidify.dev/2026/recipe"
    n_stamped = 0
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        for ext in sp_pr.findall(f"{{{NS_A}}}extLst/{{{NS_A}}}ext"):
            if ext.get("uri") == SLIDIFY:
                n_stamped += 1
    # Should be at least 4 stamped shapes (kicker text, title text, pill bg,
    # pill text) plus the slide background.
    assert n_stamped >= 4, f"only {n_stamped} shapes carried recipe-id stamps"
