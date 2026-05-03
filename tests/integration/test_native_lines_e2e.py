from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slidify.api import ConversionConfig, convert

LINE_FIXTURE = """<!doctype html>
<html>
<head>
<style>
html, body { margin:0; width:1280px; height:720px; background:#fff; }
.slide { position:relative; width:1280px; height:720px; }
hr {
  position:absolute;
  left:120px;
  top:140px;
  width:520px;
  height:0;
  margin:0;
  border:0;
  border-top:2px solid #111827;
}
.vertical {
  position:absolute;
  left:760px;
  top:180px;
  width:0;
  height:260px;
  border-left:3px solid #2563eb;
}
</style>
</head>
<body>
<div class="slide">
  <hr>
  <div class="vertical"></div>
</div>
</body>
</html>
"""


FAUX_TABLE_FIXTURE = """<!doctype html>
<html>
<head>
<style>
html, body { margin:0; width:1280px; height:720px; background:#111827; }
.slide { position:relative; width:1280px; height:720px; color:#f9fafb; font-family:Arial,sans-serif; }
.board { position:absolute; left:160px; top:120px; width:760px; }
.head, .row {
  display:grid;
  grid-template-columns:120px 1fr 140px;
  height:52px;
  align-items:center;
  border-bottom:1px solid rgba(249,250,251,0.8);
}
.head { border-bottom-width:2px; font-weight:700; }
.cell {
  height:100%;
  display:flex;
  align-items:center;
  padding:0 14px;
  border-right:1px solid rgba(249,250,251,0.65);
}
.cell:last-child { border-right:none; }
</style>
</head>
<body>
<div class="slide">
  <div class="board">
    <div class="head"><div class="cell">ID</div><div class="cell">Signal</div><div class="cell">State</div></div>
    <div class="row"><div class="cell">01</div><div class="cell">Overflow policy</div><div class="cell">Ready</div></div>
    <div class="row"><div class="cell">02</div><div class="cell">Font fallback</div><div class="cell">Open</div></div>
  </div>
</div>
</body>
</html>
"""


@pytest.mark.asyncio
async def test_zero_thickness_html_rules_emit_as_native_lines(tmp_path):
    pptx = tmp_path / "lines.pptx"
    cfg = ConversionConfig(run_oracle=False, run_tier3=False)
    await convert(LINE_FIXTURE, pptx, cfg)

    prs = Presentation(str(pptx))
    lines = [
        shape
        for shape in prs.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    ]
    assert len(lines) >= 2


@pytest.mark.asyncio
async def test_faux_table_side_borders_emit_as_native_lines(tmp_path):
    pptx = tmp_path / "faux-table.pptx"
    cfg = ConversionConfig(run_oracle=False, run_tier3=False)
    await convert(FAUX_TABLE_FIXTURE, pptx, cfg)

    prs = Presentation(str(pptx))
    lines = [
        shape
        for shape in prs.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    ]
    assert len(lines) >= 6


@pytest.mark.asyncio
async def test_dotmatrix_ladder_corpus_slide_preserves_table_rules(tmp_path):
    source = Path("_bench/corpus/slide-21-dotmatrix-ladder.html")
    pptx = tmp_path / "dotmatrix-ladder.pptx"
    cfg = ConversionConfig(run_oracle=False, run_tier3=False)
    await convert(source, pptx, cfg)

    prs = Presentation(str(pptx))
    lines = [
        shape
        for shape in prs.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    ]
    assert len(lines) >= 40
