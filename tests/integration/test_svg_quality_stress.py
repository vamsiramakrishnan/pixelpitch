"""End-to-end integration test for SVG quality features.

Exercises the three improvements landed under the "slideshow quality"
arc:
  1. SVG path arcs (A/a) translate to cubic Beziers — the donut wedges
     in the fixture use M+L+A+Z, so they only emit non-trivially if the
     arc handler does real work.
  2. viewBox + preserveAspectRatio — a checkmark icon defined in a
     `viewBox="0 0 24 24"` but rendered into an 80×80 box should land
     centered with circle radius ~33 px (instead of radius 10 in the
     wrong corner).
  3. Path-count threshold (≥ 200) — the dense bar chart contains 60+
     SVG primitives in one element, which previously rasterized whole.

The test inspects the produced PPTX to confirm:
  - The slide produces native-translated content (native_area_ratio == 1.0).
  - At least 4 custGeom freeforms appear (the donut wedges).
  - At least 60 native shapes total exist (the dense chart didn't raster).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from slidify.api import ConversionConfig, convert

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
FIXTURE = Path(__file__).resolve().parents[1] / "fixtures" / "svg_quality_stress.html"


def _count_custgeom(slide) -> int:
    """Number of shapes with `<a:custGeom>` — i.e., paths emitted natively."""
    n = 0
    for shape in slide.shapes:
        sp_pr = getattr(shape._element, "spPr", None)
        if sp_pr is None:
            continue
        if sp_pr.find(f"{{{A_NS}}}custGeom") is not None:
            n += 1
    return n


@pytest.mark.asyncio
async def test_svg_quality_fixture_emits_natively(tmp_path: Path):
    pptx_path = tmp_path / "stress.pptx"
    cfg = ConversionConfig(run_oracle=False, run_tier3=False)
    result = await convert(FIXTURE.read_text(encoding="utf-8"), pptx_path, cfg)

    assert pptx_path.exists() and pptx_path.stat().st_size > 0
    assert result.n_slides == 1
    # All three SVGs should land on the native path; nothing falls back to
    # raster. Pre-changes, the dense bar chart alone (>30 prims) forced
    # whole-region raster.
    assert result.native_area_ratio == 1.0, (
        f"expected fully-native emission, got {result.native_area_ratio}"
    )

    prs = Presentation(pptx_path)
    slide = prs.slides[0]

    # Donut: four <path> elements with arc commands → four custGeom shapes.
    # If A/a still degraded to lineTo, the wedges would emit but as straight
    # triangles; we'd still see custGeom, so we additionally check that the
    # cubicBezTo count on those shapes is high.
    custgeom_count = _count_custgeom(slide)
    assert custgeom_count >= 4, (
        f"expected ≥4 custGeom freeforms (donut wedges), got {custgeom_count}"
    )

    # Sum cubicBezTo elements across all custGeom shapes. Each donut wedge
    # contributes ≥1 cubic from the 90° arc; the inner hole circle adds
    # zero (it's a preset oval). So we expect ≥4 cubicBezTo elements
    # solely from the arc handler.
    cubic_count = 0
    for shape in slide.shapes:
        sp_pr = getattr(shape._element, "spPr", None)
        if sp_pr is None:
            continue
        cg = sp_pr.find(f"{{{A_NS}}}custGeom")
        if cg is None:
            continue
        cubic_count += len(cg.findall(f".//{{{A_NS}}}cubicBezTo"))
    assert cubic_count >= 4, (
        f"expected ≥4 cubicBezTo from donut arcs; got {cubic_count}. "
        "If 0, arcs are still falling through to lineTo."
    )

    # Dense chart: 60+ primitives must produce ≥60 native shapes total.
    # Prior to bumping SVG_NATIVE_PATH_BUDGET above 30, the entire chart
    # rasterized to a single picture and shape count plummeted.
    assert len(slide.shapes) >= 60, (
        f"expected ≥60 native shapes total (dense chart should not raster); "
        f"got {len(slide.shapes)}"
    )
