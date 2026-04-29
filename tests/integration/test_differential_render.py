"""Integration tests for the multi-pass differential render path."""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from PIL import Image

from slidify.api import ConversionConfig, convert
from slidify.renderer import Renderer

FIXTURE = Path(__file__).resolve().parents[1] / "fixtures" / "differential_demo.html"


@pytest.mark.asyncio
async def test_renderer_emits_no_text_png_when_differential_on():
    html = FIXTURE.read_text(encoding="utf-8")
    async with Renderer(differential=True) as r:
        rendered = await r.render(html)
    assert len(rendered.ground_truth_png) > 0
    assert len(rendered.no_text_png) > 0
    # The two should differ (text is in one, not the other).
    assert rendered.ground_truth_png != rendered.no_text_png


@pytest.mark.asyncio
async def test_no_text_png_is_empty_when_differential_off():
    html = FIXTURE.read_text(encoding="utf-8")
    async with Renderer(differential=False) as r:
        rendered = await r.render(html)
    assert rendered.no_text_png == b""


@pytest.mark.asyncio
async def test_differential_text_pass_has_fewer_dark_pixels():
    """The decoration-only pass should have measurably fewer dark pixels than
    the ground truth — the missing pixels are the text glyphs."""
    html = FIXTURE.read_text(encoding="utf-8")
    async with Renderer(differential=True) as r:
        rendered = await r.render(html)

    def n_white_ish(png_bytes: bytes) -> int:
        im = Image.open(io.BytesIO(png_bytes)).convert("L")
        # Count light pixels — text on dark bg means GT has FEWER light pixels
        # than the decoration-only render after text is stripped.
        return sum(1 for px in im.getdata() if px > 200)

    gt_light = n_white_ish(rendered.ground_truth_png)
    dec_light = n_white_ish(rendered.no_text_png)
    # Decoration-only pass should have at least 5,000 fewer "white text"
    # pixels because the title + paragraph are gone. (Title alone is much
    # more than 5k white pixels at 88px.)
    assert gt_light - dec_light >= 5_000, (
        f"differential pass didn't reduce text pixels enough "
        f"(gt={gt_light}, decoration={dec_light}, delta={gt_light - dec_light})"
    )


@pytest.mark.asyncio
async def test_convert_with_differential_render_succeeds(tmp_path: Path):
    """End-to-end: convert with differential_render=True. Output must contain
    a Hybrid raster shape (the url() bg) and at least one native text frame."""
    html = FIXTURE.read_text(encoding="utf-8")
    out = tmp_path / "differential.pptx"
    cfg = ConversionConfig(
        run_oracle=False, run_tier3=False, differential_render=True
    )
    result = await convert(html, out, cfg)
    assert out.exists() and out.stat().st_size > 0
    assert result.n_slides == 1
    # native_area_ratio is 1.0 because Hybrid counts as native.
    assert result.native_area_ratio > 0.5

    from pptx import Presentation

    prs = Presentation(str(out))
    slide = prs.slides[0]
    has_text = any(
        sh.has_text_frame and sh.text_frame.text.strip() for sh in slide.shapes
    )
    has_picture = any("Picture" in type(sh).__name__ for sh in slide.shapes)
    assert has_text, "no native text frame emitted"
    assert has_picture, "no Hybrid raster background emitted"
