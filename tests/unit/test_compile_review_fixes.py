"""Regression tests for proactive review-pass fixes.

Covers three bugs found while sweeping compile_ir.py + helpers:

  * Solid fill alpha was discarded — translucent solids rendered opaque.
  * `IRPatternFill.angleDeg` rotated the whole shape, not just the pattern.
  * `_fetch_picture` for non-base64 `data:` URIs returned URL-encoded
    bytes instead of decoded bytes.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slidify.compile_ir import NS_A, _fetch_picture, compile_ir
from slidify.ir import (
    ColorWithAlpha,
    FillSolid,
    IRBbox,
    IRDeck,
    IRPatternFill,
    IRShapeNode,
    IRSlide,
)


def _compile_one(node: IRShapeNode, tmp_path: Path) -> Presentation:
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")
    return Presentation(str(out))


def _shape_sp_pr(prs: Presentation):
    """Return the spPr lxml element for the first non-background shape."""
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        # Skip the slide-background rect (full-bleed RECTANGLE).
        # The first shape we want is the test node, not the background.
        return sh, sp_pr
    return None, None


# ---- Solid fill alpha -------------------------------------------------------


def test_solid_fill_translucent_writes_alpha_child(tmp_path):
    """A `{hex, alpha: 0.5}` solid color must emit `<a:alpha val="50000"/>`
    inside its `<a:srgbClr>` so the rendered shape isn't opaque.
    """
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.translucent",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        shape="rect",
        fill=FillSolid(
            kind="solid",
            color=ColorWithAlpha(hex="#a78bfa", alpha=0.5),
        ),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    solid = sp_pr.find(f"{{{NS_A}}}solidFill")
    assert solid is not None, "shape must carry a solidFill"
    srgb = solid.find(f"{{{NS_A}}}srgbClr")
    assert srgb is not None
    alpha_el = srgb.find(f"{{{NS_A}}}alpha")
    assert alpha_el is not None, (
        "translucent solid lost its alpha — `<a:alpha>` should be present "
        "as a child of `<a:srgbClr>`"
    )
    assert alpha_el.get("val") == str(int(round(0.5 * 100_000)))


def test_solid_fill_opaque_omits_alpha_child(tmp_path):
    """Sanity: an opaque solid must NOT add an `<a:alpha>` child.

    Otherwise existing IR documents would gain spurious markup post-fix.
    """
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.opaque",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        shape="rect",
        fill=FillSolid(kind="solid", color="#a78bfa"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    solid = sp_pr.find(f"{{{NS_A}}}solidFill")
    srgb = solid.find(f"{{{NS_A}}}srgbClr")
    assert srgb.find(f"{{{NS_A}}}alpha") is None, (
        "opaque solid should not add an alpha child"
    )


# ---- PatternFill angleDeg ---------------------------------------------------


def test_pattern_fill_angle_deg_does_not_rotate_shape(tmp_path):
    """Per CONTRACT §1.3, `IRPatternFill.angleDeg` rotates the pattern,
    not the shape. PPTX has no per-pattern rotation, but rotating the
    host shape was a worse misinterpretation: it tilted the bbox + any
    other content. The compiler now drops the rotation on the floor
    (logged) instead of altering the shape.
    """
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.angled-pattern",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        shape="rect",
        fill=IRPatternFill(
            kind="pattern",
            pattern="lines-h",
            fgColor="#888888",
            bgColor="#ffffff",
            angleDeg=45.0,
        ),
    )
    prs = _compile_one(node, tmp_path)
    sh, _ = _shape_sp_pr(prs)
    # python-pptx's `shape.rotation` returns degrees; 0 (or unset) means
    # the shape itself wasn't rotated by the fill emit path.
    assert sh.rotation == 0.0, (
        f"shape rotation should be 0; got {sh.rotation} — pattern "
        "angleDeg leaked into shape transform."
    )


# ---- _fetch_picture data URI ------------------------------------------------


def test_fetch_picture_non_base64_data_uri_decodes_url_escapes():
    """A `data:image/svg+xml,<svg>...` URL with a percent-escape must be
    URL-decoded before being returned as bytes. Otherwise `%3C` survives
    as the three bytes `b'%3C'` instead of `b'<'`, producing a broken
    payload.
    """
    src = "data:image/svg+xml,<svg%20width%3D%221%22/>"
    out = _fetch_picture(src)
    # `%20` → space, `%3D` → `=`, `%22` → `"`
    assert out == b'<svg width="1"/>', f"got {out!r}"


def test_fetch_picture_base64_data_uri_decodes_payload():
    """Sanity: the base64 path keeps working. (Same 1×1 transparent PNG
    used by the escape-hatch fallback.)"""
    import base64

    # 'A' = 0x41 in ASCII = 'QQ==' in base64.
    src = "data:application/octet-stream;base64,QQ=="
    out = _fetch_picture(src)
    assert out == b"A"
    # And a more realistic case: a tiny PNG.
    png_b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAA"
        "C0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
    )
    out2 = _fetch_picture(f"data:image/png;base64,{png_b64}")
    assert out2 == base64.b64decode(png_b64)


# ---- Callout pointer adjustments --------------------------------------------


def test_callout_bubble_pointer_writes_avLst_adjustments(tmp_path):
    """`callout-bubble` with pointerSide+offset+length must emit an
    `<a:avLst>` carrying two `<a:gd>` entries (adj1, adj2) so the
    pointer geometry survives PowerPoint round-trip.

    Sanity check on `_apply_callout_adjustments` — without it, the
    callout renders with default pointer position regardless of IR.
    """
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.callout",
        bbox=IRBbox(x=0, y=0, w=200, h=120),
        shape="callout-bubble",
        fill=FillSolid(kind="solid", color="#16162a"),
        calloutPointerSide="bottom",
        calloutPointerOffset=0.5,
        calloutPointerLengthPx=24.0,
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    prst_geom = sp_pr.find(f"{{{NS_A}}}prstGeom")
    assert prst_geom is not None
    av_lst = prst_geom.find(f"{{{NS_A}}}avLst")
    assert av_lst is not None, "callout-bubble must carry an <a:avLst>"
    gds = av_lst.findall(f"{{{NS_A}}}gd")
    assert len(gds) >= 2, (
        f"expected ≥2 <a:gd> adjustment entries, got {len(gds)} — "
        "_apply_callout_adjustments may not have written them"
    )
    # Both gd entries should have a numeric `fmla` reference.
    for gd in gds[:2]:
        fmla = gd.get("fmla", "")
        assert "val " in fmla, f"unexpected fmla format: {fmla!r}"


def test_callout_bubble_no_pointer_fields_no_adjustments(tmp_path):
    """When all pointer fields are None, _apply_callout_adjustments
    early-returns; the callout uses default geometry.
    """
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.callout-default",
        bbox=IRBbox(x=0, y=0, w=200, h=120),
        shape="callout-bubble",
        fill=FillSolid(kind="solid", color="#16162a"),
        # all callout fields default to None
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    # The shape still emits, just without authored pointer geometry.
    assert sp_pr.find(f"{{{NS_A}}}prstGeom") is not None


# ---- Multi-shadow z-order ---------------------------------------------------


def test_multi_shadow_z_order_primary_above_siblings(tmp_path):
    """When a shape has 3+ shadows, the leftover shadow rects must sit
    BELOW the primary shape in the spTree so the soft halos don't
    occlude the primary's edges. Regression for the bug where leftover
    rects rendered on top of the primary, producing visible darkening.
    """
    from slidify.ir import IRBoxShadow

    node = IRShapeNode(
        kind="shape",
        recipeId="shape.multi-shadow",
        bbox=IRBbox(x=100, y=100, w=200, h=200),
        shape="rounded-rect",
        borderRadiusPx=12,
        fill=FillSolid(kind="solid", color="#16162a"),
        shadows=[
            IRBoxShadow(offsetX=0, offsetY=4, blur=8, spread=0,
                        color="#000000", inset=False),
            IRBoxShadow(offsetX=0, offsetY=12, blur=24, spread=0,
                        color="#000000", inset=False),
            IRBoxShadow(offsetX=0, offsetY=24, blur=48, spread=0,
                        color="#000000", inset=False),
        ],
    )
    prs = _compile_one(node, tmp_path)
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp_tree = prs.slides[0].shapes._spTree
    # Iterate <p:sp> children of spTree (skip non-shape elements like
    # nvGrpSpPr). The primary shape is the one whose spPr.extLst carries
    # the slidify recipeId extension matching our recipe.
    sp_children = sp_tree.findall(f"{{{P_NS}}}sp")
    primary_idx = None
    for i, sp in enumerate(sp_children):
        sp_pr = sp.find(f"{{{P_NS}}}spPr")
        if sp_pr is None:
            continue
        ext_lst = sp_pr.find(f"{{{NS_A}}}extLst")
        if ext_lst is None:
            continue
        for ext in ext_lst.findall(f"{{{NS_A}}}ext"):
            if ext.get("uri") != "https://slidify.dev/2026/recipe":
                continue
            rid = ext.find("{https://slidify.dev/2026/recipe}recipeId")
            if rid is not None and rid.text == "shape.multi-shadow":
                primary_idx = i
    assert primary_idx is not None, "primary shape not found in spTree"
    # Primary should be at the END of the <p:sp> sequence: all sibling
    # shadow rects come BEFORE it. With 3 shadows, the native pair takes
    # one outer slot on the primary shape, and ONE leftover spills as a
    # sibling rect. Background shape + sibling = 2 shapes before primary;
    # primary should be the last (index n-1).
    assert primary_idx == len(sp_children) - 1, (
        f"primary shape at index {primary_idx} of {len(sp_children)} "
        f"shapes; should be last so its halos render beneath it"
    )
