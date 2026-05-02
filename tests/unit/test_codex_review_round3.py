"""Regression tests for the four issues raised in the codex review on PR #22.

Covers:
  1. IRClipPathPath fallback — pictures now actually clip via custGeom +
     blipFill (was: silently un-clipped); shapes carry an extension
     stamp + warning log when the path-clip is unimplemented.
  2. report-escape-clusters --threshold flag — was a no-op (hardcoded
     1% area-share); now reads countByIntent and compares against the
     count threshold.
  3. slidify harvest no-output JSON guard — was unreachable
     (`context.obj == {"emit_json": True}` never true); replaced with
     an explicit `--json` flag.
  4. Classifier regex — only matched double-quoted data-atom; single-
     quoted and unquoted forms now also recognized.
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from click.testing import CliRunner
from pptx import Presentation

from slidify.compile_ir import NS_A, EscapeMetering, compile_ir
from slidify.ir import (
    IRBbox,
    IRClipPathPath,
    IRDeck,
    IRPathCommand,
    IRPictureNode,
    IRShapeNode,
    IRSlide,
    FillSolid,
)


# Tiny PNG embedded as base64 — saves us a network fetch for the picture
# tests. 1x1 transparent.
_TINY_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAA"
    "C0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
)
_DATA_URI = f"data:image/png;base64,{_TINY_PNG_B64}"


def _shape_count(out_path: Path) -> int:
    """Count shapes on slide 0 (excluding the slide-background rect)."""
    prs = Presentation(str(out_path))
    n = 0
    for sh in prs.slides[0].shapes:
        # Skip the full-bleed slide-background rect.
        if (
            sh.width and sh.height
            and sh.width >= prs.slide_width
            and sh.height >= prs.slide_height
        ):
            continue
        n += 1
    return n


# ---- 1. IRClipPathPath fallback --------------------------------------------


def test_picture_with_path_clip_emits_custgeom(tmp_path):
    """A picture with `clipPath: { kind: path, commands: [...] }` must
    emit a freeform shape (custGeom) with a blipFill — not a plain,
    un-clipped picture.
    """
    diamond_commands = [
        IRPathCommand(op="M", x=50, y=0),
        IRPathCommand(op="L", x=100, y=50),
        IRPathCommand(op="L", x=50, y=100),
        IRPathCommand(op="L", x=0, y=50),
        IRPathCommand(op="Z"),
    ]
    node = IRPictureNode(
        kind="picture",
        recipeId="picture.diamond-clip",
        bbox=IRBbox(x=200, y=200, w=200, h=200),
        src=_DATA_URI,
        clipPath=IRClipPathPath(kind="path", commands=diamond_commands),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")

    prs = Presentation(str(out))
    # Find the freeform/custGeom shape carrying the recipeId stamp.
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp_tree = prs.slides[0].shapes._spTree
    custgeom_with_blipfill = 0
    for sp in sp_tree.findall(f"{{{P_NS}}}sp"):
        sp_pr = sp.find(f"{{{P_NS}}}spPr")
        if sp_pr is None:
            continue
        if (
            sp_pr.find(f"{{{NS_A}}}custGeom") is not None
            and sp_pr.find(f"{{{NS_A}}}blipFill") is not None
        ):
            custgeom_with_blipfill += 1
    assert custgeom_with_blipfill == 1, (
        f"expected exactly one custGeom-with-blipFill shape (the path-"
        f"clipped picture), found {custgeom_with_blipfill}"
    )


def test_shape_with_path_clip_stamps_fallback_extension(tmp_path):
    """Path-clip on a ShapeNode is intentionally not implemented yet
    (requires polygon intersection between shape geometry + clip path).
    Verify the compiler stamps a fallback extension on the emitted
    shape so the gap is visible to the harvester / reverse-import path
    rather than the deck silently rendering un-clipped.
    """
    diamond_commands = [
        IRPathCommand(op="M", x=50, y=0),
        IRPathCommand(op="L", x=100, y=50),
        IRPathCommand(op="L", x=50, y=100),
        IRPathCommand(op="L", x=0, y=50),
        IRPathCommand(op="Z"),
    ]
    node = IRShapeNode(
        kind="shape",
        recipeId="shape.diamond-clipped-rect",
        bbox=IRBbox(x=200, y=200, w=200, h=200),
        shape="rect",
        fill=FillSolid(kind="solid", color="#16162a"),
        clipPath=IRClipPathPath(kind="path", commands=diamond_commands),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")

    prs = Presentation(str(out))
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp_tree = prs.slides[0].shapes._spTree
    found_fallback_ext = False
    for sp in sp_tree.findall(f"{{{P_NS}}}sp"):
        sp_pr = sp.find(f"{{{P_NS}}}spPr")
        if sp_pr is None:
            continue
        ext_lst = sp_pr.find(f"{{{NS_A}}}extLst")
        if ext_lst is None:
            continue
        for ext in ext_lst.findall(f"{{{NS_A}}}ext"):
            if ext.get("uri") == "https://slidify.dev/2026/clip-fallback":
                found_fallback_ext = True
    assert found_fallback_ext, (
        "expected a clip-fallback extension stamp on the shape with "
        "an unimplemented path-clip"
    )


# ---- 2. EscapeMetering.to_report_dict carries countByIntent ----------------


def test_escape_metering_to_report_dict_includes_count_by_intent():
    """`countByIntent` must be present so `slidify report-escape-clusters
    --threshold` can apply count-based filtering instead of falling back
    to area share.
    """
    m = EscapeMetering()
    m.add_slide(1280, 720)
    m.add_escape("clip-corner", 1000.0)
    m.add_escape("clip-corner", 2000.0)
    m.add_escape("conic-gradient", 5000.0)
    d = m.to_report_dict()
    assert "countByIntent" in d, "missing countByIntent key"
    assert d["countByIntent"] == {"clip-corner": 2, "conic-gradient": 1}


def test_report_escape_clusters_threshold_uses_counts(tmp_path):
    """Build a report.json that has countByIntent, run
    `slidify report-escape-clusters --threshold 3`, and verify only
    intents with count ≥ 3 are flagged would-promote.
    """
    from slidify.cli import cli

    report = {
        "escapeRate": {
            "value": 0.05,
            "byIntent": {"common": 0.04, "rare": 0.01},
            "countByIntent": {"common": 5, "rare": 2},
            "atomCandidates": [],
        }
    }
    report_path = tmp_path / "r.json"
    report_path.write_text(json.dumps(report))

    runner = CliRunner()
    res = runner.invoke(
        cli, ["report-escape-clusters", str(report_path), "--threshold", "3"]
    )
    assert res.exit_code == 0, res.output
    # `common` (count=5) crosses the threshold; `rare` (count=2) does not.
    common_line = [ln for ln in res.output.splitlines() if "common" in ln]
    rare_line = [ln for ln in res.output.splitlines() if "rare" in ln]
    assert common_line and "would-promote" in common_line[0]
    assert rare_line and "would-promote" not in rare_line[0]


# ---- 3. slidify harvest --json flag works ----------------------------------


def test_cli_harvest_json_flag_prints_machine_readable_output(tmp_path, monkeypatch):
    """The `--json` flag should print a clusters.json payload to stdout
    so callers can pipe it. The previous unreachable guard was replaced
    with an explicit flag.
    """
    from slidify import harvester
    from slidify.cli import cli
    from slidify.harvester import (
        AtomCandidate,
        Cluster,
        ClusterExemplar,
        HarvestReport,
    )

    corpus = tmp_path / "corpus"
    corpus.mkdir()
    (corpus / "a.html").write_text(
        "<html><body><div class='tile'>x</div></body></html>",
        encoding="utf-8",
    )

    def _fake_aggregate(corpus_dir, **kw):
        return HarvestReport(
            timestamp="2026-05-02T00:00:00+00:00",
            corpus_dir=str(corpus_dir),
            decks_processed=1,
            total_unmatched=1,
            unique_signatures=1,
            clusters=[Cluster(
                id="auto-001",
                sig_hash="cafebabe",
                signature="div(bg)[][100x100]",
                instances=1,
                exemplars=[ClusterExemplar(
                    slide_ref="a.html#node-0",
                    bbox_w=100, bbox_h=100,
                    sample_text="x", sample_classes="tile",
                )],
                sample_classes=["tile"],
                sample_text=["x"],
                bbox_typical={"w_avg": 100, "h_avg": 100,
                              "w_min": 100, "w_max": 100,
                              "h_min": 100, "h_max": 100},
            )],
            candidates={"auto-001": AtomCandidate(
                candidate_atom_id="surf.tile",
                candidate_axis="surf",
                candidate_props={},
                confidence=0.5,
                reason="",
            )},
        )

    monkeypatch.setattr(harvester, "aggregate_corpus", _fake_aggregate)

    runner = CliRunner()
    res = runner.invoke(cli, ["harvest", str(corpus), "--json"])
    assert res.exit_code == 0, res.output
    # The output should contain a parseable JSON object somewhere.
    # The CLI also emits human-readable text, so we extract the JSON
    # block by finding the first `{` and parsing forward.
    out = res.output
    json_start = out.find("{")
    assert json_start >= 0, f"no JSON in output:\n{out}"
    # The clusters.json payload spans multiple lines; just ensure it
    # contains the cluster id we built.
    assert "auto-001" in out
    assert "harvest_run" in out


# ---- 4. Classifier regex handles single-quoted + unquoted data-atom -------


def _run_classifier(html: str) -> dict:
    """Run `python -m slidify.patterns --classify-stdin` with the given HTML."""
    proc = subprocess.run(
        [sys.executable, "-m", "slidify.patterns", "--classify-stdin"],
        input=html.encode("utf-8"),
        capture_output=True,
        timeout=10,
    )
    # Exit code 0 = matched; 1 = unmatched (still valid output).
    out = proc.stdout.decode("utf-8").strip()
    return json.loads(out)


def test_classifier_double_quoted_data_atom():
    """Sanity: the original double-quoted form still works."""
    res = _run_classifier('<div data-atom="bg.mesh">x</div>')
    assert res.get("atom_id") == "bg.mesh"


def test_classifier_single_quoted_data_atom():
    """Single quotes were silently treated as unmatched before the fix."""
    res = _run_classifier("<div data-atom='bg.mesh'>x</div>")
    assert res.get("atom_id") == "bg.mesh", (
        f"single-quoted data-atom should match; got {res!r}"
    )


def test_classifier_unquoted_data_atom():
    """HTML5 allows unquoted attribute values for tokens without
    whitespace; the classifier should accept them too.
    """
    res = _run_classifier("<div data-atom=bg.mesh>x</div>")
    assert res.get("atom_id") == "bg.mesh", (
        f"unquoted data-atom should match; got {res!r}"
    )


def test_classifier_no_data_atom_returns_null():
    """Sanity: HTML without any data-atom returns atom_id=None."""
    res = _run_classifier("<div class='card'>x</div>")
    assert res.get("atom_id") is None


# ---- Round-4 codex feedback ------------------------------------------------


def test_path_clip_picture_anchored_to_node_bbox(tmp_path):
    """The freeform shape must be sized to `node.bbox`, not to the clip
    path's tight bbox. Otherwise an inset / off-center / extends-beyond
    clip path squashes or shifts the picture into the clip's bounds.

    Regression for codex-r3176316127. Builds a picture at bbox (200,200,
    400,400) with a clip path that lives at (300,300,300,300) — entirely
    inside the picture's bbox. The freeform shape's xfrm must read off
    of the picture bbox, not the clip's tight bounds.
    """
    diamond = [
        IRPathCommand(op="M", x=350, y=300),
        IRPathCommand(op="L", x=600, y=450),
        IRPathCommand(op="L", x=350, y=600),
        IRPathCommand(op="L", x=300, y=450),
        IRPathCommand(op="Z"),
    ]
    node = IRPictureNode(
        kind="picture",
        recipeId="picture.inset-clip",
        bbox=IRBbox(x=200, y=200, w=400, h=400),
        src=_DATA_URI,
        clipPath=IRClipPathPath(kind="path", commands=diamond),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")

    prs = Presentation(str(out))
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    for sp in prs.slides[0].shapes._spTree.findall(f"{{{P_NS}}}sp"):
        sp_pr = sp.find(f"{{{P_NS}}}spPr")
        if sp_pr is None or sp_pr.find(f"{{{NS_A}}}custGeom") is None:
            continue
        if sp_pr.find(f"{{{NS_A}}}blipFill") is None:
            continue
        xfrm = sp_pr.find(f"{{{NS_A}}}xfrm")
        assert xfrm is not None, "freeform shape must carry an xfrm"
        off = xfrm.find(f"{{{NS_A}}}off")
        ext = xfrm.find(f"{{{NS_A}}}ext")
        assert off is not None and ext is not None
        # 1 px = 9525 EMU. Picture bbox is (200,200,400,400) px.
        # Expect off.x = 200 * 9525 = 1_905_000, off.y same, ext.cx = 400*9525.
        EMU = 9525
        x_emu = int(off.get("x"))
        y_emu = int(off.get("y"))
        cx_emu = int(ext.get("cx"))
        cy_emu = int(ext.get("cy"))
        # Allow ±1 EMU jitter from rounding.
        assert abs(x_emu - 200 * EMU) <= EMU, (
            f"shape origin x={x_emu} EMU; expected node.bbox.x ({200*EMU})"
        )
        assert abs(y_emu - 200 * EMU) <= EMU
        assert abs(cx_emu - 400 * EMU) <= EMU, (
            f"shape width={cx_emu}; expected node.bbox.w ({400*EMU})"
        )
        assert abs(cy_emu - 400 * EMU) <= EMU
        return  # found the right shape; success
    raise AssertionError("never found a custGeom+blipFill shape on the slide")


def test_harvest_json_produces_pure_stdout(tmp_path, monkeypatch):
    """`slidify harvest <corpus> --json` must emit JSON-only on stdout
    so `... --json | jq ...` works. Status / progress text is routed
    to stderr.

    Regression for codex-r3176316129.
    """
    from slidify import harvester
    from slidify.cli import cli
    from slidify.harvester import (
        AtomCandidate,
        Cluster,
        ClusterExemplar,
        HarvestReport,
    )

    corpus = tmp_path / "corpus"
    corpus.mkdir()
    (corpus / "a.html").write_text("<html><body><div class='tile'>x</div></body></html>",
                                   encoding="utf-8")

    def _fake(corpus_dir, **kw):
        return HarvestReport(
            timestamp="2026-05-02T00:00:00+00:00",
            corpus_dir=str(corpus_dir),
            decks_processed=1,
            total_unmatched=1,
            unique_signatures=1,
            clusters=[Cluster(
                id="auto-001",
                sig_hash="cafebabe",
                signature="div(bg)[][100x100]",
                instances=1,
                exemplars=[ClusterExemplar(
                    slide_ref="a.html#node-0",
                    bbox_w=100, bbox_h=100,
                    sample_text="x", sample_classes="tile",
                )],
                sample_classes=["tile"],
                sample_text=["x"],
                bbox_typical={"w_avg": 100, "h_avg": 100,
                              "w_min": 100, "w_max": 100,
                              "h_min": 100, "h_max": 100},
            )],
            candidates={"auto-001": AtomCandidate(
                candidate_atom_id="surf.tile",
                candidate_axis="surf",
                candidate_props={},
                confidence=0.5,
                reason="",
            )},
        )

    monkeypatch.setattr(harvester, "aggregate_corpus", _fake)

    runner = CliRunner()
    # Recent Click versions split stdout/stderr by default; result exposes
    # them on `.stdout` and `.stderr`.
    res = runner.invoke(cli, ["harvest", str(corpus), "--json"])
    assert res.exit_code == 0, res.output
    # stdout MUST be parseable JSON, with no human chatter mixed in.
    payload = json.loads(res.stdout)
    assert "harvest_run" in payload
    assert payload["clusters"][0]["id"] == "auto-001"
    # Status text should be visible on stderr (best-effort: at least the
    # "Harvesting …" header).
    assert "Harvesting" in res.stderr, (
        f"expected status text on stderr; got: {res.stderr!r}"
    )
