"""Tests for `IRPatternFill` → PPTX compilation (CONTRACT §3.2).

Per-pattern routing:
  * `lines-h` / `lines-v` / `lines-grid` / `diagonal` / `crosshatch` →
    `<a:pattFill prst="…">` on the host shape.
  * `dots` → host shape goes transparent + a tiled group of small ovals.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from slidify.compile_ir import NS_A, compile_ir
from slidify.ir import (
    IRBbox,
    IRDeck,
    IRPatternFill,
    IRShapeNode,
    IRSlide,
)


def _deck_with_pattern(pattern: str, **patt_kwargs) -> IRDeck:
    return IRDeck(
        version=2,
        slides=[
            IRSlide(
                index=0,
                nodes=[
                    IRShapeNode(
                        kind="shape",
                        recipeId=f"pattern.{pattern}",
                        bbox=IRBbox(x=20, y=20, w=200, h=200),
                        shape="rect",
                        fill=IRPatternFill(
                            kind="pattern",
                            pattern=pattern,
                            fgColor="#0000ff",
                            bgColor="#ffffff",
                            **patt_kwargs,
                        ),
                    )
                ],
            )
        ],
    )


def _first_pattfill_prst(prs: Presentation) -> str | None:
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        patt = sp_pr.find(f"{{{NS_A}}}pattFill")
        if patt is not None:
            return patt.get("prst")
    return None


def test_pattern_lines_h_emits_lt_horz(tmp_path):
    out = compile_ir(_deck_with_pattern("lines-h"), tmp_path / "lh.pptx")
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "ltHorz"


def test_pattern_lines_v_emits_lt_vert(tmp_path):
    out = compile_ir(_deck_with_pattern("lines-v"), tmp_path / "lv.pptx")
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "ltVert"


def test_pattern_lines_grid_emits_sm_grid(tmp_path):
    out = compile_ir(_deck_with_pattern("lines-grid"), tmp_path / "lg.pptx")
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "smGrid"


def test_pattern_diagonal_emits_lt_up_diag(tmp_path):
    out = compile_ir(_deck_with_pattern("diagonal"), tmp_path / "diag.pptx")
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "ltUpDiag"


def test_pattern_crosshatch_emits_diag_cross(tmp_path):
    out = compile_ir(_deck_with_pattern("crosshatch"), tmp_path / "cross.pptx")
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "diagCross"


def test_pattern_heavy_feature_size_promotes_to_dark_variant(tmp_path):
    out = compile_ir(
        _deck_with_pattern("lines-h", featureSizePx=3.0),
        tmp_path / "lh-dk.pptx",
    )
    prs = Presentation(str(out))
    assert _first_pattfill_prst(prs) == "dkHorz"


def test_pattern_dots_emits_oval_grid(tmp_path):
    """`dots` → host shape transparent + N×M sibling ovals.

    With tile = 25×25 and bbox 200×200, expect 8×8 = 64 dots.
    """
    deck = _deck_with_pattern(
        "dots", tileWidthPx=25.0, tileHeightPx=25.0, featureSizePx=2.0
    )
    out = compile_ir(deck, tmp_path / "dots.pptx")
    prs = Presentation(str(out))
    # Count ovals (auto_shape_type == OVAL = 9)
    n_ovals = sum(
        1
        for sh in prs.slides[0].shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 9
    )
    assert n_ovals == 64, f"expected 64 dots, got {n_ovals}"


def test_pattern_dots_host_shape_has_no_fill(tmp_path):
    """The host shape under `dots` should be transparent (noFill)."""
    deck = _deck_with_pattern(
        "dots", tileWidthPx=50.0, tileHeightPx=50.0
    )
    out = compile_ir(deck, tmp_path / "dots-bg.pptx")
    prs = Presentation(str(out))
    # First non-oval shape is the host rect; check it has noFill.
    for sh in prs.slides[0].shapes:
        if (
            getattr(sh, "auto_shape_type", None) is not None
            and sh.auto_shape_type != 9
        ):
            sp_pr = sh._element.spPr
            assert sp_pr.find(f"{{{NS_A}}}noFill") is not None
            return
    pytest.fail("no host rectangle found")


def test_pattern_fill_round_trip_reopens(tmp_path):
    """python-pptx must re-open every pattern variant."""
    for variant in (
        "lines-h", "lines-v", "lines-grid", "diagonal", "crosshatch", "dots"
    ):
        deck = _deck_with_pattern(variant)
        out = compile_ir(deck, tmp_path / f"{variant}.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 1, variant


def test_pattern_fg_bg_colors_propagate(tmp_path):
    """fgClr + bgClr <a:srgbClr val=…> match IR colors."""
    deck = _deck_with_pattern("lines-h")
    out = compile_ir(deck, tmp_path / "colors.pptx")
    prs = Presentation(str(out))
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        patt = sp_pr.find(f"{{{NS_A}}}pattFill")
        if patt is None:
            continue
        fg_clr = patt.find(f"{{{NS_A}}}fgClr/{{{NS_A}}}srgbClr")
        bg_clr = patt.find(f"{{{NS_A}}}bgClr/{{{NS_A}}}srgbClr")
        assert fg_clr is not None and fg_clr.get("val") == "0000FF"
        assert bg_clr is not None and bg_clr.get("val") == "FFFFFF"
        return
    pytest.fail("no pattFill emitted")
