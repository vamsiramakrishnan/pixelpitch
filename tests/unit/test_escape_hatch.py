"""Tests for the chrome.escape-hatch first-class metered escape valve.

Covers:
  * EscapeMetering bookkeeping (area share + per-intent breakdown).
  * Compile path: a deck with one EscapeHatch raster produces escapeRate>0.
  * Per-intent attribution: escape area lands in the right bucket.
  * native_area_ratio exclusion: escape-hatch rasters don't drag down the
    denominator (per CONTRACT-v1 §9.5).
  * PPTX extension stamping: cssPayload + intent survive to disk so the
    harvester can find them.

These tests intentionally use a 1×1 transparent placeholder PNG (provided
by the IR) rather than invoking Chromium — they validate the metering and
metadata stamping, not the screenshot machinery.
"""

from __future__ import annotations

import base64
import json
from pathlib import Path

from pptx import Presentation

from slidify.compile_ir import (
    NS_A,
    SLIDIFY_ESCAPE_NS,
    SLIDIFY_RECIPE_NS,
    EscapeMetering,
    compile_ir,
)
from slidify.ir import (
    FillSolid,
    IRBbox,
    IRDeck,
    IRRasterNode,
    IRShapeNode,
    IRSlide,
)

# 1×1 transparent PNG — gives the test a valid raster payload without
# pulling Chromium into the unit-test path.
_TRANSPARENT_1PX_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
)


def _escape_node(
    *,
    bbox: IRBbox,
    intent: str = "non-rect-clip",
    attempted: str | None = "surf.glass",
    css_payload: str = "background: red; clip-path: polygon(0 0, 100% 0, 50% 100%);",
    body: str = "",
) -> IRRasterNode:
    """Build a `chrome.escape-hatch` IR raster matching the TSX emitter shape."""
    return IRRasterNode(
        kind="raster",
        recipeId="chrome.escape-hatch",
        bbox=bbox,
        zOrder=0,
        metadata={
            "role": "escape-hatch",
            "cssPayload": css_payload,
            "body": body,
            "intent": intent,
            "attempted": attempted,
            "excludeFromNativeRatio": True,
        },
        pngBase64=_TRANSPARENT_1PX_PNG_B64,
    )


def _slide_with(node: IRRasterNode | list[IRRasterNode], idx: int = 0) -> IRSlide:
    nodes = node if isinstance(node, list) else [node]
    return IRSlide(
        index=idx,
        bbox=IRBbox(x=0, y=0, w=1280, h=720),
        background=FillSolid(kind="solid", color="#ffffff"),
        nodes=list(nodes),
    )


# ---------------------------------------------------------------------------
# EscapeMetering unit tests
# ---------------------------------------------------------------------------


def test_metering_tracks_area_share() -> None:
    """The total escape-rate is escape_area / total_slide_area."""
    m = EscapeMetering()
    m.add_slide(1280, 720)  # 921_600 px²
    m.add_escape("non-rect-clip", 92_160)  # exactly 10% of the slide
    report = m.to_report_dict()
    assert abs(report["value"] - 0.10) < 1e-9
    assert abs(report["byIntent"]["non-rect-clip"] - 0.10) < 1e-9
    assert report["atomCandidates"] == []  # M4 fills this later


def test_metering_buckets_by_intent() -> None:
    """Multiple escapes with different intents land in distinct buckets."""
    m = EscapeMetering()
    m.add_slide(1000, 1000)  # 1_000_000
    m.add_escape("non-rect-clip", 100_000)
    m.add_escape("complex-text-warp", 200_000)
    m.add_escape("non-rect-clip", 50_000)  # accumulates into the same bucket
    report = m.to_report_dict()
    assert abs(report["value"] - 0.35) < 1e-9
    assert abs(report["byIntent"]["non-rect-clip"] - 0.15) < 1e-9
    assert abs(report["byIntent"]["complex-text-warp"] - 0.20) < 1e-9


def test_metering_zero_division_safe() -> None:
    """No slides → returns the empty shape, not a divide-by-zero."""
    m = EscapeMetering()
    m.add_escape("intent", 1000)
    report = m.to_report_dict()
    assert report == {"value": 0.0, "byIntent": {}, "atomCandidates": []}


# ---------------------------------------------------------------------------
# End-to-end: compile a deck with one escape hatch
# ---------------------------------------------------------------------------


def test_single_escape_hatch_produces_positive_escape_rate(tmp_path: Path) -> None:
    """Compile a deck with one EscapeHatch covering 1/4 of the slide and
    confirm the metering reports the expected area share."""
    # 640×360 escape on a 1280×720 slide = 25% of slide area.
    deck = IRDeck(
        version=1,
        slides=[
            _slide_with(
                _escape_node(bbox=IRBbox(x=0, y=0, w=640, h=360))
            )
        ],
    )
    out, metering = compile_ir(
        deck, tmp_path / "esc.pptx", return_metering=True
    )
    assert out.exists() and out.stat().st_size > 0

    report = metering.to_report_dict()
    assert report["value"] > 0
    assert abs(report["value"] - 0.25) < 1e-9
    assert "non-rect-clip" in report["byIntent"]
    assert abs(report["byIntent"]["non-rect-clip"] - 0.25) < 1e-9


def test_escape_area_attributed_to_intent(tmp_path: Path) -> None:
    """Two escape hatches with different intents → distinct buckets."""
    deck = IRDeck(
        version=1,
        slides=[
            _slide_with([
                _escape_node(
                    bbox=IRBbox(x=0, y=0, w=320, h=360),  # 320*360 = 115_200
                    intent="non-rect-clip",
                ),
                _escape_node(
                    bbox=IRBbox(x=320, y=0, w=320, h=360),
                    intent="complex-text-warp",
                ),
            ])
        ],
    )
    _, metering = compile_ir(
        deck, tmp_path / "two.pptx", return_metering=True
    )
    report = metering.to_report_dict()
    # Each escape is 115_200 / 921_600 = 0.125 of the slide.
    assert abs(report["byIntent"]["non-rect-clip"] - 0.125) < 1e-9
    assert abs(report["byIntent"]["complex-text-warp"] - 0.125) < 1e-9
    assert abs(report["value"] - 0.25) < 1e-9


def test_no_escape_hatch_produces_zero_rate(tmp_path: Path) -> None:
    """A plain shape-only deck has escapeRate.value == 0."""
    deck = IRDeck(
        version=1,
        slides=[
            IRSlide(
                index=0,
                nodes=[
                    IRShapeNode(
                        kind="shape",
                        recipeId="x",
                        bbox=IRBbox(x=0, y=0, w=100, h=100),
                        shape="rect",
                        fill=FillSolid(kind="solid", color="#ff0000"),
                    )
                ],
            )
        ],
    )
    _, metering = compile_ir(
        deck, tmp_path / "plain.pptx", return_metering=True
    )
    report = metering.to_report_dict()
    assert report["value"] == 0.0
    assert report["byIntent"] == {}


# ---------------------------------------------------------------------------
# native_area_ratio exclusion (CONTRACT-v1 §9.5)
# ---------------------------------------------------------------------------


def test_escape_hatch_marked_excluded_from_native_ratio(tmp_path: Path) -> None:
    """The IR metadata carries `excludeFromNativeRatio: True` — the flag the
    `native_area_ratio` computation reads to skip the raster's bbox.

    This test asserts the flag survives compile (it's preserved on the IR
    node and stamped into the PPTX extension payload) — not the ratio math
    itself, which lives in slidify/emitter.py and operates on different
    types entirely (EmitOps from the HTML→PPTX path, not IRNodes).
    """
    bbox = IRBbox(x=100, y=100, w=400, h=200)
    node = _escape_node(bbox=bbox)
    assert node.metadata["excludeFromNativeRatio"] is True

    # Compile and confirm the escape ext payload carries the flag's
    # justification (intent + cssPayload) so a downstream native-ratio
    # consumer can identify the picture as a deliberate native skip.
    deck = IRDeck(version=1, slides=[_slide_with(node)])
    out, _ = compile_ir(deck, tmp_path / "exclude.pptx", return_metering=True)

    prs = Presentation(str(out))
    found_payload = False
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        for ext in sp_pr.findall(f"{{{NS_A}}}extLst/{{{NS_A}}}ext"):
            if ext.get("uri") == SLIDIFY_ESCAPE_NS:
                payload_el = ext.find(f"{{{SLIDIFY_ESCAPE_NS}}}escapePayload")
                assert payload_el is not None and payload_el.text
                payload = json.loads(payload_el.text)
                assert payload["intent"] == "non-rect-clip"
                assert payload["attempted"] == "surf.glass"
                assert "clip-path" in payload["cssPayload"]
                found_payload = True
    assert found_payload, "escape-hatch payload extension not found in slide"


# ---------------------------------------------------------------------------
# PPTX stamping
# ---------------------------------------------------------------------------


def test_escape_hatch_stamps_recipe_id(tmp_path: Path) -> None:
    """The picture frame gets the standard `recipeId` extension AND the
    extra escape-hatch payload extension."""
    deck = IRDeck(
        version=1,
        slides=[_slide_with(_escape_node(bbox=IRBbox(x=0, y=0, w=200, h=200)))],
    )
    out = compile_ir(deck, tmp_path / "stamp.pptx")
    assert isinstance(out, Path)

    prs = Presentation(str(out))
    n_recipe = 0
    n_escape = 0
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        for ext in sp_pr.findall(f"{{{NS_A}}}extLst/{{{NS_A}}}ext"):
            uri = ext.get("uri")
            if uri == SLIDIFY_RECIPE_NS:
                n_recipe += 1
            elif uri == SLIDIFY_ESCAPE_NS:
                n_escape += 1
    assert n_recipe >= 1, "recipe-id extension missing"
    assert n_escape == 1, "escape-hatch payload extension missing"


def test_escape_hatch_recognized_by_recipe_id_fallback(tmp_path: Path) -> None:
    """A node with recipeId='chrome.escape-hatch' but no `role` metadata
    is still treated as an escape-hatch (defensive fallback for
    handcrafted IR fixtures)."""
    node = IRRasterNode(
        kind="raster",
        recipeId="chrome.escape-hatch",
        bbox=IRBbox(x=0, y=0, w=640, h=720),
        zOrder=0,
        metadata={"intent": "ad-hoc-grid"},  # no `role` key
        pngBase64=_TRANSPARENT_1PX_PNG_B64,
    )
    deck = IRDeck(version=1, slides=[_slide_with(node)])
    _, metering = compile_ir(
        deck, tmp_path / "fallback.pptx", return_metering=True
    )
    report = metering.to_report_dict()
    assert report["value"] > 0
    assert "ad-hoc-grid" in report["byIntent"]


# ---------------------------------------------------------------------------
# atoms.yaml: the row exists and is loadable
# ---------------------------------------------------------------------------


def test_atoms_yaml_contains_chrome_escape_hatch() -> None:
    """The matcher catalog must include the `chrome.escape-hatch` row so
    HTML decks tagged `data-atom="chrome.escape-hatch"` route through it."""
    from slidify.patterns.matcher import get_default_patterns

    pats = get_default_patterns()
    matches = [p for p in pats if p.id == "chrome.escape-hatch"]
    assert len(matches) == 1, "exactly one chrome.escape-hatch row expected"
    p = matches[0]
    assert p.priority == 90
    # Lives in the 80–94 band: after the typography/atom-specific rules,
    # before the catch-all `atom-unrecognised` (priority 95).
    assert p.match.get("anchor.data_atom_id") == "chrome.escape-hatch"
    assert p.emit.get("kind") == "Raster"
    assert p.emit.get("metadata", {}).get("excludeFromNativeRatio") is True


def test_decode_placeholder_png_is_valid() -> None:
    """The 1×1 placeholder PNG used as a fallback decodes to valid PNG."""
    raw = base64.b64decode(_TRANSPARENT_1PX_PNG_B64)
    assert raw.startswith(b"\x89PNG\r\n\x1a\n")
