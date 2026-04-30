"""Tests for the layered native-shape decoration system.

Decorations let one source element emit a stack of N shapes (mesh glow,
hairline, rim highlight, inset glow, drop shadow) instead of a single
flat rectangle — the foundation for layered Linear/Vercel/Stripe-style
visual sophistication.
"""

from __future__ import annotations

from pptx import Presentation

from slidify.decorations import (
    DecorationStack,
    DropShadowLayer,
    Hairline,
    InsetGlow,
    MeshGlow,
    RimHighlight,
    derive_decorations,
)
from slidify.models import BoundingBox, DomElement
from slidify.shadows import BoxShadow

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _slide():
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bbox():
    return BoundingBox(x=0, y=0, w=400, h=300)


def _el(hint: str = "") -> DomElement:
    return DomElement(
        id=0,
        parent_id=None,
        depth=0,
        tag="DIV",
        bbox=_bbox(),
        decorate_hint=hint,
        stable_selector="#x",
    )


def test_decoration_stack_emits_below_then_above_in_z_order():
    stack = DecorationStack()
    stack.add(Hairline())            # z=2
    stack.add(MeshGlow(palette=["6366f1"]))  # z=-10
    stack.add(RimHighlight())        # z=1
    slide = _slide()
    below = stack.emit_below(slide, _bbox())
    above = stack.emit_above(slide, _bbox())
    # Mesh hero = 4 blobs below.
    assert len(below) >= 1
    # RimHighlight + Hairline = 2 above.
    assert len(above) == 2


def test_meshglow_emits_hero_blobs():
    slide = _slide()
    m = MeshGlow(palette=["6366f1", "a855f7", "ec4899", "f59e0b"], style="hero")
    shapes = m.emit(slide, _bbox())
    assert len(shapes) == 4


def test_rim_highlight_emits_one_top_band_with_gradient():
    slide = _slide()
    rh = RimHighlight(color_hex="ffffff", opacity=0.2, height_pct=0.5)
    shapes = rh.emit(slide, _bbox())
    assert len(shapes) == 1
    sp_pr = shapes[0]._element.spPr
    grad = sp_pr.find(f"{{{NS_A}}}gradFill")
    assert grad is not None
    stops = grad.findall(f"{{{NS_A}}}gsLst/{{{NS_A}}}gs")
    assert len(stops) >= 2


def test_hairline_emits_transparent_rect_with_alpha_stroke():
    slide = _slide()
    h = Hairline(color_hex="ffffff", opacity=0.10)
    shapes = h.emit(slide, _bbox())
    assert len(shapes) == 1
    sp_pr = shapes[0]._element.spPr
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    alpha = ln.find(f".//{{{NS_A}}}alpha")
    assert alpha is not None
    assert int(alpha.get("val")) == 10000


def test_inset_glow_attaches_inner_shadow():
    slide = _slide()
    ig = InsetGlow(color_hex="ffffff", blur_px=12, opacity=0.5)
    shapes = ig.emit(slide, _bbox())
    sp_pr = shapes[0]._element.spPr
    eff = sp_pr.find(f"{{{NS_A}}}effectLst")
    assert eff is not None
    inner = eff.find(f"{{{NS_A}}}innerShdw")
    assert inner is not None


def test_drop_shadow_layer_attaches_multiple_outer_shadows():
    slide = _slide()
    layers = [
        BoxShadow(False, 0, 10, 15, 0, "000000", 0.10),
        BoxShadow(False, 0, 4, 6, 0, "000000", 0.20),
    ]
    dsl = DropShadowLayer(layers=layers)
    shapes = dsl.emit(slide, _bbox())
    sp_pr = shapes[0]._element.spPr
    eff = sp_pr.find(f"{{{NS_A}}}effectLst")
    assert eff is not None
    outer = eff.findall(f"{{{NS_A}}}outerShdw")
    assert len(outer) == 2


def test_derive_decorations_hero_hint_returns_meshglow_and_hairline():
    el = _el("hero")
    stack = derive_decorations(el)
    assert any(isinstance(d, MeshGlow) for d in stack.layers)
    assert any(isinstance(d, Hairline) for d in stack.layers)


def test_derive_decorations_glass_hint_returns_full_glass_stack():
    el = _el("glass")
    stack = derive_decorations(el)
    types = {type(d).__name__ for d in stack.layers}
    assert "RimHighlight" in types
    assert "Hairline" in types
    assert "InsetGlow" in types


def test_derive_decorations_recessed_hint_returns_dark_inset_glow():
    el = _el("recessed")
    stack = derive_decorations(el)
    assert len(stack.layers) == 1
    glow = stack.layers[0]
    assert isinstance(glow, InsetGlow)
    assert glow.color_hex == "000000"


def test_derive_decorations_no_hint_returns_empty_stack():
    el = _el("")
    stack = derive_decorations(el)
    assert stack.is_empty()


def test_derive_decorations_unknown_hint_returns_empty_stack():
    el = _el("nonsense-hint")
    stack = derive_decorations(el)
    assert stack.is_empty()


def test_decoration_stack_round_trip_save_load(tmp_path):
    """A full layered stack must save+load via python-pptx without error."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    stack = DecorationStack()
    stack.add(MeshGlow(palette=["6366f1", "ec4899"], style="hero"))
    stack.add(RimHighlight())
    stack.add(Hairline())
    stack.emit_below(slide, _bbox())
    stack.emit_above(slide, _bbox())
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    reloaded = Presentation(str(out))
    assert len(reloaded.slides) == 1
