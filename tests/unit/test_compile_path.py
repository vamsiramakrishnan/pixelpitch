"""Tests for `IRPathShapeNode` → `<a:custGeom>` compilation (CONTRACT §3.1).

Covers each PathCommand op (M, L, C, Q, A, Z) plus arrowheads, dash
patterns, and the recipeId stamp.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

from slidify.compile_ir import NS_A, compile_ir
from slidify.ir import (
    FillSolid,
    IRArrowhead,
    IRBbox,
    IRDeck,
    IRPathCommand,
    IRPathShapeNode,
    IRSlide,
)


def _compile_one(node: IRPathShapeNode, tmp_path: Path) -> Presentation:
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")
    return Presentation(str(out))


def _shape_sp_pr(prs: Presentation):
    """Return the spPr lxml element for the first non-background shape."""
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        # Skip the slide-background rect (it has prstGeom rect, not custGeom).
        if sp_pr.find(f"{{{NS_A}}}custGeom") is not None:
            return sh, sp_pr
    return None, None


def test_path_emits_custgeom_with_moveto(tmp_path):
    """`M 100 100 L 200 200 Z` → custGeom with moveTo + lnTo + close."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.line",
        bbox=IRBbox(x=100, y=100, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=100, y=100),
            IRPathCommand(op="L", x=200, y=200),
            IRPathCommand(op="Z"),
        ],
        fill=FillSolid(kind="solid", color="#ff0000"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    assert sp_pr is not None
    cust = sp_pr.find(f"{{{NS_A}}}custGeom")
    assert cust is not None
    path = cust.find(f"{{{NS_A}}}pathLst/{{{NS_A}}}path")
    assert path is not None
    move = path.find(f"{{{NS_A}}}moveTo")
    line = path.find(f"{{{NS_A}}}lnTo")
    close = path.find(f"{{{NS_A}}}close")
    assert move is not None
    assert line is not None
    assert close is not None


def test_path_cubic_bezier(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.cubic",
        bbox=IRBbox(x=0, y=0, w=200, h=200),
        commands=[
            IRPathCommand(op="M", x=0, y=0),
            IRPathCommand(op="C", x1=50, y1=50, x2=150, y2=150, x=200, y=200),
        ],
        fill=FillSolid(kind="solid", color="#00ff00"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    cubic = sp_pr.find(
        f"{{{NS_A}}}custGeom/{{{NS_A}}}pathLst/{{{NS_A}}}path/{{{NS_A}}}cubicBezTo"
    )
    assert cubic is not None
    pts = cubic.findall(f"{{{NS_A}}}pt")
    assert len(pts) == 3  # 2 controls + endpoint


def test_path_quadratic_bezier(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.quad",
        bbox=IRBbox(x=0, y=0, w=200, h=200),
        commands=[
            IRPathCommand(op="M", x=0, y=0),
            IRPathCommand(op="Q", x1=100, y1=100, x=200, y=0),
        ],
        fill=FillSolid(kind="solid", color="#0000ff"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    quad = sp_pr.find(
        f"{{{NS_A}}}custGeom/{{{NS_A}}}pathLst/{{{NS_A}}}path/{{{NS_A}}}quadBezTo"
    )
    assert quad is not None
    pts = quad.findall(f"{{{NS_A}}}pt")
    assert len(pts) == 2  # control + endpoint


def test_path_arc_flattens_to_cubics(tmp_path):
    """SVG `A` arcs lower to ≤4 cubic Bezier segments (no native arcTo)."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.arc",
        bbox=IRBbox(x=0, y=0, w=200, h=200),
        commands=[
            IRPathCommand(op="M", x=0, y=100),
            IRPathCommand(
                op="A",
                rx=100,
                ry=100,
                xAxisRotationDeg=0,
                largeArc=False,
                sweep=True,
                x=200,
                y=100,
            ),
        ],
        fill=FillSolid(kind="solid", color="#ff00ff"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    cubics = sp_pr.findall(
        f"{{{NS_A}}}custGeom/{{{NS_A}}}pathLst/{{{NS_A}}}path/{{{NS_A}}}cubicBezTo"
    )
    assert 1 <= len(cubics) <= 4
    # Bbox of the arc should still hold the cubic endpoints inside [0..units].
    units_max = 100_000
    for cb in cubics:
        for pt in cb.findall(f"{{{NS_A}}}pt"):
            assert 0 <= int(pt.get("x")) <= units_max
            assert 0 <= int(pt.get("y")) <= units_max


def test_path_round_trip_reopens(tmp_path):
    """python-pptx must re-open the emitted file (round-trip safety)."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.roundtrip",
        bbox=IRBbox(x=10, y=10, w=300, h=300),
        commands=[
            IRPathCommand(op="M", x=10, y=10),
            IRPathCommand(op="L", x=310, y=10),
            IRPathCommand(op="L", x=310, y=310),
            IRPathCommand(op="L", x=10, y=310),
            IRPathCommand(op="Z"),
        ],
        fill=FillSolid(kind="solid", color="#abcdef"),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "rt.pptx")
    # Re-open directly from disk.
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    # Re-open from bytes too.
    with open(out, "rb") as f:
        Presentation(io.BytesIO(f.read()))


def test_path_arrowhead_emits_head_and_tail_end(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.arrowed",
        bbox=IRBbox(x=0, y=0, w=200, h=10),
        commands=[
            IRPathCommand(op="M", x=0, y=5),
            IRPathCommand(op="L", x=200, y=5),
        ],
        strokeWidthPx=2,
        strokeColor="#000000",
        markerStart=IRArrowhead(kind="dot", size="sm"),
        markerEnd=IRArrowhead(kind="arrow", size="lg"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    head = ln.find(f"{{{NS_A}}}headEnd")
    tail = ln.find(f"{{{NS_A}}}tailEnd")
    assert head is not None and head.get("type") == "oval"
    assert head.get("w") == "sm" and head.get("len") == "sm"
    assert tail is not None and tail.get("type") == "triangle"
    assert tail.get("w") == "lg" and tail.get("len") == "lg"


def test_path_dash_pattern_maps_to_prst_dash(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.dashed",
        bbox=IRBbox(x=0, y=0, w=100, h=10),
        commands=[
            IRPathCommand(op="M", x=0, y=5),
            IRPathCommand(op="L", x=100, y=5),
        ],
        strokeWidthPx=1,
        strokeColor="#000000",
        strokeDasharray=[4, 2],
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    prst = ln.find(f"{{{NS_A}}}prstDash")
    assert prst is not None
    assert prst.get("val") == "dash"


def test_path_custom_dash_pattern_falls_back_to_custdash(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.custdash",
        bbox=IRBbox(x=0, y=0, w=100, h=10),
        commands=[
            IRPathCommand(op="M", x=0, y=5),
            IRPathCommand(op="L", x=100, y=5),
        ],
        strokeWidthPx=1,
        strokeColor="#000000",
        strokeDasharray=[7, 3, 1, 3],
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    cust = ln.find(f"{{{NS_A}}}custDash")
    assert cust is not None
    ds_entries = cust.findall(f"{{{NS_A}}}ds")
    assert len(ds_entries) == 2  # two (d, sp) pairs


def test_path_all_ops_round_trip(tmp_path):
    """A single path exercising M, L, C, Q, A, Z all in one shape."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.kitchen-sink",
        bbox=IRBbox(x=0, y=0, w=400, h=400),
        commands=[
            IRPathCommand(op="M", x=0, y=0),
            IRPathCommand(op="L", x=100, y=0),
            IRPathCommand(op="Q", x1=200, y1=50, x=200, y=100),
            IRPathCommand(op="C", x1=200, y1=200, x2=300, y2=300, x=400, y=400),
            IRPathCommand(
                op="A",
                rx=50,
                ry=50,
                xAxisRotationDeg=0,
                largeArc=False,
                sweep=True,
                x=300,
                y=400,
            ),
            IRPathCommand(op="Z"),
        ],
        fill=FillSolid(kind="solid", color="#102030"),
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    assert sp_pr is not None
    path = sp_pr.find(
        f"{{{NS_A}}}custGeom/{{{NS_A}}}pathLst/{{{NS_A}}}path"
    )
    assert path is not None
    # M(1) + L(1) + Q(1) + C(1) + A→cubics(≥1) + Z(1)
    children_tags = [etree.QName(c).localname for c in path]
    assert "moveTo" in children_tags
    assert "lnTo" in children_tags
    assert "quadBezTo" in children_tags
    assert "cubicBezTo" in children_tags
    assert "close" in children_tags


def test_path_recipe_id_stamped(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.with-id",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=0, y=0),
            IRPathCommand(op="L", x=100, y=100),
            IRPathCommand(op="Z"),
        ],
        fill=FillSolid(kind="solid", color="#111111"),
    )
    prs = _compile_one(node, tmp_path)
    sh, sp_pr = _shape_sp_pr(prs)
    SLIDIFY = "https://slidify.dev/2026/recipe"
    found = False
    for ext in sp_pr.findall(f"{{{NS_A}}}extLst/{{{NS_A}}}ext"):
        if ext.get("uri") == SLIDIFY:
            found = True
    assert found


def test_path_bbox_omitted_is_inferred(tmp_path):
    """When bbox is omitted, the compiler infers it from path_bbox()."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.inferred-bbox",
        commands=[
            IRPathCommand(op="M", x=50, y=50),
            IRPathCommand(op="L", x=150, y=50),
            IRPathCommand(op="L", x=150, y=150),
            IRPathCommand(op="Z"),
        ],
        fill=FillSolid(kind="solid", color="#202020"),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    # Should not crash even though bbox is None.
    out = compile_ir(deck, tmp_path / "inferred.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_path_linecap_round(tmp_path):
    """strokeLinecap='round' → <a:ln cap="rnd">."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.cap-round",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=10, y=50),
            IRPathCommand(op="L", x=90, y=50),
        ],
        strokeWidthPx=4,
        strokeColor="#888888",
        strokeLinecap="round",
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    assert ln.get("cap") == "rnd"


def test_path_linecap_square(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.cap-square",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=10, y=50),
            IRPathCommand(op="L", x=90, y=50),
        ],
        strokeWidthPx=4,
        strokeColor="#888888",
        strokeLinecap="square",
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    assert ln.get("cap") == "sq"


def test_path_linejoin_round(tmp_path):
    """strokeLinejoin='round' → <a:ln><a:round/>."""
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.join-round",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=10, y=10),
            IRPathCommand(op="L", x=90, y=10),
            IRPathCommand(op="L", x=90, y=90),
        ],
        strokeWidthPx=4,
        strokeColor="#888888",
        strokeLinejoin="round",
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    assert ln.find(f"{{{NS_A}}}round") is not None


def test_path_linejoin_bevel(tmp_path):
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.join-bevel",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=10, y=10),
            IRPathCommand(op="L", x=90, y=10),
            IRPathCommand(op="L", x=90, y=90),
        ],
        strokeWidthPx=4,
        strokeColor="#888888",
        strokeLinejoin="bevel",
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    assert ln.find(f"{{{NS_A}}}bevel") is not None


def test_path_odd_length_dasharray_doubles_per_svg_spec(tmp_path):
    """`[3, 2, 1]` should expand to `[3, 2, 1, 3, 2, 1]` (3 ds children).

    Per SVG/CSS spec, odd-length dash arrays repeat once to become even.
    The previous loop dropped the trailing element silently.
    """
    node = IRPathShapeNode(
        kind="path",
        recipeId="path.dash-odd",
        bbox=IRBbox(x=0, y=0, w=100, h=100),
        commands=[
            IRPathCommand(op="M", x=10, y=50),
            IRPathCommand(op="L", x=90, y=50),
        ],
        strokeWidthPx=2,
        strokeColor="#888888",
        strokeDasharray=[3.0, 2.0, 1.0],  # odd length
    )
    prs = _compile_one(node, tmp_path)
    _, sp_pr = _shape_sp_pr(prs)
    ln = sp_pr.find(f"{{{NS_A}}}ln")
    assert ln is not None
    cust = ln.find(f"{{{NS_A}}}custDash")
    assert cust is not None
    # [3,2,1] → doubled to [3,2,1,3,2,1] → 3 (dash, gap) pairs
    ds_children = cust.findall(f"{{{NS_A}}}ds")
    assert len(ds_children) == 3, f"expected 3 <a:ds> children, got {len(ds_children)}"
    # Verify the values: pair 1 = (3, 2), pair 2 = (1, 3), pair 3 = (2, 1)
    expected = [(300_000, 200_000), (100_000, 300_000), (200_000, 100_000)]
    actual = [(int(c.get("d")), int(c.get("sp"))) for c in ds_children]
    assert actual == expected, f"dash sequence wrong: {actual}"
