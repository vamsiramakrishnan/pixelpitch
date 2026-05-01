"""Tests for slidify.preset_shapes."""

from __future__ import annotations

import io
import math

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from slidify.models import BoundingBox, DomElement
from slidify.preset_shapes import (
    KNOWN_PRESETS,
    PresetMatch,
    classify_polygon,
    clip_path_to_preset,
    detect_preset_shape,
    emit_preset,
    parse_circle_clip_path,
    parse_inset_clip_path,
    parse_polygon_clip_path,
)


def _el(**overrides) -> DomElement:
    base = dict(
        id=1,
        parent_id=None,
        depth=0,
        tag="div",
        bbox=BoundingBox(x=0, y=0, w=200, h=100),
    )
    base.update(overrides)
    return DomElement(**base)


def _slide():
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return prs, slide


# 1.
def test_chevron_clip_path():
    el = _el(
        clip_path="polygon(0% 0%, 75% 0%, 100% 50%, 75% 100%, 0% 100%, 25% 50%)",
        bbox=BoundingBox(x=0, y=0, w=200, h=80),
    )
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.CHEVRON


# 2.
def test_arrow_right_class():
    el = _el(cls="arrow-right cta", bbox=BoundingBox(x=0, y=0, w=240, h=80))
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.RIGHT_ARROW


# 3.
def test_pill_class_and_radius():
    el1 = _el(cls="pill button", bbox=BoundingBox(x=0, y=0, w=200, h=40))
    m1 = detect_preset_shape(el1)
    assert m1 is not None
    assert m1.preset == MSO_SHAPE.ROUNDED_RECTANGLE

    el2 = _el(border_radius="9999px", bbox=BoundingBox(x=0, y=0, w=200, h=40))
    m2 = detect_preset_shape(el2)
    assert m2 is not None
    assert m2.preset == MSO_SHAPE.ROUNDED_RECTANGLE


# 4.
def test_oval_from_radius_50_on_square():
    el = _el(border_radius="50%", bbox=BoundingBox(x=0, y=0, w=80, h=80))
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.OVAL


# 5.
def test_plain_rectangle_returns_none():
    el = _el(
        cls="card",
        border_radius="0px",
        clip_path="none",
        bbox=BoundingBox(x=0, y=0, w=300, h=200),
    )
    assert detect_preset_shape(el) is None


# 6 + 7.
def test_emit_preset_adds_shape_and_round_trips(tmp_path):
    prs, slide = _slide()
    bbox = BoundingBox(x=10, y=10, w=200, h=120)
    match = PresetMatch(preset=MSO_SHAPE.HEXAGON, confidence=1.0, reason="test")
    shape = emit_preset(slide, bbox, match)
    assert shape is not None
    assert shape.auto_shape_type == MSO_SHAPE.HEXAGON

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    reopened = Presentation(buf)
    assert len(reopened.slides) == 1
    shapes = list(reopened.slides[0].shapes)
    assert len(shapes) >= 1
    assert shapes[0].auto_shape_type == MSO_SHAPE.HEXAGON


# 8.
def test_polygon_signature_hexagon():
    pts = []
    for i in range(6):
        ang = 2 * math.pi * i / 6
        pts.append((0.5 + 0.5 * math.cos(ang), 0.5 + 0.5 * math.sin(ang)))
    match = classify_polygon(pts)
    assert match is not None
    assert match.preset == MSO_SHAPE.HEXAGON


# 9.
def test_polygon_signature_5_point_star():
    pts = []
    for i in range(10):
        ang = -math.pi / 2 + math.pi * i / 5
        r = 0.5 if i % 2 == 0 else 0.22
        pts.append((0.5 + r * math.cos(ang), 0.5 + r * math.sin(ang)))
    match = classify_polygon(pts)
    assert match is not None
    assert match.preset == MSO_SHAPE.STAR_5_POINT


# 10.
def test_known_presets_telemetry_has_30_plus():
    assert len(KNOWN_PRESETS) >= 30


# Additional sanity tests
def test_parse_polygon_clip_path_basic():
    pts = parse_polygon_clip_path("polygon(0% 0%, 100% 0%, 50% 100%)")
    assert pts == [(0.0, 0.0), (1.0, 0.0), (0.5, 1.0)]


def test_parse_polygon_clip_path_invalid():
    assert parse_polygon_clip_path("none") is None
    assert parse_polygon_clip_path("circle(50%)") is None


def test_class_chevron():
    el = _el(cls="chevron-right step", bbox=BoundingBox(x=0, y=0, w=200, h=60))
    m = detect_preset_shape(el)
    assert m is not None
    assert m.preset == MSO_SHAPE.CHEVRON


def test_class_lightning():
    el = _el(cls="lightning glow")
    m = detect_preset_shape(el)
    assert m is not None
    assert m.preset == MSO_SHAPE.LIGHTNING_BOLT


def test_polygon_pentagon():
    pts = []
    for i in range(5):
        ang = -math.pi / 2 + 2 * math.pi * i / 5
        pts.append((0.5 + 0.5 * math.cos(ang), 0.5 + 0.5 * math.sin(ang)))
    match = classify_polygon(pts)
    assert match is not None
    assert match.preset == MSO_SHAPE.PENTAGON


def test_polygon_octagon():
    pts = []
    for i in range(8):
        ang = math.pi / 8 + 2 * math.pi * i / 8
        pts.append((0.5 + 0.5 * math.cos(ang), 0.5 + 0.5 * math.sin(ang)))
    match = classify_polygon(pts)
    assert match is not None
    assert match.preset == MSO_SHAPE.OCTAGON


def test_emit_preset_with_rotation():
    _, slide = _slide()
    bbox = BoundingBox(x=10, y=10, w=120, h=120)
    match = PresetMatch(preset=MSO_SHAPE.STAR_5_POINT, rotation_deg=45.0)
    shape = emit_preset(slide, bbox, match)
    assert abs(shape.rotation - 45.0) < 0.01


# clip-path: inset(...) ------------------------------------------------------


def test_inset_uniform_percent_to_rectangle():
    bbox = BoundingBox(x=0, y=0, w=200, h=100)
    inset = parse_inset_clip_path("inset(10%)", bbox)
    assert inset is not None
    new_bbox, radius = inset
    assert radius == 0.0
    assert new_bbox.x == 20.0  # 10% of 200
    assert new_bbox.y == 10.0  # 10% of 100
    assert new_bbox.w == 160.0
    assert new_bbox.h == 80.0


def test_inset_two_value_percent():
    bbox = BoundingBox(x=0, y=0, w=200, h=100)
    inset = parse_inset_clip_path("inset(10% 20%)", bbox)
    assert inset is not None
    new_bbox, _ = inset
    # top/bottom = 10% of 100, left/right = 20% of 200
    assert new_bbox.x == 40.0
    assert new_bbox.y == 10.0
    assert new_bbox.w == 120.0  # 200 - 40 - 40
    assert new_bbox.h == 80.0


def test_inset_four_value_pixel():
    bbox = BoundingBox(x=100, y=200, w=300, h=200)
    inset = parse_inset_clip_path("inset(5px 10px 15px 20px)", bbox)
    assert inset is not None
    new_bbox, _ = inset
    assert new_bbox.x == 120.0  # 100 + 20px left
    assert new_bbox.y == 205.0  # 200 + 5px top
    assert new_bbox.w == 270.0  # 300 - 10 right - 20 left
    assert new_bbox.h == 180.0  # 200 - 5 top - 15 bottom


def test_inset_round_yields_rounded_rectangle():
    el = _el(
        clip_path="inset(10% round 8px)",
        bbox=BoundingBox(x=0, y=0, w=200, h=100),
    )
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.ROUNDED_RECTANGLE
    assert match.bbox_override is not None
    assert match.bbox_override.w == 160.0


def test_inset_no_round_yields_rectangle():
    el = _el(
        clip_path="inset(20px)",
        bbox=BoundingBox(x=0, y=0, w=200, h=100),
    )
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.RECTANGLE
    assert match.bbox_override is not None
    assert match.bbox_override.w == 160.0
    assert match.bbox_override.h == 60.0


def test_inset_invalid_returns_none():
    bbox = BoundingBox(x=0, y=0, w=200, h=100)
    assert parse_inset_clip_path("none", bbox) is None
    assert parse_inset_clip_path("inset()", bbox) is None
    assert parse_inset_clip_path("circle(50%)", bbox) is None


# clip-path: circle(...) -----------------------------------------------------


def test_circle_default_inscribes_in_square():
    bbox = BoundingBox(x=0, y=0, w=100, h=100)
    visible = parse_circle_clip_path("circle()", bbox)
    assert visible is not None
    assert visible.x == 0.0
    assert visible.y == 0.0
    assert visible.w == 100.0
    assert visible.h == 100.0


def test_circle_default_inscribes_in_rect():
    bbox = BoundingBox(x=0, y=0, w=200, h=100)
    visible = parse_circle_clip_path("circle()", bbox)
    assert visible is not None
    # closest-side from center (100, 50) → min(100, 100, 50, 50) = 50
    assert visible.w == 100.0
    assert visible.h == 100.0
    assert visible.x == 50.0  # centered
    assert visible.y == 0.0


def test_circle_pixel_radius_at_center():
    bbox = BoundingBox(x=10, y=20, w=200, h=100)
    visible = parse_circle_clip_path("circle(30px at 50% 50%)", bbox)
    assert visible is not None
    assert visible.w == 60.0
    assert visible.h == 60.0
    assert visible.x == 80.0  # 10 + 100 - 30
    assert visible.y == 40.0  # 20 + 50 - 30


def test_circle_yields_oval_preset():
    el = _el(
        clip_path="circle(50%)",
        bbox=BoundingBox(x=0, y=0, w=80, h=80),
    )
    match = detect_preset_shape(el)
    assert match is not None
    assert match.preset == MSO_SHAPE.OVAL
    assert match.bbox_override is not None
    # circle(50%) on a square: radius = 50% * sqrt(2*w^2)/sqrt(2) = w/2
    assert abs(match.bbox_override.w - 80.0) < 1e-6


# clip_path_to_preset entry point -------------------------------------------


def test_clip_path_to_preset_polygon_dispatches():
    bbox = BoundingBox(x=0, y=0, w=100, h=100)
    pts = []
    for i in range(6):
        ang = 2 * math.pi * i / 6
        pts.append((0.5 + 0.5 * math.cos(ang), 0.5 + 0.5 * math.sin(ang)))
    poly = "polygon(" + ", ".join(f"{p[0] * 100}% {p[1] * 100}%" for p in pts) + ")"
    match = clip_path_to_preset(poly, bbox)
    assert match is not None
    assert match.preset == MSO_SHAPE.HEXAGON


def test_clip_path_to_preset_unknown_returns_none():
    bbox = BoundingBox(x=0, y=0, w=100, h=100)
    assert clip_path_to_preset("path('M 0 0 L 1 1')", bbox) is None
    assert clip_path_to_preset("url(#mask)", bbox) is None
    assert clip_path_to_preset("none", bbox) is None
    assert clip_path_to_preset("", bbox) is None
