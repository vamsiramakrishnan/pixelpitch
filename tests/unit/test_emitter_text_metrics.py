"""Integration test: verify the emitter bakes <a:normAutofit fontScale=...>
into the slide XML for a too-wide title that overflows in the fallback font."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Emu

from slidify.emitter import Emitter, _apply_explicit_autofit, _union_line_box
from slidify.models import (
    BoundingBox,
    Decision,
    DecisionKind,
    DomElement,
    EmitOp,
    TextRun,
    UnitKind,
    VisualUnit,
)
from slidify.text_metrics import get_inter_font_path

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_text_element(
    eid: int, text: str, bbox: BoundingBox, font_size_px: float
) -> DomElement:
    return DomElement(
        id=eid,
        parent_id=None,
        depth=0,
        tag="H1",
        bbox=bbox,
        text=text,
        font_family="Inter, sans-serif",
        font_size=f"{font_size_px}px",
        font_weight="700",
        color="rgb(0, 0, 0)",
        stable_selector=f"#e{eid}",
    )


@pytest.mark.asyncio
async def test_emitter_shrinks_overflowing_title(tmp_path):
    """Multi-line wrapping textbox: assert <a:normAutofit fontScale="..."> is
    baked. Single-line ``wrap=none`` boxes bake the scale into rPr instead
    (LibreOffice ignores normAutofit when wrap=none) — that path is covered
    by ``test_single_line_title_bakes_size_into_rPr``.
    """
    if get_inter_font_path() is None:
        pytest.skip("Inter not installed — autofit pre-compute requires it")
    # Long title that, in any reasonable font, overflows a tight bbox.
    # Use a tall bbox so the wrap=true / multi-line branch is selected and
    # we can assert the normAutofit element specifically.
    long_text = "An Extraordinarily Long Presentation Title That Will Overflow"
    bbox = BoundingBox(x=10, y=10, w=200, h=300)
    el = _make_text_element(1, long_text, bbox, font_size_px=32.0)
    unit = VisualUnit(
        id="u1",
        kind=UnitKind.Title,
        bbox=bbox,
        elements=[el],
    )
    op = EmitOp(
        unit_id="u1",
        decision=Decision(kind=DecisionKind.NativeText),
        z_order=0,
        bbox=bbox,
    )

    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    em._emit_native_text(slide, unit, op)

    out = tmp_path / "out.pptx"
    em.save(out)
    em.close()

    # Reopen and inspect the XML.
    prs = Presentation(str(out))
    s = prs.slides[0]
    found_scales: list[int] = []
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        bodyPr = shape.text_frame._txBody.bodyPr
        for autofit in bodyPr.findall(f"{{{A_NS}}}normAutofit"):
            scale = autofit.get("fontScale")
            if scale is not None:
                found_scales.append(int(scale))
    assert found_scales, "expected at least one <a:normAutofit fontScale=...>"
    assert any(s < 100_000 for s in found_scales), (
        f"expected a fontScale < 100000 (shrink); got {found_scales}"
    )


@pytest.mark.asyncio
async def test_single_line_title_bakes_size_into_rPr(tmp_path):
    """Bug A regression: single-line wrap=none titles must bake the
    autofit scale into per-run ``<a:rPr sz=...>`` rather than emit
    ``<a:normAutofit fontScale="...">``. LibreOffice silently drops the
    fontScale attribute when bodyPr.wrap="none", which left slide-13's
    72pt title clipped at the slide edge with the substituted font even
    though the scale had been computed correctly.
    """
    if get_inter_font_path() is None:
        pytest.skip("Inter not installed — autofit pre-compute requires it")
    # 72pt headline in a single-line, just-too-narrow box (mirrors slide-13:
    # ``"Pro, in the smallest possible frame."`` at 72px font in ~1126px wide).
    long_text = "Pro, in the smallest possible frame."
    bbox = BoundingBox(x=10, y=10, w=900, h=80)
    el = _make_text_element(1, long_text, bbox, font_size_px=72.0)
    unit = VisualUnit(
        id="u1", kind=UnitKind.Title, bbox=bbox, elements=[el]
    )
    op = EmitOp(
        unit_id="u1",
        decision=Decision(kind=DecisionKind.NativeText),
        z_order=0,
        bbox=bbox,
    )

    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    em._emit_native_text(slide, unit, op)
    out = tmp_path / "out.pptx"
    em.save(out)
    em.close()

    # Reopen and inspect the title textbox's bodyPr + rPr.
    prs = Presentation(str(out))
    title_shape = next(
        sh for sh in prs.slides[0].shapes
        if sh.has_text_frame and "Pro," in sh.text_frame.text
    )
    bodyPr = title_shape.text_frame._txBody.bodyPr
    # 1) wrap=none was set (single-line guard)
    assert bodyPr.get("wrap") == "none"
    # 2) NO normAutofit fontScale on a wrap=none box (LibreOffice ignores it).
    autofits = bodyPr.findall(f"{{{A_NS}}}normAutofit")
    for af in autofits:
        # fontScale absent OR == 100000 (no-op) is acceptable; any < 100000 is
        # the bug we're guarding against.
        fs = af.get("fontScale")
        assert fs is None or int(fs) >= 100_000, (
            f"single-line wrap=none box should not rely on normAutofit; "
            f"got fontScale={fs}"
        )
    # 3) The per-run rPr sz must be smaller than the source 72pt → 7200.
    rprs = title_shape.text_frame._txBody.findall(
        f".//{{{A_NS}}}r/{{{A_NS}}}rPr"
    )
    sizes = [int(r.get("sz")) for r in rprs if r.get("sz") is not None]
    assert sizes, "expected at least one <a:rPr sz=...>"
    assert min(sizes) < 7200, (
        f"expected a baked-down sz < 7200; got {sizes}"
    )


def test_union_line_box_includes_leaf_text_without_runs():
    """Bug B regression: _union_line_box must fall back to el.bbox for
    text-bearing elements that lack `runs` (the walker only populates
    runs for multi-run text containers). When a unit holds both a leaf
    `<h1>` and a multi-run `.sub`, the union must encompass BOTH —
    otherwise slide-17's 180px "Thank you." headline gets squashed into
    the 29px tall sub-line bbox.
    """
    title = DomElement(
        id=1,
        parent_id=None,
        depth=0,
        tag="H1",
        bbox=BoundingBox(x=80, y=400, w=928, h=166),
        text="Thank you.",
        font_size="180px",
        stable_selector="#title",
    )
    sub = DomElement(
        id=2,
        parent_id=None,
        depth=0,
        tag="DIV",
        bbox=BoundingBox(x=80, y=512, w=372, h=31),
        text=None,
        is_text_container=True,
        runs=[
            TextRun(
                text="Available everywhere on May 14.",
                font_size="24px",
                line_boxes=[BoundingBox(x=80, y=513, w=372, h=29)],
            )
        ],
        stable_selector="#sub",
    )
    unit = VisualUnit(
        id="u1",
        kind=UnitKind.Generic,
        bbox=BoundingBox(x=80, y=400, w=928, h=166),
        elements=[title, sub],
    )
    union = _union_line_box(unit)
    assert union is not None
    # Must span from the title's top (~400) to at least the sub's bottom
    # (~542). Pre-fix this returned (80, 513, 372, 29).
    assert union.y <= 410, f"union should start at title top, got y={union.y}"
    assert union.y + union.h >= 540, (
        f"union should reach sub bottom, got y2={union.y + union.h}"
    )
    assert union.w >= 900, (
        f"union should be at least the title width, got w={union.w}"
    )


def test_union_line_box_returns_none_when_no_text():
    """Sanity: a unit with zero text elements still returns None
    (caller falls back to op.bbox)."""
    el = DomElement(
        id=1,
        parent_id=None,
        depth=0,
        tag="DIV",
        bbox=BoundingBox(x=0, y=0, w=100, h=100),
        background_color="rgb(0, 128, 255)",
        stable_selector="#bg",
    )
    unit = VisualUnit(
        id="u1",
        kind=UnitKind.Generic,
        bbox=BoundingBox(x=0, y=0, w=100, h=100),
        elements=[el],
    )
    assert _union_line_box(unit) is None


def test_apply_explicit_autofit_replaces_existing():
    # Build a textbox, then call _apply_explicit_autofit twice — confirm only
    # one normAutofit remains and it has the latest values.
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
    tf = tb.text_frame
    tf.text = "x"

    _apply_explicit_autofit(tf, 90_000, 0)
    _apply_explicit_autofit(tf, 75_000, 5_000)

    bodyPr = tf._txBody.bodyPr
    autofits = bodyPr.findall(f"{{{A_NS}}}normAutofit")
    assert len(autofits) == 1
    assert autofits[0].get("fontScale") == "75000"
    assert autofits[0].get("lnSpcReduction") == "5000"
    # No spAutoFit/noAutofit left behind either.
    assert bodyPr.find(f"{{{A_NS}}}spAutoFit") is None
    assert bodyPr.find(f"{{{A_NS}}}noAutofit") is None
