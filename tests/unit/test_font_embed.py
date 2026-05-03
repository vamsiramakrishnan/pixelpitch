"""Tests for font embedding into the .pptx package.

Verifies that after `embed_fonts_in_pptx` runs:
  * the font binary appears at /ppt/fonts/font*.fntdata inside the zip
  * presentation.xml.rels contains a relationship of type `…/font`
  * presentation.xml gains a <p:embeddedFontLst> with the typeface
  * [Content_Types].xml declares the obfuscatedFont default extension

We don't check that PowerPoint actually USES the font (that requires
PowerPoint); the verification here is that the OOXML package is well-formed
and the font part is referenced from the right places.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

from slidify.font_embed import (
    FontVariants,
    audit_font_bindings,
    discover_inter,
    embed_default_fonts,
    embed_fonts_in_pptx,
)

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
FONT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
)


def _minimal_pptx(tmp_path: Path) -> Path:
    """Generate a minimal blank .pptx via python-pptx."""
    out = tmp_path / "blank.pptx"
    p = Presentation()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(out))
    return out


def _fonted_pptx(tmp_path: Path, typeface: str) -> Path:
    out = tmp_path / "fonted.pptx"
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1_000_000, 300_000)
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Editorial"
    run.font.name = typeface
    p.save(str(out))
    return out


def test_discover_inter_finds_otfs():
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    assert fv.typeface == "Inter"
    assert fv.regular.exists()
    # Other variants may be missing on minimal systems — only Regular is required.


def test_embed_inter_adds_font_part(tmp_path: Path):
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    pptx = _minimal_pptx(tmp_path)
    embed_fonts_in_pptx(pptx, [fv])
    with zipfile.ZipFile(pptx, "r") as zf:
        names = set(zf.namelist())
        # Font binaries land under /ppt/fonts/font{N}.fntdata
        font_parts = [n for n in names if n.startswith("ppt/fonts/font")]
        assert font_parts, "no font part written"
        # Each font part should be non-empty
        for fp in font_parts:
            assert len(zf.read(fp)) > 1024, f"{fp} suspiciously small"


def test_embed_inter_registers_relationship(tmp_path: Path):
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    pptx = _minimal_pptx(tmp_path)
    embed_fonts_in_pptx(pptx, [fv])
    with zipfile.ZipFile(pptx, "r") as zf:
        rels = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    font_rels = [
        r
        for r in rels.findall(f"{{{NS_REL}}}Relationship")
        if r.get("Type") == FONT_REL_TYPE
    ]
    assert len(font_rels) >= 1, "no font relationship in presentation rels"
    # Each relationship should target a fonts/font*.fntdata path
    for r in font_rels:
        assert r.get("Target", "").startswith("fonts/font")


def test_embed_inter_adds_embedded_font_lst(tmp_path: Path):
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    pptx = _minimal_pptx(tmp_path)
    embed_fonts_in_pptx(pptx, [fv])
    with zipfile.ZipFile(pptx, "r") as zf:
        pres = etree.fromstring(zf.read("ppt/presentation.xml"))
    elst = pres.find(f"{{{NS_P}}}embeddedFontLst")
    assert elst is not None, "<p:embeddedFontLst> missing"
    fonts = elst.findall(f"{{{NS_P}}}embeddedFont")
    assert len(fonts) == 1
    name = fonts[0].find(f"{{{NS_P}}}font")
    assert name is not None and name.get("typeface") == "Inter"
    # Each declared variant <p:regular>/<p:bold>/<p:italic>/<p:boldItalic>
    # should reference a real relationship id.
    for tag in ("regular", "bold", "italic", "boldItalic"):
        v = fonts[0].find(f"{{{NS_P}}}{tag}")
        if v is not None:
            assert v.get(f"{{{NS_R}}}id"), f"no r:id on {tag}"


def test_audit_font_bindings_reports_non_core_missing_embed(tmp_path: Path):
    pptx = _fonted_pptx(tmp_path, "Playfair Display")
    audit = audit_font_bindings(pptx)
    assert "Playfair Display" in audit.referenced
    assert "Playfair Display" in audit.missing_embeds


def test_audit_font_bindings_passes_after_embedding_requested_face(tmp_path: Path):
    pptx = _fonted_pptx(tmp_path, "Playfair Display")
    fake_font = tmp_path / "fake.ttf"
    fake_font.write_bytes(b"not-a-real-font-but-a-font-part")
    embed_fonts_in_pptx(
        pptx,
        [FontVariants(typeface="Playfair Display", regular=fake_font)],
    )
    audit = audit_font_bindings(pptx)
    assert "Playfair Display" in audit.embedded
    assert "Playfair Display" not in audit.missing_embeds


def test_embed_inter_adds_content_type(tmp_path: Path):
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    pptx = _minimal_pptx(tmp_path)
    embed_fonts_in_pptx(pptx, [fv])
    with zipfile.ZipFile(pptx, "r") as zf:
        ct = etree.fromstring(zf.read("[Content_Types].xml"))
    found = False
    for d in ct.findall(f"{{{NS_CT}}}Default"):
        if d.get("Extension") == "fntdata":
            found = True
            assert "obfuscatedFont" in d.get("ContentType", "")
    assert found, "<Default Extension='fntdata'> missing from Content_Types"


def test_embedded_pptx_still_opens_via_python_pptx(tmp_path: Path):
    """Re-opening the patched .pptx through python-pptx must still work —
    that's a fast check that we didn't break the package structure."""
    fv = discover_inter()
    if fv is None:
        pytest.skip("Inter not installed on this system")
    pptx = _minimal_pptx(tmp_path)
    embed_fonts_in_pptx(pptx, [fv])
    p = Presentation(str(pptx))
    assert len(p.slides) == 1


def test_embed_default_fonts_returns_false_when_inter_absent(
    tmp_path: Path, monkeypatch
):
    """If Inter isn't on the system, embed_default_fonts is a no-op (returns
    False) — never raises, never corrupts the .pptx."""
    pptx = _minimal_pptx(tmp_path)
    # Force discover_inter to return None.
    import slidify.font_embed as fe

    monkeypatch.setattr(fe, "discover_inter", lambda: None)
    assert embed_default_fonts(pptx) is False
    # Original .pptx still opens.
    Presentation(str(pptx))
