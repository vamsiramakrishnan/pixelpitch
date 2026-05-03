"""Tests for the pre-save sanitization pass that prevents PowerPoint's
"repair dialog" from popping at open time.

Patterns covered:
  * Empty <a:p/> gets <a:endParaRPr lang="en-US"/> inserted.
  * Whitespace-bearing <a:t> gains xml:space="preserve" if missing.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slidify.emitter import Emitter, _rotation_degrees

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


def test_save_inserts_endpararpr_on_empty_paragraphs(tmp_path: Path):
    em = Emitter()
    layout = em.prs.slide_layouts[6]
    slide = em.prs.slides.add_slide(layout)
    tb = slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000)
    tf = tb.text_frame
    # A paragraph with content.
    tf.paragraphs[0].add_run().text = "x"
    # An intentionally empty paragraph (acts as visual spacer).
    tf.add_paragraph()

    out = tmp_path / "deck.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    slide_el = prs.slides[0]._element
    paragraphs = list(slide_el.iter(f"{{{NS_A}}}p"))
    assert any(
        len(p.findall(f"{{{NS_A}}}r")) == 0
        and p.find(f"{{{NS_A}}}endParaRPr") is not None
        for p in paragraphs
    ), "empty <a:p> still missing <a:endParaRPr>"


def test_save_preserves_whitespace_in_runs(tmp_path: Path):
    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 5_000_000, 1_000_000)
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "  leading and trailing  "

    out = tmp_path / "deck.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    slide_el = prs.slides[0]._element
    ts = list(slide_el.iter(f"{{{NS_A}}}t"))
    relevant = [t for t in ts if t.text and t.text.strip() == "leading and trailing"]
    assert relevant, "test text not found"
    for t in relevant:
        assert t.get(XML_SPACE) == "preserve", f"missing xml:space on '{t.text}'"


def test_sanitize_is_idempotent(tmp_path: Path):
    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000)
    tb.text_frame.paragraphs[0].add_run().text = "x"
    tb.text_frame.add_paragraph()

    em._sanitize_for_repair_dialog()
    em._sanitize_for_repair_dialog()

    # Should not have stacked multiple endParaRPr children.
    slide_el = slide._element
    for p in slide_el.iter(f"{{{NS_A}}}p"):
        assert len(p.findall(f"{{{NS_A}}}endParaRPr")) <= 1
    em.close()


def test_save_roundtrip_still_valid(tmp_path: Path):
    """The hygiene pass must not break the package — save+load should
    succeed and the slide count must match."""
    em = Emitter()
    for _ in range(3):
        em.prs.slides.add_slide(em.prs.slide_layouts[6])
    out = tmp_path / "deck.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_rotation_degrees_extracts_css_2d_matrix():
    assert round(_rotation_degrees("matrix(0, 1, -1, 0, 10, 20)"), 1) == 90.0
    assert round(_rotation_degrees("rotate(-12deg)"), 1) == -12.0


def test_emit_op_skips_off_canvas_shapes(tmp_path: Path):
    """Source HTML can naturally overflow the slide bbox (e.g. a 5-row
    grid taller than 720px). Without an off-canvas guard, every shape past
    the slide bottom got clamped to h=1 and emitted invisibly — bloating
    shape count and confusing region-attribution. ``_emit_op`` now skips
    them entirely."""
    import asyncio

    from slidify.models import (
        BoundingBox,
        Decision,
        DecisionKind,
        DomElement,
        EmitOp,
        UnitKind,
        VisualUnit,
    )
    from slidify.renderer import RenderedSlide

    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])

    on_bbox = BoundingBox(x=64, y=400, w=200, h=40)
    off_bbox = BoundingBox(x=64, y=900, w=200, h=120)
    on_el = DomElement(
        id=1, parent_id=None, depth=0, tag="DIV", bbox=on_bbox,
        text="visible", stable_selector="#on",
    )
    off_el = DomElement(
        id=2, parent_id=None, depth=0, tag="DIV", bbox=off_bbox,
        text="off-canvas", stable_selector="#off",
    )
    on_unit = VisualUnit(id="u_on", kind=UnitKind.Generic, bbox=on_bbox, elements=[on_el])
    off_unit = VisualUnit(id="u_off", kind=UnitKind.Generic, bbox=off_bbox, elements=[off_el])
    units_by_id = {"u_on": on_unit, "u_off": off_unit}

    on_op = EmitOp(
        unit_id="u_on", decision=Decision(kind=DecisionKind.NativeText),
        z_order=0, bbox=on_bbox,
    )
    off_op = EmitOp(
        unit_id="u_off", decision=Decision(kind=DecisionKind.NativeText),
        z_order=1, bbox=off_bbox,
    )

    rendered = RenderedSlide(
        html="",
        elements=[on_el, off_el],
        viewport_w=1280,
        viewport_h=720,
        ground_truth_png=b"",
    )

    async def run():
        await em._emit_op(slide, on_op, units_by_id, rendered, None)
        await em._emit_op(slide, off_op, units_by_id, rendered, None)

    asyncio.run(run())

    out = tmp_path / "deck.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    s = prs.slides[0]
    texts = [sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
    assert "visible" in "".join(texts)
    assert "off-canvas" not in "".join(texts), (
        f"off-canvas shape should have been skipped; got {texts}"
    )
