"""Tests for CSS box-shadow parsing + multi-layer OOXML emission.

Tailwind's `shadow-xl` / `shadow-2xl` / etc decompose into 2+ comma-
separated layers in computed style. Emitting only the first layer loses
the second (often a tighter, sharper highlight that gives the shadow
its dimensional shape).
"""

from __future__ import annotations

from lxml import etree
from pptx import Presentation

from slidify.shadows import (
    NS_A,
    apply_shadows,
    parse_box_shadow,
    parse_box_shadows,
)


def test_parse_single_shadow_unchanged():
    sh = parse_box_shadow("0px 4px 6px -1px rgba(0, 0, 0, 0.1)")
    assert sh is not None
    assert sh.offset_y_px == 4
    assert sh.blur_px == 6
    assert sh.color_hex.lower() == "000000"


def test_parse_box_shadows_returns_all_layers():
    """Tailwind shadow-lg = `0 10px 15px -3px rgba(0,0,0,0.1),
                            0 4px 6px -4px rgba(0,0,0,0.1)`"""
    layers = parse_box_shadows(
        "0px 10px 15px -3px rgba(0, 0, 0, 0.1), "
        "0px 4px 6px -4px rgba(0, 0, 0, 0.1)"
    )
    assert len(layers) == 2
    assert layers[0].blur_px == 15
    assert layers[1].blur_px == 6


def test_parse_box_shadows_skips_unparseable_layers():
    layers = parse_box_shadows(
        "garbage, 0px 4px 6px rgba(0, 0, 0, 0.2), nonsense too"
    )
    # The 1-numeric "garbage" / "nonsense" don't have 2+ px values, so
    # only the middle (real) layer survives.
    assert len(layers) == 1
    assert layers[0].blur_px == 6


def test_parse_box_shadows_handles_none_and_empty():
    assert parse_box_shadows(None) == []
    assert parse_box_shadows("") == []
    assert parse_box_shadows("none") == []


def test_apply_shadows_emits_multiple_outerShdw_in_one_effectLst():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000)

    layers = parse_box_shadows(
        "0px 10px 15px -3px rgba(0, 0, 0, 0.1), "
        "0px 4px 6px -4px rgba(0, 0, 0, 0.2)"
    )
    assert apply_shadows(tb, layers) is True

    sp_pr = tb._element.spPr
    effect_lst = sp_pr.find(f"{{{NS_A}}}effectLst")
    assert effect_lst is not None
    shadows = effect_lst.findall(f"{{{NS_A}}}outerShdw")
    assert len(shadows) == 2

    # Order preserved (CSS-source order = OOXML-child order).
    blurs = [int(s.get("blurRad")) for s in shadows]
    # 15px → 15*9525 = 142875 ; 6px → 57150
    assert blurs == [142875, 57150]

    # Alpha preserved (0.1 → 10000, 0.2 → 20000).
    alphas = [
        s.find(f"{{{NS_A}}}srgbClr").find(f"{{{NS_A}}}alpha")
        for s in shadows
    ]
    assert alphas[0] is not None and alphas[1] is not None
    assert int(alphas[0].get("val")) == 10000
    assert int(alphas[1].get("val")) == 20000


def test_apply_shadows_with_inset_layer_uses_innerShdw():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000)

    layers = parse_box_shadows(
        "0px 4px 6px rgba(0,0,0,0.1), inset 0px 2px 4px rgba(0,0,0,0.2)"
    )
    apply_shadows(tb, layers)
    effect_lst = tb._element.spPr.find(f"{{{NS_A}}}effectLst")
    assert effect_lst.find(f"{{{NS_A}}}outerShdw") is not None
    assert effect_lst.find(f"{{{NS_A}}}innerShdw") is not None
