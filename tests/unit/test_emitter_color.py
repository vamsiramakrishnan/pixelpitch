"""Unit tests for emitter color resolution — specifically the
gradient-clipped-text fallback.

CSS pattern:
    .accent {
      background: linear-gradient(...);
      -webkit-background-clip: text;
      color: transparent;
    }

We cannot natively reproduce gradient-clipped text in PPTX, so the emitter
substitutes the gradient's first stop color when the run is otherwise
transparent.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor

from slidify.emitter import _resolve_run_color
from slidify.models import BoundingBox, DomElement


def _el(color: str = "rgb(255, 255, 255)") -> DomElement:
    return DomElement(
        id=0,
        parent_id=None,
        depth=0,
        tag="DIV",
        bbox=BoundingBox(x=0, y=0, w=10, h=10),
        color=color,
        stable_selector="#x",
    )


def test_resolve_run_color_returns_solid():
    spec = {"text": "hi", "color": "rgb(99, 102, 241)"}
    rgb = _resolve_run_color(spec, _el())
    assert rgb == RGBColor(99, 102, 241)


def test_resolve_run_color_with_transparent_falls_through_to_gradient_first_stop():
    """Gradient-clipped text: color=transparent + bg-image=gradient →
    use the FIRST gradient stop's color as a solid fallback."""
    spec = {
        "text": "presentations",
        "color": "rgba(0, 0, 0, 0)",
        "background_image": "linear-gradient(135deg, rgb(129, 140, 248) 0%, rgb(192, 132, 252) 50%, rgb(244, 114, 182) 100%)",
    }
    rgb = _resolve_run_color(spec, _el())
    # First stop = #818cf8 (indigo-400 in v3 / similar in v4)
    assert rgb == RGBColor(0x81, 0x8c, 0xf8)


def test_resolve_run_color_transparent_no_gradient_returns_none():
    """No bg-image gradient available → return None (caller leaves color unset)."""
    spec = {"text": "x", "color": "rgba(0, 0, 0, 0)", "background_image": "none"}
    rgb = _resolve_run_color(spec, _el(color="rgba(0, 0, 0, 0)"))
    assert rgb is None


def test_resolve_run_color_falls_back_to_element_color_if_spec_missing():
    spec = {"text": "x"}
    rgb = _resolve_run_color(spec, _el(color="rgb(50, 60, 70)"))
    assert rgb == RGBColor(50, 60, 70)


def test_native_gradient_text_fill_is_emitted_for_bg_clip_text():
    """When a run has color: transparent + a parseable gradient bg-image,
    the emitter should attach a native <a:gradFill> to the run via
    Font.fill.gradient(). End result: PowerPoint renders the gradient
    through the glyph silhouettes, matching the source's
    `background-clip: text` design."""
    from pathlib import Path

    from pptx import Presentation

    from slidify.emitter import _try_apply_gradient_text_fill

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "presentations"

    spec = {
        "text": "presentations",
        "color": "rgba(0, 0, 0, 0)",
        "background_image": "linear-gradient(135deg, rgb(129, 140, 248), rgb(244, 114, 182))",
    }
    applied = _try_apply_gradient_text_fill(run.font, spec)
    assert applied is True

    # Verify the rPr now contains a <a:gradFill> with two stops.
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    rpr = run.font._rPr
    assert rpr is not None
    grad = rpr.find(f"{{{NS_A}}}gradFill")
    assert grad is not None
    stops = grad.findall(f"{{{NS_A}}}gsLst/{{{NS_A}}}gs")
    assert len(stops) >= 2  # may be auto-extended to 3 with terminal stop


def test_native_gradient_text_fill_skips_non_transparent_runs():
    """A run with a real solid color should NOT trigger gradient text fill
    even if the run carries a background_image — it's not a bg-clip text."""
    from pptx import Presentation

    from slidify.emitter import _try_apply_gradient_text_fill

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "x"
    spec = {
        "text": "x",
        "color": "rgb(100, 100, 100)",
        "background_image": "linear-gradient(135deg, #abc, #def)",
    }
    assert _try_apply_gradient_text_fill(run.font, spec) is False
