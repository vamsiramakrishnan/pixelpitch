"""Tests for the post-emit editability round-trip check.

The check re-opens the produced .pptx and asserts that the per-slide count
of editable primitives matches what the emitter was asked to produce.
A drop indicates python-pptx (or our sanitization pass) silently lost a
shape — invisible to the SSIM oracle (which just compares pixels).
"""

from __future__ import annotations

from pathlib import Path

from pptx.util import Emu

from slidify.emitter import Emitter
from slidify.models import (
    BoundingBox,
    Decision,
    DecisionKind,
    EmitOp,
)
from slidify.roundtrip import (
    _intended_counts,
    check_pptx_editability,
)


def _op(kind: DecisionKind, x: int = 0, y: int = 0, w: int = 100, h: int = 50) -> EmitOp:
    return EmitOp(
        unit_id=f"u_{x}_{y}",
        decision=Decision(kind=kind, confidence=1.0, source_tier="tier1"),
        z_order=0,
        bbox=BoundingBox(x=x, y=y, w=w, h=h),
    )


def test_intended_counts_split_by_primitive():
    ops = [
        _op(DecisionKind.NativeText),
        _op(DecisionKind.NativeText),
        _op(DecisionKind.NativeShape),
        _op(DecisionKind.NativePicture),
        _op(DecisionKind.NativeTable),
        _op(DecisionKind.NativeBullet),
        _op(DecisionKind.Raster),  # not editable
        _op(DecisionKind.Skip),    # not editable
        _op(DecisionKind.Hybrid),  # editable (decoration + native children)
    ]
    counts = _intended_counts(ops)
    # Editable = NativeText*2 + NativeShape + NativePicture + NativeTable +
    #            NativeBullet + Hybrid = 7 (NativeSvg is also editable but
    #            unused here). Raster + Skip excluded.
    assert counts["editable"] == 7
    assert counts["pictures"] == 1
    assert counts["tables"] == 1
    assert counts["text_frames"] == 3  # 2 NativeText + 1 NativeBullet


def test_check_passes_when_actual_meets_intent(tmp_path: Path):
    em = Emitter()
    layout = em.prs.slide_layouts[6]
    slide = em.prs.slides.add_slide(layout)
    # Manually drop a textbox to simulate a NativeText emit.
    tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
    tb.text_frame.text = "hello"
    out = tmp_path / "ok.pptx"
    em.save(out)
    em.close()

    ops = [_op(DecisionKind.NativeText)]
    report = check_pptx_editability(out, [ops])
    assert report.n_slides == 1
    assert report.n_passed == 1
    assert report.passed is True
    assert report.per_slide[0].actual_text_frames == 1


def test_check_flags_dropped_shape(tmp_path: Path):
    em = Emitter()
    em.prs.slides.add_slide(em.prs.slide_layouts[6])  # blank slide, no shapes
    out = tmp_path / "drop.pptx"
    em.save(out)
    em.close()

    # Intent: 3 native text frames + 1 picture; actual: 0.
    ops = [
        _op(DecisionKind.NativeText),
        _op(DecisionKind.NativeText),
        _op(DecisionKind.NativeText),
        _op(DecisionKind.NativePicture),
    ]
    report = check_pptx_editability(out, [ops])
    assert report.n_slides == 1
    assert report.n_passed == 0
    assert report.passed is False
    s0 = report.per_slide[0]
    assert s0.intended_editable == 4
    assert s0.actual_editable == 0
    assert "editable shapes" in s0.notes
    assert "pictures" in s0.notes


def test_extra_shapes_pass_with_drop_only_failure(tmp_path: Path):
    """Decoration stacks legitimately add MORE shapes than ops; that's fine."""
    em = Emitter()
    layout = em.prs.slide_layouts[6]
    slide = em.prs.slides.add_slide(layout)
    # Two textboxes (decoration overlay + the actual text), but intent is 1.
    for txt in ("decoration", "real text"):
        tb = slide.shapes.add_textbox(
            Emu(0), Emu(0), Emu(1_000_000), Emu(500_000)
        )
        tb.text_frame.text = txt
    out = tmp_path / "extra.pptx"
    em.save(out)
    em.close()

    ops = [_op(DecisionKind.NativeText)]
    report = check_pptx_editability(out, [ops])
    assert report.passed is True
    assert report.per_slide[0].actual_text_frames == 2


def test_skip_only_slide_passes_with_zero_intent(tmp_path: Path):
    em = Emitter()
    em.prs.slides.add_slide(em.prs.slide_layouts[6])
    out = tmp_path / "skip.pptx"
    em.save(out)
    em.close()

    ops = [_op(DecisionKind.Skip), _op(DecisionKind.Raster)]
    report = check_pptx_editability(out, [ops])
    assert report.passed is True
    assert report.per_slide[0].intended_editable == 0
