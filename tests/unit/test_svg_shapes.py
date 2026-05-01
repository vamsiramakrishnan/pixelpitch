"""Tests for slidify.svg_shapes — SVG primitive → PPTX shape translation.

Focused on the fill/stroke opacity behavior because regressions there
(silently rendering at full alpha) wash out organic-blob slides.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu

from slidify.svg_shapes import _build_mapping, emit_svg_shapes

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _new_slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_circle_with_fill_opacity_emits_alpha():
    """Bug C regression: SVG `fill-opacity="0.55"` must reach the PPTX
    shape as `<a:srgbClr><a:alpha val="55000"/></a:srgbClr>`. Without
    this the 55% blobs on slide-43 / slide-44 painted at full opacity
    and drowned the rest of the canvas.
    """
    slide = _new_slide()
    n = emit_svg_shapes(
        slide,
        [
            {
                "tag": "circle",
                "cx": 100,
                "cy": 100,
                "r": 50,
                "fill": "#a78bfa",
                "fill_opacity": 0.55,
                "stroke": "none",
                "stroke_width": 0,
                "_svgRect": {"x": 0, "y": 0, "w": 1280, "h": 720},
            }
        ],
    )
    assert n == 1
    shape = next(iter(slide.shapes))
    # Verify alpha child on the solidFill > srgbClr
    sp_pr = shape._element.spPr
    solid = sp_pr.find(f"{{{A_NS}}}solidFill")
    assert solid is not None
    srgb = solid.find(f"{{{A_NS}}}srgbClr")
    assert srgb is not None
    alpha = srgb.find(f"{{{A_NS}}}alpha")
    assert alpha is not None, "expected <a:alpha> baked in for fill_opacity=0.55"
    # 0.55 → 55_000 (in 1/1000 of a percent)
    assert int(alpha.get("val")) == 55_000


def test_path_with_fill_opacity_emits_alpha():
    """Same guarantee for closed cubic-bezier blob paths (slide-43 ritual
    background blob, slide-44 corner blobs)."""
    slide = _new_slide()
    d = (
        "M 60 60 C 220 28, 460 50, 540 140 "
        "C 620 230, 540 270, 360 256 "
        "C 180 242, 30 220, 60 60 Z"
    )
    n = emit_svg_shapes(
        slide,
        [
            {
                "tag": "path",
                "d": d,
                "fill": "#f7b2a2",
                "fill_opacity": 0.32,
                "stroke": "none",
                "stroke_width": 0,
                "_svgRect": {"x": 0, "y": 0, "w": 1280, "h": 720},
            }
        ],
    )
    assert n == 1
    shape = next(iter(slide.shapes))
    sp_pr = shape._element.spPr
    # Path emit replaces prstGeom with custGeom — verify that and that
    # the alpha is baked.
    assert sp_pr.find(f"{{{A_NS}}}custGeom") is not None
    solid = sp_pr.find(f"{{{A_NS}}}solidFill")
    assert solid is not None
    srgb = solid.find(f"{{{A_NS}}}srgbClr")
    alpha = srgb.find(f"{{{A_NS}}}alpha")
    assert alpha is not None
    assert int(alpha.get("val")) == 32_000


def test_fill_opacity_one_does_not_emit_alpha():
    """No alpha child for fully-opaque fills — keeps the XML tidy and
    avoids fooling renderers that special-case the absence of `<a:alpha>`.
    """
    slide = _new_slide()
    emit_svg_shapes(
        slide,
        [
            {
                "tag": "rect",
                "x": 0,
                "y": 0,
                "w": 100,
                "h": 100,
                "fill": "#000",
                "fill_opacity": 1.0,
                "stroke": "none",
                "stroke_width": 0,
                "_svgRect": {"x": 0, "y": 0, "w": 100, "h": 100},
            }
        ],
    )
    shape = next(iter(slide.shapes))
    sp_pr = shape._element.spPr
    solid = sp_pr.find(f"{{{A_NS}}}solidFill")
    srgb = solid.find(f"{{{A_NS}}}srgbClr") if solid is not None else None
    if srgb is not None:
        assert srgb.find(f"{{{A_NS}}}alpha") is None


# --- _build_mapping (viewBox + preserveAspectRatio) -----------------------


def test_mapping_no_viewbox_is_identity():
    """SVGs without a viewBox have userspace == viewport pixels (1:1)."""
    m = _build_mapping({"x": 100, "y": 50, "w": 800, "h": 600})
    assert m.scale_x == 1.0 and m.scale_y == 1.0
    assert m.x(0) == 100 and m.y(0) == 50
    assert m.x(40) == 140 and m.y(60) == 110


def test_mapping_uniform_viewbox_scales_uniformly():
    """viewBox `0 0 100 100` rendered at 200×200 → 2× uniform scale."""
    m = _build_mapping(
        {"x": 0, "y": 0, "w": 200, "h": 200, "viewBox": [0, 0, 100, 100]}
    )
    assert m.scale_x == 2.0 and m.scale_y == 2.0
    assert m.x(50) == 100 and m.y(50) == 100


def test_mapping_default_preserve_aspect_meet_centers_on_wide_viewport():
    """Default preserveAspectRatio is `xMidYMid meet`. A square viewBox
    rendered into a wider viewport keeps the smaller scale and centers
    horizontally — slight letterboxing on the sides."""
    # viewBox 100×100 → viewport 200×100. scale = min(2, 1) = 1.
    # Extra width = 200 - 100*1 = 100, centered → ox = 50.
    m = _build_mapping(
        {"x": 0, "y": 0, "w": 200, "h": 100, "viewBox": [0, 0, 100, 100]}
    )
    assert m.scale_x == 1.0 and m.scale_y == 1.0
    assert m.x(0) == 50 and m.y(0) == 0
    assert m.x(100) == 150 and m.y(100) == 100


def test_mapping_preserve_aspect_none_stretches_independently():
    """preserveAspectRatio="none" stretches each axis independently."""
    m = _build_mapping(
        {
            "x": 0,
            "y": 0,
            "w": 200,
            "h": 100,
            "viewBox": [0, 0, 100, 100],
            "preserveAspectRatio": "none",
        }
    )
    assert m.scale_x == 2.0 and m.scale_y == 1.0
    assert m.x(50) == 100 and m.y(50) == 50


def test_mapping_viewbox_with_nonzero_origin_translates_correctly():
    """A viewBox `-10 -20 100 100` shifts userspace: a userspace point
    (-10, -20) maps to the viewport origin."""
    m = _build_mapping(
        {"x": 0, "y": 0, "w": 100, "h": 100, "viewBox": [-10, -20, 100, 100]}
    )
    assert m.x(-10) == 0 and m.y(-20) == 0
    assert m.x(0) == 10 and m.y(0) == 20


def test_mapping_xminymin_meet_anchors_top_left():
    """preserveAspectRatio="xMinYMin meet" keeps the smaller scale but
    anchors at the top-left instead of centering."""
    m = _build_mapping(
        {
            "x": 0,
            "y": 0,
            "w": 200,
            "h": 100,
            "viewBox": [0, 0, 100, 100],
            "preserveAspectRatio": "xMinYMin meet",
        }
    )
    assert m.scale_x == 1.0
    assert m.x(0) == 0  # no horizontal centering offset
    assert m.x(100) == 100


def test_mapping_offset_viewport_with_viewbox():
    """SVGs are usually placed at non-zero positions on the slide; the
    viewport offset must compose correctly with the viewBox-derived scale."""
    m = _build_mapping(
        {"x": 64, "y": 32, "w": 640, "h": 360, "viewBox": [0, 0, 100, 100]}
    )
    # Square viewBox in 640×360 viewport, default meet → scale = min(6.4, 3.6) = 3.6.
    # Extra width = 640 - 100*3.6 = 280, centered → ox = 140.
    assert m.scale_x == 3.6 and m.scale_y == 3.6
    assert m.x(0) == 64 + 140
    assert m.y(0) == 32


def test_mapping_invalid_viewbox_falls_back_to_identity():
    """Malformed viewBox (zero/negative dims, wrong arity) silently degrades
    to 1:1 mapping rather than producing NaN/division-by-zero coords."""
    assert _build_mapping(
        {"x": 0, "y": 0, "w": 100, "h": 100, "viewBox": [0, 0, 0, 100]}
    ).scale_x == 1.0
    assert _build_mapping(
        {"x": 0, "y": 0, "w": 100, "h": 100, "viewBox": [0, 0, 100]}
    ).scale_x == 1.0
    assert _build_mapping(
        {"x": 0, "y": 0, "w": 0, "h": 100, "viewBox": [0, 0, 100, 100]}
    ).scale_x == 1.0


def test_stroke_opacity_emits_alpha_on_line_color():
    """slide-44 has `<path stroke-opacity="0.40">` for the wavy divider.
    The opacity must land on the line's color, not just the fill.
    """
    slide = _new_slide()
    emit_svg_shapes(
        slide,
        [
            {
                "tag": "path",
                "d": "M 0 0 C 60 -20, 120 20, 180 0",
                "fill": "none",
                "fill_opacity": 1.0,
                "stroke": "#264653",
                "stroke_width": 1.2,
                "stroke_opacity": 0.40,
                "_svgRect": {"x": 0, "y": 0, "w": 1280, "h": 720},
            }
        ],
    )
    shape = next(iter(slide.shapes))
    sp_pr = shape._element.spPr
    ln = sp_pr.find(f"{{{A_NS}}}ln")
    assert ln is not None
    solid = ln.find(f"{{{A_NS}}}solidFill")
    assert solid is not None
    srgb = solid.find(f"{{{A_NS}}}srgbClr")
    alpha = srgb.find(f"{{{A_NS}}}alpha")
    assert alpha is not None, "expected <a:alpha> on stroke for stroke-opacity"
    assert int(alpha.get("val")) == 40_000
