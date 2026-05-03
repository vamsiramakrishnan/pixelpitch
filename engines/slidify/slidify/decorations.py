"""Layered native-shape decoration system.

Today the emitter writes one PPTX shape per DOM element. Modern visual
sophistication (Linear, Vercel, Stripe, Anthropic homepage) is *layered*:
a base + a soft glow + a noise overlay + a highlight + a hairline, each a
separate translucent shape stacked along z.

A `Decoration` paints one such layer into a slide at a given bbox. A
`DecorationStack` is an ordered collection: layers with `z_offset < 0`
emit BELOW the main content shape; `z_offset >= 0` emit ABOVE. Within
each band, lower z_offset is painted first (so higher z appears on top).

Decorations are opt-in. Heuristics never silently inflate shape count.
Two ways to opt in:

  1. HTML: add `data-slidify-decorate="hero"` (or "glass" / "tactile" /
     "recessed" / "aurora") to any element. The DOM walker captures it
     as `el.decorate_hint`; the emitter calls `derive_decorations(el)`.

  2. IR: composite templates (StatCardWithDepth, GlassPanel, etc.)
     attach a DecorationStack directly when emitting.

This module declares the abstract base + concrete decorations that wrap
the layered primitives shipping in sibling modules:
  - `MeshGlow`   wraps `slidify.mesh`
  - `RimHighlight` / `Hairline` / `InsetGlow` / `GrainOverlay` are local
"""

from __future__ import annotations

from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import TYPE_CHECKING

from slidify.geom import px_to_emu
from slidify.models import BoundingBox

if TYPE_CHECKING:
    from pptx.slide import Slide

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _emu_rect(bbox: BoundingBox) -> tuple[int, int, int, int]:
    return (
        px_to_emu(bbox.x),
        px_to_emu(bbox.y),
        px_to_emu(bbox.w),
        px_to_emu(bbox.h),
    )


# ---------------------------------------------------------------------------
# Abstract base
# ---------------------------------------------------------------------------


class Decoration(ABC):
    """One layered shape painted into a slide at a target bbox.

    Subclasses set `z_offset` (negative = below main content; positive =
    above) and implement `emit`.
    """

    z_offset: int = -1

    @abstractmethod
    def emit(self, slide: Slide, bbox: BoundingBox) -> list:
        """Paint shapes onto `slide` covering `bbox`. Returns the shapes
        added so callers can apply additional treatment if needed."""
        raise NotImplementedError


@dataclass
class DecorationStack:
    """Ordered collection of decorations applied around a content shape."""

    layers: list[Decoration] = field(default_factory=list)

    def add(self, decoration: Decoration) -> DecorationStack:
        self.layers.append(decoration)
        return self

    def emit_below(self, slide: Slide, bbox: BoundingBox) -> list:
        out: list = []
        for d in sorted(self.layers, key=lambda d: d.z_offset):
            if d.z_offset < 0:
                out.extend(d.emit(slide, bbox))
        return out

    def emit_above(self, slide: Slide, bbox: BoundingBox) -> list:
        out: list = []
        for d in sorted(self.layers, key=lambda d: d.z_offset):
            if d.z_offset >= 0:
                out.extend(d.emit(slide, bbox))
        return out

    def is_empty(self) -> bool:
        return not self.layers


# ---------------------------------------------------------------------------
# Concrete decorations
# ---------------------------------------------------------------------------


@dataclass
class MeshGlow(Decoration):
    """Stacked translucent radial gradients — the Stripe/Vercel hero look.

    Wires `slidify.mesh` into the decoration system. The blob palette
    is derived from the slide's brand colors via
    `slidify.mesh.derive_mesh_from_palette`.
    """

    palette: list[str]
    style: str = "hero"
    z_offset: int = -10

    def emit(self, slide, bbox: BoundingBox) -> list:
        # Local import — keeps decorations.py importable even if mesh.py
        # is not on disk yet (during agent integration).
        try:
            from slidify.mesh import derive_mesh_from_palette, emit_mesh
        except ImportError:
            return []
        mesh = derive_mesh_from_palette(self.palette, self.style)
        return emit_mesh(slide, bbox, mesh)


@dataclass
class RimHighlight(Decoration):
    """Soft inner highlight along the top edge of a card.

    Gives the surface a subtle sense of depth — like specular highlight
    on a curved button. Emits one rect with a top-down linear gradient
    going from `color@opacity` to transparent over `height_pct`.
    """

    color_hex: str = "ffffff"
    opacity: float = 0.15
    height_pct: float = 0.5
    z_offset: int = 1

    def emit(self, slide, bbox: BoundingBox) -> list:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        from slidify.gradients import (
            GradientStop,
            LinearGradient,
            apply_gradient_fill,
        )

        h = max(2.0, bbox.h * self.height_pct)
        x, y, w, _ = _emu_rect(bbox)
        h_emu = px_to_emu(h)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h_emu)
        try:
            shape.line.fill.background()
        except Exception:
            pass
        grad = LinearGradient(
            angle_deg=180.0,  # top → bottom
            stops=[
                GradientStop(color_hex=self.color_hex, alpha=self.opacity, position=0.0),
                GradientStop(color_hex=self.color_hex, alpha=0.0, position=1.0),
            ],
        )
        apply_gradient_fill(shape, grad)
        return [shape]


@dataclass
class Hairline(Decoration):
    """1-px translucent border — the crisp rim that doesn't look like
    Office defaults. Emits a transparent rect with only a line stroke."""

    color_hex: str = "ffffff"
    opacity: float = 0.10
    width_px: float = 1.0
    z_offset: int = 2

    def emit(self, slide, bbox: BoundingBox) -> list:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        x, y, w, h = _emu_rect(bbox)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        try:
            shape.fill.background()
        except Exception:
            pass
        try:
            shape.line.color.rgb = RGBColor(
                int(self.color_hex[0:2], 16),
                int(self.color_hex[2:4], 16),
                int(self.color_hex[4:6], 16),
            )
            shape.line.width = Emu(px_to_emu(self.width_px))
        except Exception:
            pass
        # Apply alpha to the stroke via lxml — python-pptx doesn't expose it.
        try:
            from lxml import etree

            sp_pr = shape._element.spPr  # noqa: SLF001
            ln = sp_pr.find(f"{{{NS_A}}}ln")
            if ln is not None:
                solid = ln.find(f"{{{NS_A}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{NS_A}}}srgbClr")
                    if srgb is not None:
                        for existing in srgb.findall(f"{{{NS_A}}}alpha"):
                            srgb.remove(existing)
                        etree.SubElement(
                            srgb,
                            f"{{{NS_A}}}alpha",
                            attrib={"val": str(int(round(self.opacity * 100_000)))},
                        )
        except Exception:
            pass
        return [shape]


@dataclass
class InsetGlow(Decoration):
    """Soft inner glow — gives 'recessed' or 'inset' surfaces tactile depth.

    Emits a transparent shape carrying a multi-layer innerShdw effect.
    """

    color_hex: str = "ffffff"
    blur_px: float = 12.0
    opacity: float = 0.5
    z_offset: int = 3

    def emit(self, slide, bbox: BoundingBox) -> list:
        from pptx.enum.shapes import MSO_SHAPE

        from slidify.shadows import BoxShadow, apply_shadows

        x, y, w, h = _emu_rect(bbox)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        try:
            shape.fill.background()
            shape.line.fill.background()
        except Exception:
            pass
        apply_shadows(
            shape,
            [
                BoxShadow(
                    inset=True,
                    offset_x_px=0.0,
                    offset_y_px=0.0,
                    blur_px=self.blur_px,
                    spread_px=0.0,
                    color_hex=self.color_hex.upper(),
                    alpha=self.opacity,
                )
            ],
        )
        return [shape]


@dataclass
class DropShadowLayer(Decoration):
    """Multi-layer drop shadow attached to a transparent silhouette shape.

    Useful when you want the shadow to extend a card's footprint without
    binding it to the card's own shape (e.g. when the card itself is a
    glass overlay whose outerShdw would be partially clipped)."""

    layers: list  # list[BoxShadow]
    z_offset: int = -1

    def emit(self, slide, bbox: BoundingBox) -> list:
        from pptx.enum.shapes import MSO_SHAPE

        from slidify.shadows import apply_shadows

        x, y, w, h = _emu_rect(bbox)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        try:
            shape.fill.background()
            shape.line.fill.background()
        except Exception:
            pass
        if self.layers:
            apply_shadows(shape, self.layers)
        return [shape]


# ---------------------------------------------------------------------------
# HTML opt-in: derive a stack from a `data-slidify-decorate="..."` hint
# ---------------------------------------------------------------------------


_HINT_PALETTES_FALLBACK = [
    "6366f1",  # indigo
    "a855f7",  # violet
    "ec4899",  # pink
    "f59e0b",  # amber
]


def derive_decorations(
    el,
    *,
    palette: list[str] | None = None,
) -> DecorationStack:
    """Build a DecorationStack from an element's `decorate_hint` attribute.

    Hints recognized:
      "hero"     — MeshGlow (4 corner blobs) + Hairline
      "spotlight"— MeshGlow (1 huge centered blob)
      "aurora"   — MeshGlow (3 horizontal bands)
      "glass"    — RimHighlight + Hairline + InsetGlow (white)
      "tactile"  — RimHighlight + Hairline (button-like)
      "recessed" — InsetGlow (dark)

    Returns an empty DecorationStack when no recognized hint.
    """
    stack = DecorationStack()
    hint = (getattr(el, "decorate_hint", "") or "").strip().lower()
    if not hint:
        return stack
    pal = palette or _HINT_PALETTES_FALLBACK
    if hint in ("hero", "spotlight", "aurora", "orbit"):
        stack.add(MeshGlow(palette=pal, style=hint, z_offset=-10))
        if hint == "hero":
            stack.add(Hairline(color_hex="ffffff", opacity=0.06))
    elif hint == "glass":
        stack.add(RimHighlight(color_hex="ffffff", opacity=0.20, height_pct=0.45))
        stack.add(Hairline(color_hex="ffffff", opacity=0.18))
        stack.add(InsetGlow(color_hex="ffffff", opacity=0.10, blur_px=18))
    elif hint == "tactile":
        stack.add(RimHighlight(color_hex="ffffff", opacity=0.30, height_pct=0.5))
        stack.add(Hairline(color_hex="ffffff", opacity=0.15))
    elif hint == "recessed":
        stack.add(InsetGlow(color_hex="000000", opacity=0.35, blur_px=10))
    return stack


__all__ = [
    "Decoration",
    "DecorationStack",
    "DropShadowLayer",
    "Hairline",
    "InsetGlow",
    "MeshGlow",
    "RimHighlight",
    "derive_decorations",
]
