"""Tests for `_emit_picture_in_rounded_rect` failure semantics.

Verifies that a rounded-rect-clipped picture whose source URL is
unreachable does NOT leave an orphan auto-shape on the slide. The
non-clipped path returns before adding anything; the clipped path
should match that contract.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slidify.compile_ir import compile_ir
from slidify.ir import (
    IRBbox,
    IRClipPathRoundedRect,
    IRDeck,
    IRPictureNode,
    IRSlide,
)


def _shape_count(out: Path) -> int:
    """Count non-background shapes on slide 0."""
    prs = Presentation(str(out))
    n = 0
    for sh in prs.slides[0].shapes:
        # Skip the slide-background rect (always present, full-bleed).
        if getattr(sh, "width", None) and sh.width >= prs.slide_width:
            if getattr(sh, "height", None) and sh.height >= prs.slide_height:
                continue
        n += 1
    return n


def test_rounded_rect_clip_unreachable_src_leaves_no_orphan_shape(tmp_path):
    """When _fetch_picture fails, the rounded-rect clip path must NOT
    add the auto-shape — otherwise the slide gets an empty rounded
    rectangle where the picture should have been.
    """
    node = IRPictureNode(
        kind="picture",
        recipeId="picture.unreachable-clipped",
        bbox=IRBbox(x=100, y=100, w=400, h=400),
        src="https://nonexistent.invalid/never/reaches/here.png",
        clipPath=IRClipPathRoundedRect(kind="rounded-rect", radiusPx=24),
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")
    assert out.exists()
    assert _shape_count(out) == 0, (
        "expected zero shapes on slide after failed picture fetch; "
        "the orphan rounded rectangle is the bug being guarded."
    )


def test_unclipped_unreachable_src_also_leaves_no_orphan(tmp_path):
    """Sanity: the non-clipped path already had this contract — verify
    we haven't regressed it while fixing the clipped path.
    """
    node = IRPictureNode(
        kind="picture",
        recipeId="picture.unreachable-flat",
        bbox=IRBbox(x=100, y=100, w=400, h=400),
        src="https://nonexistent.invalid/missing.png",
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "p.pptx")
    assert out.exists()
    assert _shape_count(out) == 0
