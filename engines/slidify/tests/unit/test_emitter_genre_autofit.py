"""Bug B (slide-54): emitter must measure the substituted-font width
with a *genre-matched* fallback (DejaVu Serif for serif primaries),
otherwise a serif headline that LOOKS like it fits in Liberation Sans
gets ``wrap=none`` with a baked size of 100%, then overflows the
bbox horizontally when LibreOffice picks a wider serif fallback.

Slide-54 symptom: "Procurement is consolidating — six vendors absorb
the budget that twenty-two used to share." kept its 12.75pt size and
ran past the 784px insight column into the right-rail stat number.
"""

from __future__ import annotations

import pytest
from pptx import Presentation

from slidify.emitter import Emitter
from slidify.models import (
    BoundingBox,
    Decision,
    DecisionKind,
    DomElement,
    EmitOp,
    UnitKind,
    VisualUnit,
)
from slidify.text_metrics import (
    estimate_wrapped_lines,
    genre_for_family,
    get_genre_fallback_font_path,
    get_inter_font_path,
)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def test_genre_for_family_classifies_serif_and_mono():
    assert genre_for_family("'Source Serif Pro', Georgia, serif") == "serif"
    assert genre_for_family("Georgia, serif") == "serif"
    assert genre_for_family("'JetBrains Mono', Consolas, monospace") == "mono"
    assert genre_for_family("'Inter', sans-serif") == "sans"
    assert genre_for_family("") == "sans"
    assert genre_for_family(None) == "sans"


def test_genre_fallback_font_path_returns_serif_for_serif():
    """Without genre routing, a serif primary's measurement uses
    Liberation Sans (~15% narrower than DejaVu Serif), under-estimating
    overflow by the same factor.
    """
    serif = get_genre_fallback_font_path("serif")
    if serif is None:
        pytest.skip("No serif fallback installed on this system")
    name = str(serif).lower()
    assert "serif" in name or "freeserif" in name


def test_estimate_wrapped_lines_overflows_at_serif_width():
    """The "Procurement..." headline at 12.75pt fits in 1 line at 784px
    in Liberation Sans (~700px) but NOT in DejaVu Serif Bold (~900px).
    The estimator must report > 1 line for serif so the emitter shrinks
    the font instead of leaving the textbox at wrap=none + sz=100%.
    """
    if get_genre_fallback_font_path("serif") is None:
        pytest.skip("No serif fallback installed on this system")
    text = (
        "Procurement is consolidating — six vendors absorb the budget that "
        "twenty-two used to share."
    )
    n = estimate_wrapped_lines(
        text, bbox_w_px=784, font_size_pt=12.75, genre="serif"
    )
    assert n >= 2, (
        f"serif-genre measurement should report >=2 lines for this text "
        f"in 784px; got {n}"
    )


def _make_h3(text: str, bbox: BoundingBox) -> DomElement:
    return DomElement(
        id=1,
        parent_id=None,
        depth=0,
        tag="H3",
        bbox=bbox,
        text=text,
        font_family="'Source Serif Pro', Georgia, serif",
        font_size="17px",  # = 12.75pt
        font_weight="600",
        color="rgb(26, 26, 26)",
        stable_selector="#h3",
    )


@pytest.mark.asyncio
async def test_serif_h3_in_single_line_bbox_shrinks_font(tmp_path):
    """Bug B regression: a serif h3 in a 1-line-tall bbox whose text
    overflows the bbox width with the serif fallback must bake a
    smaller per-run ``sz`` (LibreOffice ignores ``normAutofit`` when
    ``wrap=none``). Without genre-aware measurement the autofit
    silently no-op'd because Liberation Sans was narrow enough.
    """
    if get_inter_font_path() is None:
        pytest.skip("Inter not installed — autofit pre-compute requires it")
    if get_genre_fallback_font_path("serif") is None:
        pytest.skip("No serif fallback installed")

    text = (
        "Procurement is consolidating — six vendors absorb the budget that "
        "twenty-two used to share."
    )
    bbox = BoundingBox(x=168, y=241, w=784, h=22)
    el = _make_h3(text, bbox)
    unit = VisualUnit(id="u1", kind=UnitKind.Generic, bbox=bbox, elements=[el])
    op = EmitOp(
        unit_id="u1",
        decision=Decision(kind=DecisionKind.NativeText),
        z_order=0,
        bbox=bbox,
    )

    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    em._emit_native_text(slide, unit, op)
    out = tmp_path / "out.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    target = next(
        sh for sh in prs.slides[0].shapes
        if sh.has_text_frame and "Procurement" in sh.text_frame.text
    )
    rprs = target.text_frame._txBody.findall(
        f".//{{{A_NS}}}r/{{{A_NS}}}rPr"
    )
    sizes = [int(r.get("sz")) for r in rprs if r.get("sz") is not None]
    assert sizes, "expected at least one <a:rPr sz=...>"
    # 12.75pt = 1275; the autofit should shrink it below the source size.
    # If the genre routing didn't fire, sz would be 1275 (no shrink).
    assert min(sizes) < 1275, (
        f"expected shrunk sz < 1275 (12.75pt source); got {sizes}. "
        "Likely cause: emitter measured the serif headline against a "
        "sans-serif fallback and missed the overflow."
    )
