"""Tests for slidify.mesh — palette → mesh-spec presets and OOXML emission.

The module fakes a Stripe/Vercel/Linear-style mesh gradient by stacking
3-5 translucent radial-gradient rectangles. These tests verify the preset
shapes (corner counts, center placement, vertical spread) and that the
emitted shapes carry valid radial `<a:gradFill>` payloads with alpha.
"""

from __future__ import annotations

import io

from pptx import Presentation

from slidify.mesh import (
    GlowBlob,
    MeshSpec,
    derive_mesh_from_palette,
    emit_mesh,
)
from slidify.models import BoundingBox

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PALETTE = ["#6366f1", "#a855f7", "#ec4899", "#f59e0b"]


# ---- derive_mesh_from_palette ------------------------------------------------


def test_hero_returns_four_corner_blobs():
    spec = derive_mesh_from_palette(PALETTE, "hero")
    assert len(spec.blobs) == 4
    # The four blobs should sit in distinct quadrants — top-left, top-right,
    # bottom-left, bottom-right. We don't pin exact percentages so the preset
    # can evolve, but each corner quadrant must have one center.
    quadrants = {(b.cx_pct < 0.5, b.cy_pct < 0.5) for b in spec.blobs}
    assert quadrants == {(True, True), (False, True), (True, False), (False, False)}
    # Each blob inherits a palette color (no leading '#', uppercase hex).
    for b in spec.blobs:
        assert len(b.color_hex) == 6
        assert b.color_hex == b.color_hex.upper()


def test_spotlight_returns_one_large_centered_blob():
    spec = derive_mesh_from_palette(PALETTE, "spotlight")
    assert len(spec.blobs) == 1
    only = spec.blobs[0]
    assert abs(only.cx_pct - 0.5) < 0.05
    assert abs(only.cy_pct - 0.5) < 0.05
    # "Huge" — at least 80% of the long side.
    assert only.radius_pct >= 0.8


def test_aurora_returns_three_blobs_at_varying_y():
    spec = derive_mesh_from_palette(PALETTE, "aurora")
    assert len(spec.blobs) == 3
    ys = sorted(b.cy_pct for b in spec.blobs)
    # Top, middle, bottom — y values must be monotonically distinct.
    assert ys[0] < 0.35
    assert 0.35 <= ys[1] <= 0.65
    assert ys[2] > 0.65


def test_orbit_returns_five_blobs():
    spec = derive_mesh_from_palette(PALETTE, "orbit")
    assert len(spec.blobs) == 5


def test_short_palette_wraps_around():
    """A 1-color palette should still produce a full hero spec — colors wrap."""
    spec = derive_mesh_from_palette(["#6366f1"], "hero")
    assert len(spec.blobs) == 4
    assert all(b.color_hex == "6366F1" for b in spec.blobs)


# ---- emit_mesh ---------------------------------------------------------------


def _new_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    return prs, slide


def test_emit_mesh_adds_one_shape_per_blob_with_radial_gradfill():
    prs, slide = _new_slide()
    bbox = BoundingBox(x=0, y=0, w=1280, h=720)
    spec = derive_mesh_from_palette(PALETTE, "hero")

    n_before = len(slide.shapes)
    shapes = emit_mesh(slide, bbox, spec)
    n_after = len(slide.shapes)

    assert len(shapes) == len(spec.blobs)
    assert n_after - n_before == len(spec.blobs)

    for shape in shapes:
        sp_pr = shape._element.spPr
        grad = sp_pr.find(f"{{{NS_A}}}gradFill")
        assert grad is not None, "every blob must use a:gradFill"
        path = grad.find(f"{{{NS_A}}}path")
        assert path is not None and path.get("path") == "circle", (
            "radial direction is encoded as <a:path path='circle'>"
        )


def test_emit_mesh_each_stop_carries_alpha_child():
    prs, slide = _new_slide()
    bbox = BoundingBox(x=0, y=0, w=1280, h=720)
    spec = MeshSpec(
        blobs=[
            GlowBlob(color_hex="6366F1", cx_pct=0.5, cy_pct=0.5, radius_pct=0.5, alpha_max=0.6),
        ]
    )
    shapes = emit_mesh(slide, bbox, spec)
    assert len(shapes) == 1

    sp_pr = shapes[0]._element.spPr
    grad = sp_pr.find(f"{{{NS_A}}}gradFill")
    stops = grad.find(f"{{{NS_A}}}gsLst").findall(f"{{{NS_A}}}gs")
    # Both the center and the edge stop must carry an explicit <a:alpha>:
    # center alpha < 1.0 (alpha_max=0.6) and edge alpha = 0.0.
    for gs in stops:
        srgb = gs.find(f"{{{NS_A}}}srgbClr")
        assert srgb is not None
        assert srgb.get("val") == "6366F1"
        alpha = srgb.find(f"{{{NS_A}}}alpha")
        assert alpha is not None, "translucency requires an explicit a:alpha child"


def test_emit_mesh_empty_spec_is_noop():
    prs, slide = _new_slide()
    bbox = BoundingBox(x=0, y=0, w=1280, h=720)
    shapes = emit_mesh(slide, bbox, MeshSpec())
    assert shapes == []


def test_emit_mesh_round_trips_through_pptx_save_load():
    """Sanity check: a mesh-decorated deck must still open in python-pptx."""
    prs, slide = _new_slide()
    bbox = BoundingBox(x=0, y=0, w=1280, h=720)
    spec = derive_mesh_from_palette(PALETTE, "hero")
    emit_mesh(slide, bbox, spec)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    reopened = Presentation(buf)
    # Slide 1 (after the blank-layout add_slide) still has our blobs.
    reopened_shapes = list(reopened.slides[0].shapes)
    assert len(reopened_shapes) == len(spec.blobs)
