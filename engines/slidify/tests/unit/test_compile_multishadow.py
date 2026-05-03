"""Tests for the multi-shadow stack split (CONTRACT §3.3).

PPTX shape effects support exactly **one** outer + one inner shadow per
shape via `<a:effectLst>`. F3 splits any IR `shadows: [outer, outer, inner,
…]` into:
  * the first outer + first inner → native `<a:outerShdw>` + `<a:innerShdw>`
    on the primary shape
  * remaining shadows → transparent sibling rects with their own native
    shadow.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from slidify.compile_ir import NS_A, compile_ir
from slidify.ir import (
    FillSolid,
    IRBbox,
    IRBoxShadow,
    IRDeck,
    IRShapeNode,
    IRSlide,
)


def _shape_with_shadows(shadows: list[IRBoxShadow]) -> IRShapeNode:
    return IRShapeNode(
        kind="shape",
        recipeId="multi-shadow",
        bbox=IRBbox(x=100, y=100, w=200, h=120),
        shape="rect",
        fill=FillSolid(kind="solid", color="#3b82f6"),
        shadows=shadows,
    )


def _compile(node, tmp_path: Path) -> Presentation:
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / "ms.pptx")
    return Presentation(str(out))


def _count_shadows_on(sp_pr) -> tuple[int, int]:
    """Return (n_outerShdw, n_innerShdw) attached to the spPr's effectLst."""
    if sp_pr is None:
        return (0, 0)
    eff = sp_pr.find(f"{{{NS_A}}}effectLst")
    if eff is None:
        return (0, 0)
    return (
        len(eff.findall(f"{{{NS_A}}}outerShdw")),
        len(eff.findall(f"{{{NS_A}}}innerShdw")),
    )


def test_two_outers_plus_inner_splits_to_one_primary_one_sibling(tmp_path):
    """Per spec: shadows [outer1, outer2, inner1] → primary native (outer1
    + inner1) + ONE sibling rect carrying outer2."""
    node = _shape_with_shadows([
        IRBoxShadow(offsetY=2, blur=8, color="#000000"),       # outer1
        IRBoxShadow(offsetY=12, blur=24, color="#3b82f6"),     # outer2
        IRBoxShadow(offsetY=0, blur=4, color="#ffffff", inset=True),  # inner1
    ])
    prs = _compile(node, tmp_path)
    shapes = [
        sh for sh in prs.slides[0].shapes
        if getattr(sh._element, "spPr", None) is not None
    ]
    # 1 primary + 1 sibling rect (no slide background since fill is white-by-default).
    rect_shapes = [
        sh for sh in shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 1  # MSO_SHAPE.RECTANGLE
    ]
    # We expect at least 2 RECTANGLE shapes: the primary + 1 sibling.
    assert len(rect_shapes) == 2, (
        f"expected 2 rectangles (primary + 1 sibling), got {len(rect_shapes)}"
    )

    # Primary should carry 1 outer + 1 inner natively.
    primary_outer = primary_inner = sibling_outer = sibling_inner = 0
    for sh in rect_shapes:
        n_outer, n_inner = _count_shadows_on(sh._element.spPr)
        if n_outer + n_inner >= 2:
            primary_outer += n_outer
            primary_inner += n_inner
        else:
            sibling_outer += n_outer
            sibling_inner += n_inner
    assert primary_outer == 1 and primary_inner == 1
    assert sibling_outer == 1 and sibling_inner == 0


def test_single_outer_only_uses_native_no_siblings(tmp_path):
    """A 1-shadow stack must NOT create any sibling shapes."""
    node = _shape_with_shadows([
        IRBoxShadow(offsetY=4, blur=8, color="#000000"),
    ])
    prs = _compile(node, tmp_path)
    rect_shapes = [
        sh for sh in prs.slides[0].shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 1
    ]
    assert len(rect_shapes) == 1
    n_outer, n_inner = _count_shadows_on(rect_shapes[0]._element.spPr)
    assert n_outer == 1 and n_inner == 0


def test_two_outer_one_inner_one_inner_caps_to_native_pair_plus_two_siblings(
    tmp_path,
):
    """[outer1, outer2, inner1, inner2] → native (outer1+inner1) + 2 siblings
    (outer2, inner2)."""
    node = _shape_with_shadows([
        IRBoxShadow(offsetY=2, blur=4, color="#000000"),  # outer1
        IRBoxShadow(offsetY=8, blur=16, color="#3333ff"),  # outer2
        IRBoxShadow(offsetY=0, blur=2, color="#ffffff", inset=True),  # inner1
        IRBoxShadow(offsetY=0, blur=4, color="#ff00ff", inset=True),  # inner2
    ])
    prs = _compile(node, tmp_path)
    rect_shapes = [
        sh for sh in prs.slides[0].shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 1
    ]
    assert len(rect_shapes) == 3  # primary + 2 siblings


def test_legacy_shadow_field_still_compiles(tmp_path):
    """Old code that sets `shadow=…` (singular) instead of `shadows=[…]`
    must continue to work (deprecated but valid)."""
    node = IRShapeNode(
        kind="shape",
        recipeId="legacy-shadow",
        bbox=IRBbox(x=10, y=10, w=100, h=80),
        shape="rect",
        fill=FillSolid(kind="solid", color="#222222"),
        shadow=IRBoxShadow(offsetY=4, blur=8, color="#000000"),
    )
    prs = _compile(node, tmp_path)
    rect_shapes = [
        sh for sh in prs.slides[0].shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 1
    ]
    assert len(rect_shapes) == 1
    n_outer, n_inner = _count_shadows_on(rect_shapes[0]._element.spPr)
    assert n_outer == 1 and n_inner == 0


def test_shadows_overrides_legacy_shadow_when_both_present(tmp_path):
    """Per §1.2: if both `shadow` AND `shadows` are set, `shadows` wins."""
    node = IRShapeNode(
        kind="shape",
        recipeId="both",
        bbox=IRBbox(x=10, y=10, w=100, h=80),
        shape="rect",
        fill=FillSolid(kind="solid", color="#222222"),
        shadow=IRBoxShadow(offsetY=4, blur=8, color="#000000"),
        shadows=[
            IRBoxShadow(offsetY=2, blur=2, color="#ff0000"),
            IRBoxShadow(offsetY=8, blur=16, color="#0000ff"),
        ],
    )
    prs = _compile(node, tmp_path)
    rect_shapes = [
        sh for sh in prs.slides[0].shapes
        if getattr(sh, "auto_shape_type", None) is not None
        and sh.auto_shape_type == 1
    ]
    # 2 outers in `shadows` → primary + 1 sibling.
    assert len(rect_shapes) == 2
