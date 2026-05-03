"""Tests for slidify.theme — programmatic clrScheme rewrites and the
deck-color → accent palette derivation.

Coverage:
  * derive_accents_from_elements ranks gradient stops above one-off
    solids and filters out pure black/white.
  * set_theme_accents patches accent1..accent6 in the master theme part
    in-place, surviving a save+reload roundtrip.
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation

from slidify.models import BoundingBox, DomElement
from slidify.theme import (
    NS_A,
    derive_accents_from_elements,
    set_theme_accents,
)


def _el(
    *,
    color: str = "rgb(255, 255, 255)",
    background_color: str = "rgba(0, 0, 0, 0)",
    background_image: str = "none",
) -> DomElement:
    return DomElement(
        id=0,
        parent_id=None,
        depth=0,
        tag="DIV",
        bbox=BoundingBox(x=0, y=0, w=10, h=10),
        color=color,
        background_color=background_color,
        background_image=background_image,
        stable_selector="#x",
    )


def test_derive_accents_returns_top_n_in_order():
    elements = [
        _el(background_color="rgb(99, 102, 241)"),  # indigo-500 ×3
        _el(background_color="rgb(99, 102, 241)"),
        _el(background_color="rgb(99, 102, 241)"),
        _el(color="rgb(236, 72, 153)"),  # pink-500 ×2
        _el(color="rgb(236, 72, 153)"),
        _el(background_color="rgb(52, 211, 153)"),  # green ×1
    ]
    accents = derive_accents_from_elements(elements)
    assert accents[0] == "#6366f1"
    assert accents[1] == "#ec4899"
    assert "#34d399" in accents


def test_derive_accents_weights_gradient_stops_higher_than_solids():
    """A gradient stop appearing once should outrank a solid-fill color
    appearing once, because gradients are designed brand colors."""
    elements = [
        _el(background_color="rgb(50, 50, 50)"),  # solid ×1
        _el(
            background_image=(
                "linear-gradient(135deg, rgb(129, 140, 248), rgb(244, 114, 182))"
            )
        ),  # 2 stops, weight 2 each
    ]
    accents = derive_accents_from_elements(elements)
    # Both gradient stops should rank above the single solid grey
    assert accents[0] in ("#818cf8", "#f472b6")
    assert accents[1] in ("#818cf8", "#f472b6")
    assert accents[0] != accents[1]


def test_derive_accents_skips_pure_black_white_and_neutrals():
    elements = [
        _el(color="rgb(255, 255, 255)") for _ in range(20)
    ] + [
        _el(background_color="rgb(0, 0, 0)") for _ in range(20)
    ] + [
        _el(background_color="rgb(99, 102, 241)"),
    ]
    accents = derive_accents_from_elements(elements)
    assert "#ffffff" not in accents
    assert "#000000" not in accents
    assert accents[0] == "#6366f1"


def test_derive_accents_caps_at_six():
    elements = [
        _el(background_color=f"rgb({r}, 0, 0)") for r in range(1, 21)
    ]
    accents = derive_accents_from_elements(elements)
    assert len(accents) <= 6


def test_set_theme_accents_patches_clrscheme(tmp_path: Path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])

    ok = set_theme_accents(
        prs,
        primary="#112233",
        secondary="#445566",
        accents=["#778899", "#aabbcc", "#ddeeff", "#1a2b3c"],
    )
    assert ok is True

    out = tmp_path / "themed.pptx"
    prs.save(str(out))

    prs2 = Presentation(str(out))
    master_part = prs2.slide_masters[0].part
    theme_part = next(
        rel.target_part
        for rel in master_part.rels.values()
        if "theme" in rel.reltype.lower()
    )
    scheme = etree.fromstring(theme_part.blob).find(
        f".//{{{NS_A}}}themeElements/{{{NS_A}}}clrScheme"
    )
    expected = {
        "accent1": "112233",
        "accent2": "445566",
        "accent3": "778899",
        "accent4": "AABBCC",
        "accent5": "DDEEFF",
        "accent6": "1A2B3C",
    }
    for slot, want in expected.items():
        slot_el = scheme.find(f"{{{NS_A}}}{slot}")
        assert slot_el is not None
        clr = slot_el.find(f"{{{NS_A}}}srgbClr")
        assert clr is not None
        assert clr.get("val").upper() == want.upper()


def test_set_theme_accents_returns_false_when_no_master():
    """Defensive: brand-new Presentation with no slide master shouldn't crash."""

    class FakePrs:
        slide_masters: list = []

    assert set_theme_accents(FakePrs(), primary="#000000") is False
