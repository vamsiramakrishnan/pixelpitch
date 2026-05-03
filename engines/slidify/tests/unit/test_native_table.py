"""Tier-1 routing + emitter coverage for the HTML <table> → PPTX table path.

The DOM walker captures translatable tables into ``DomElement.table_data`` and
``rule_native_table`` routes them to ``DecisionKind.NativeTable``. The emitter
lays a real PPTX ``GraphicFrame`` table — editable cells, not floating
text frames.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from slidify.classifier.tier1 import classify_tier1, rule_native_table
from slidify.emitter import Emitter
from slidify.models import (
    BoundingBox,
    Decision,
    DecisionKind,
    DomElement,
    EmitOp,
    UnitKind,
    VisualUnit,
)


def _table_data(
    rows: list[list[dict]], header_rows: int = 0, n_cols: int | None = None
) -> dict:
    if n_cols is None:
        n_cols = sum(c.get("colspan", 1) for c in rows[0])
    return {"rows": rows, "header_rows": header_rows, "n_cols": n_cols}


def _table_unit(table_data: dict, bbox: BoundingBox | None = None) -> VisualUnit:
    bbox = bbox or BoundingBox(x=0, y=0, w=600, h=200)
    el = DomElement(
        id=1,
        parent_id=None,
        depth=0,
        tag="TABLE",
        bbox=bbox,
        is_table=True,
        table_data=table_data,
    )
    return VisualUnit(id="u_t", kind=UnitKind.Table, bbox=bbox, elements=[el])


def _cell(text: str, **overrides) -> dict:
    base = {
        "text": text,
        "is_header": False,
        "colspan": 1,
        "rowspan": 1,
        "color": "rgb(0,0,0)",
        "background_color": "rgba(0,0,0,0)",
        "font_family": "Inter",
        "font_size": "14px",
        "font_weight": "400",
        "text_align": "start",
        "italic": False,
        "underline": False,
    }
    base.update(overrides)
    return base


def test_rule_native_table_fires_for_translatable_table():
    rows = [
        [_cell("Q", is_header=True), _cell("Revenue", is_header=True)],
        [_cell("Q1"), _cell("$3.2M")],
        [_cell("Q2"), _cell("$3.9M")],
    ]
    unit = _table_unit(_table_data(rows, header_rows=1, n_cols=2))
    d = rule_native_table(unit)
    assert d is not None
    assert d.kind == DecisionKind.NativeTable
    assert d.source_tier == "tier1"
    assert d.metadata["table_element_id"] == 1


def test_rule_native_table_returns_none_when_no_table_data():
    bbox = BoundingBox(x=0, y=0, w=10, h=10)
    el = DomElement(
        id=1, parent_id=None, depth=0, tag="DIV", bbox=bbox, is_table=False
    )
    unit = VisualUnit(id="u", kind=UnitKind.Generic, bbox=bbox, elements=[el])
    assert rule_native_table(unit) is None


def test_classify_tier1_picks_native_table_for_table_unit():
    rows = [[_cell("a"), _cell("b")], [_cell("c"), _cell("d")]]
    unit = _table_unit(_table_data(rows, n_cols=2))
    d = classify_tier1(unit)
    assert d is not None and d.kind == DecisionKind.NativeTable


def test_emitter_writes_a_real_pptx_table(tmp_path: Path):
    rows = [
        [_cell("Quarter", is_header=True), _cell("Revenue", is_header=True)],
        [_cell("Q1"), _cell("$3.2M")],
        [_cell("Q2"), _cell("$3.9M")],
    ]
    unit = _table_unit(
        _table_data(rows, header_rows=1, n_cols=2),
        BoundingBox(x=80, y=120, w=600, h=200),
    )
    op = EmitOp(
        unit_id=unit.id,
        decision=Decision(
            kind=DecisionKind.NativeTable,
            confidence=1.0,
            reason="html_table",
            source_tier="tier1",
        ),
        z_order=0,
        bbox=unit.bbox,
    )
    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    em._emit_native_table(slide, unit, op)
    out = tmp_path / "table.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    s = prs.slides[0]
    tables = [shp for shp in s.shapes if shp.has_table]
    assert len(tables) == 1, "expected exactly one PPTX table primitive"
    t = tables[0].table
    assert len(t.rows) == 3
    assert len(t.columns) == 2
    # Cell text round-trips and header bolding survived.
    assert t.cell(0, 0).text_frame.text == "Quarter"
    assert t.cell(2, 1).text_frame.text == "$3.9M"
    header_run = t.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert header_run.font.bold is True


def test_emitter_handles_colspan(tmp_path: Path):
    rows = [
        [_cell("Title", colspan=2, is_header=True)],
        [_cell("Left"), _cell("Right")],
    ]
    unit = _table_unit(_table_data(rows, header_rows=1, n_cols=2))
    op = EmitOp(
        unit_id=unit.id,
        decision=Decision(
            kind=DecisionKind.NativeTable,
            confidence=1.0,
            reason="html_table",
            source_tier="tier1",
        ),
        z_order=0,
        bbox=unit.bbox,
    )
    em = Emitter()
    slide = em.prs.slides.add_slide(em.prs.slide_layouts[6])
    em._emit_native_table(slide, unit, op)
    out = tmp_path / "merged.pptx"
    em.save(out)
    em.close()

    prs = Presentation(str(out))
    t = next(shp for shp in prs.slides[0].shapes if shp.has_table).table
    assert len(t.rows) == 2 and len(t.columns) == 2
    # The merged origin cell holds the title text, the merged-away cell is
    # marked as merged in PPTX semantics.
    assert t.cell(0, 0).text_frame.text == "Title"
    assert t.cell(0, 0).is_merge_origin
    # Row-2 cells are independent.
    assert t.cell(1, 0).text_frame.text == "Left"
    assert t.cell(1, 1).text_frame.text == "Right"


def test_native_table_counted_in_native_area_ratio():
    from slidify.emitter import native_area_ratio

    rows = [[_cell("a")]]
    unit = _table_unit(
        _table_data(rows, n_cols=1), BoundingBox(x=0, y=0, w=640, h=360)
    )
    op = EmitOp(
        unit_id=unit.id,
        decision=Decision(
            kind=DecisionKind.NativeTable, confidence=1.0, source_tier="tier1"
        ),
        z_order=0,
        bbox=unit.bbox,
    )
    ratio = native_area_ratio([op])
    # 640*360 of 1280*720 == 1/4
    assert ratio == pytest.approx(0.25, abs=1e-3)
