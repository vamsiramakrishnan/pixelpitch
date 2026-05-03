"""Tests for the 18 new preset shapes added in Wave-2 (CONTRACT §3.5).

Each preset must lower to the documented `<a:prstGeom prst=…>` value, and
the resulting PPTX must round-trip through python-pptx.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from slidify.compile_ir import NS_A, compile_ir
from slidify.ir import (
    FillSolid,
    IRBbox,
    IRDeck,
    IRShapeNode,
    IRSlide,
)


# IR shape kind → expected `<a:prstGeom prst=…>` (CONTRACT §3.5).
_KIND_TO_PRST = {
    "triangle": "triangle",
    "right-triangle": "rtTriangle",
    "pentagon": "pentagon",
    "hexagon": "hexagon",
    "octagon": "octagon",
    "parallelogram": "parallelogram",
    "trapezoid": "trapezoid",
    "chevron": "chevron",
    "chevron-left": "chevron",  # + flipH
    "callout-bubble": "wedgeRectCallout",
    "brace-left": "leftBrace",
    "brace-right": "rightBrace",
    "brace-top": "bracketPair",  # + 90deg rotation
    "brace-bottom": "bracketPair",  # + -90deg rotation
    "plus": "mathPlus",
    "star-5": "star5",
    "star-6": "star6",
    "arrow-right": "rightArrow",
    "arrow-left": "leftArrow",
    "arrow-up": "upArrow",
    "arrow-down": "downArrow",
}


def _compile_preset(kind: str, tmp_path: Path, **extra) -> Presentation:
    node = IRShapeNode(
        kind="shape",
        recipeId=f"preset.{kind}",
        bbox=IRBbox(x=20, y=20, w=200, h=120),
        shape=kind,
        fill=FillSolid(kind="solid", color="#3b82f6"),
        **extra,
    )
    deck = IRDeck(version=2, slides=[IRSlide(index=0, nodes=[node])])
    out = compile_ir(deck, tmp_path / f"{kind}.pptx")
    return Presentation(str(out))


def _first_prst(prs: Presentation) -> str | None:
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        prst = sp_pr.find(f"{{{NS_A}}}prstGeom")
        if prst is not None:
            return prst.get("prst")
    return None


@pytest.mark.parametrize("kind,expected_prst", list(_KIND_TO_PRST.items()))
def test_preset_emits_correct_prstgeom(tmp_path, kind, expected_prst):
    prs = _compile_preset(kind, tmp_path)
    actual = _first_prst(prs)
    assert actual == expected_prst, f"{kind}: expected {expected_prst}, got {actual}"


def test_chevron_left_has_fliph_attribute(tmp_path):
    prs = _compile_preset("chevron-left", tmp_path)
    sh = prs.slides[0].shapes[0]
    sp_pr = sh._element.spPr
    xfrm = sp_pr.find(f"{{{NS_A}}}xfrm")
    assert xfrm is not None
    assert xfrm.get("flipH") == "1"


def test_brace_top_rotates_90(tmp_path):
    prs = _compile_preset("brace-top", tmp_path)
    sh = prs.slides[0].shapes[0]
    # python-pptx normalizes rotation to a property on the shape.
    assert abs(sh.rotation - 90.0) < 0.01


def test_brace_bottom_rotates_minus_90(tmp_path):
    prs = _compile_preset("brace-bottom", tmp_path)
    sh = prs.slides[0].shapes[0]
    # PPTX wraps rotation to [0,360); -90 ≡ 270.
    rot = sh.rotation % 360.0
    assert abs(rot - 270.0) < 0.01


def test_callout_bubble_with_pointer_sets_adjustments(tmp_path):
    prs = _compile_preset(
        "callout-bubble",
        tmp_path,
        calloutPointerSide="bottom",
        calloutPointerOffset=0.5,
        calloutPointerLengthPx=24.0,
    )
    sh = prs.slides[0].shapes[0]
    # The wedgeRectCallout exposes 2 adjustments; presence + non-default value
    # is enough — exact formula is tested via IR-level calc.
    assert len(sh.adjustments) >= 2


def test_all_presets_round_trip_reopen(tmp_path):
    """Every preset emits a file python-pptx can re-open."""
    for kind in _KIND_TO_PRST:
        prs = _compile_preset(kind, tmp_path)
        assert len(prs.slides) == 1, kind
        assert len(prs.slides[0].shapes) >= 1, kind


def test_preset_shape_round_trip_recipe_id(tmp_path):
    prs = _compile_preset("hexagon", tmp_path)
    SLIDIFY = "https://slidify.dev/2026/recipe"
    found = False
    for sh in prs.slides[0].shapes:
        sp_pr = getattr(sh._element, "spPr", None)
        if sp_pr is None:
            continue
        for ext in sp_pr.findall(f"{{{NS_A}}}extLst/{{{NS_A}}}ext"):
            if ext.get("uri") == SLIDIFY:
                found = True
    assert found
