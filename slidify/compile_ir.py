"""IR → PPTX compiler.

Consumes a JSON IR document (produced by `@slidify/components`) and emits a
native PPTX. No browser, no clustering, no classifier — every primitive in the
IR has a known emission path. Native ratio is 100% by construction; only
`raster` nodes ever hit a picture.

Each emitted shape carries a `slidify:recipeId` extension list entry so the
reverse path (PPTX → IR → JSX) can recognize what to reconstitute.

Wave-2 additions (CONTRACT §3):
  * _emit_path: PathShape → `<a:custGeom>` (§3.1).
  * _emit_shape: routes the 18 new preset shapes (§3.5).
  * _emit_shape_with_shadows: multi-shadow split (§3.3).
  * clipPath: rounded-rect → ROUNDED_RECTANGLE (picture); path → raster
    fallback (§3.4).
  * mask: alphaModFix-equivalent rendering (§3.4 spirit).
  * Pattern fill → `<a:pattFill>` or tiled ovals (§3.2).
  * Deck.version ∈ {1, 2} (§9.6).
"""

from __future__ import annotations

import asyncio
import base64
import io
import json
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path

import structlog
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from slidify.colors import parse_color
from slidify.effects import apply_shadow_stack
from slidify.fonts import resolve as resolve_font
from slidify.geom import (
    SLIDE_H_EMU,
    SLIDE_H_PX,
    SLIDE_W_EMU,
    SLIDE_W_PX,
    px_to_emu,
)
from slidify.gradients import (
    GradientStop as GGradientStop,
)
from slidify.gradients import (
    LinearGradient as GLinearGradient,
)
from slidify.gradients import (
    RadialGradient as GRadialGradient,
)
from slidify.gradients import (
    apply_gradient_fill,
)
from slidify.ir import (
    Fill,
    FillLinearGradient,
    FillRadialGradient,
    FillSolid,
    IRBbox,
    IRClipPathPath,
    IRClipPathRoundedRect,
    IRDeck,
    IRGroupNode,
    IRMaskLinearGradient,
    IRMaskRadialGradient,
    IRNode,
    IRParagraph,
    IRPathShapeNode,
    IRPatternFill,
    IRPictureNode,
    IRRasterNode,
    IRShapeNode,
    IRSlide,
    IRTextNode,
    _color_to_hex_alpha,
    normalize_shadows,
)
from slidify.path import (
    detect_prst_txwarp,
    make_custgeom_xml,
    path_bbox,
)
from slidify.pattern_fills import apply_pattern_fill
from slidify.shadows import (
    BoxShadow as ShBoxShadow,
)
from slidify.shadows import (
    apply_shadow,
)

log = structlog.get_logger(__name__)


# OOXML namespace + slidify recipe-id extension URI.
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
SLIDIFY_RECIPE_NS = "https://slidify.dev/2026/recipe"
# Distinct extension URI for escape-hatch metering payload — keeps the
# harvester's parsing job trivial: scan for this URI, decode the JSON.
SLIDIFY_ESCAPE_NS = "https://slidify.dev/2026/escape-hatch"

# Default slide canvas (px) — kept in sync with IRSlide.bbox default
# in slidify/ir.py. Used as the denominator for escapeRate when the slide
# itself omits a bbox.
_DEFAULT_SLIDE_W_PX = 1280
_DEFAULT_SLIDE_H_PX = 720

# 1×1 transparent PNG — fallback for environments where Playwright
# isn't available (CI without Chromium, unit tests). The .pptx still
# embeds a picture so the slide stays editable; the harvester reads the
# CSS payload from the metadata extension and re-renders if needed.
_TRANSPARENT_1PX_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
)

_TEXT_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "right": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER,
    "justify": PP_ALIGN.JUSTIFY,
}


# §3.5 IR shape kind → MSO_SHAPE preset. None entries are handled separately
# (e.g. chevron-left is CHEVRON + flipH; brace-top/bottom is DOUBLE_BRACKET
# rotated). Unknown kinds default to RECTANGLE.
_SHAPE_KIND_MAP: dict[str, MSO_SHAPE] = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded-rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "line": MSO_SHAPE.RECTANGLE,
    # Wave-2 additions (§1.7).
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right-triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "chevron": MSO_SHAPE.CHEVRON,
    "chevron-left": MSO_SHAPE.CHEVRON,  # + flipH=1 below
    "callout-bubble": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "brace-left": MSO_SHAPE.LEFT_BRACE,
    "brace-right": MSO_SHAPE.RIGHT_BRACE,
    # OOXML's `bracketPair` preset is exposed as DOUBLE_BRACKET in python-pptx.
    "brace-top": MSO_SHAPE.DOUBLE_BRACKET,  # + rotation 90
    "brace-bottom": MSO_SHAPE.DOUBLE_BRACKET,  # + rotation -90
    "plus": MSO_SHAPE.MATH_PLUS,
    "star-5": MSO_SHAPE.STAR_5_POINT,
    "star-6": MSO_SHAPE.STAR_6_POINT,
    "arrow-right": MSO_SHAPE.RIGHT_ARROW,
    "arrow-left": MSO_SHAPE.LEFT_ARROW,
    "arrow-up": MSO_SHAPE.UP_ARROW,
    "arrow-down": MSO_SHAPE.DOWN_ARROW,
}

# Arrowhead kind/size → OOXML <a:headEnd type=…> attribute.
_ARROWHEAD_TYPE_MAP = {
    "none": "none",
    "arrow": "triangle",
    "dot": "oval",
    "diamond": "diamond",
    "bar": "stealth",  # closest visual to a flat bar terminator
}
_ARROWHEAD_SIZE_MAP = {"sm": "sm", "md": "med", "lg": "lg"}

# IR strokeDasharray → OOXML `<a:prstDash val=…>`. Arbitrary dash patterns
# fall through to `<a:custDash>` emission.
_DASH_PRESET_MAP: dict[tuple[float, ...], str] = {
    (4.0, 2.0): "dash",
    (1.0, 2.0): "dot",
    (6.0, 2.0, 1.0, 2.0): "dashDot",
    (6.0, 2.0, 1.0, 2.0, 1.0, 2.0): "lgDashDotDot",
}


# ---- Escape-hatch metering --------------------------------------------------


@dataclass
class EscapeMetering:
    """Per-deck running totals for `report.escapeRate`.

    Tracks the area share each `chrome.escape-hatch` raster consumed,
    bucketed by its free-text `intent` label. Read by `compile_ir(...)`
    as part of the public return value so callers can write the metering
    block straight into `report.json`.

    Per CONTRACT-v2 §F.4. Fields:

    * `total_slide_area`  — sum of slide.bbox area across the deck.
    * `escape_area`       — sum of escape-hatch raster bbox areas.
    * `area_by_intent`    — escape area summed per intent label.
    * `count_by_intent`   — escape-hatch instance count per intent label.
    """

    total_slide_area: float = 0.0
    escape_area: float = 0.0
    area_by_intent: dict[str, float] = field(default_factory=lambda: defaultdict(float))
    count_by_intent: dict[str, int] = field(default_factory=lambda: defaultdict(int))

    def add_slide(self, slide_w: float, slide_h: float) -> None:
        self.total_slide_area += max(0.0, slide_w) * max(0.0, slide_h)

    def add_escape(self, intent: str, area: float) -> None:
        if area <= 0:
            return
        self.escape_area += area
        self.area_by_intent[intent] += area
        self.count_by_intent[intent] += 1

    def to_report_dict(self) -> dict:
        """Serialize to the `report.escapeRate` shape per CONTRACT-v2 §F.4.

        Includes both `byIntent` (area-share fractions) AND `countByIntent`
        (raw instance counts) so downstream tools — notably
        `slidify report-escape-clusters` — can apply count-based
        thresholds. Without `countByIntent` the CLI's `--threshold`
        flag was a no-op (forced to compare against area share).
        """
        denom = self.total_slide_area
        if denom <= 0:
            return {
                "value": 0.0,
                "byIntent": {},
                "countByIntent": {},
                "atomCandidates": [],
            }
        return {
            "value": self.escape_area / denom,
            "byIntent": {k: v / denom for k, v in sorted(self.area_by_intent.items())},
            "countByIntent": dict(sorted(self.count_by_intent.items())),
            # M4 harvester populates atomCandidates when run later — see
            # CONTRACT-v2 §F.5. Empty during single-deck compile.
            "atomCandidates": [],
        }


# ---- Public entry points -----------------------------------------------------


def compile_ir(
    deck_json: str | dict | IRDeck,
    output_path: str | Path,
    *,
    return_metering: bool = False,
) -> Path | tuple[Path, EscapeMetering]:
    """Compile an IR deck (JSON string, dict, or IRDeck) to PPTX. Returns path.

    When ``return_metering=True`` returns ``(path, EscapeMetering)`` so the
    caller can attach `report.escapeRate` to its conversion report. Default
    return shape is unchanged for backward compatibility.
    """
    if isinstance(deck_json, IRDeck):
        deck = deck_json
    elif isinstance(deck_json, dict):
        deck = IRDeck.model_validate(deck_json)
    else:
        deck = IRDeck.model_validate(json.loads(deck_json))
    out = Path(output_path)
    compiler = _IRCompiler()
    compiler.compile(deck, out)
    if return_metering:
        return out, compiler.metering
    return out


def compile_ir_file(json_path: str | Path, output_path: str | Path) -> Path:
    """Read an IR JSON file and compile to PPTX."""
    text = Path(json_path).read_text(encoding="utf-8")
    return compile_ir(text, output_path)


# ---- Compiler ----------------------------------------------------------------


class _IRCompiler:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W_EMU)
        self.prs.slide_height = Emu(SLIDE_H_EMU)
        # Per-deck escape-hatch metering, surfaced via `report.escapeRate`.
        self.metering = EscapeMetering()

    def compile(self, deck: IRDeck, out_path: Path) -> None:
        # Deck.version is constrained at the model level (Literal[1, 2]); a
        # raised ValidationError on parse already enforces §9.6 lockstep.
        for slide_ir in deck.slides:
            self._emit_slide(slide_ir)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out_path))
        log.info(
            "compile_ir.done",
            path=str(out_path),
            n_slides=len(deck.slides),
            escape_rate=(
                self.metering.escape_area / self.metering.total_slide_area
                if self.metering.total_slide_area
                else 0.0
            ),
        )

    def _emit_slide(self, slide_ir: IRSlide) -> None:
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)

        # Slide background: emit as a full-bleed shape behind everything.
        if not _is_white_solid(slide_ir.background):
            self._emit_background_shape(slide, slide_ir.background)

        # Track slide area for the escape-rate denominator BEFORE any node
        # emission — gives accurate accounting even on slides with zero
        # escape-hatch usage (the denominator still grows).
        slide_w = slide_ir.bbox.w if slide_ir.bbox else _DEFAULT_SLIDE_W_PX
        slide_h = slide_ir.bbox.h if slide_ir.bbox else _DEFAULT_SLIDE_H_PX
        self.metering.add_slide(slide_w, slide_h)

        for node in slide_ir.nodes:
            self._emit_node(slide, node)

        if slide_ir.notes:
            try:
                slide.notes_slide.notes_text_frame.text = slide_ir.notes
            except Exception as e:
                log.warning("compile_ir.notes_failed", error=str(e))

    # ---- Node dispatch -------------------------------------------------------

    def _emit_node(self, slide, node: IRNode) -> None:
        if isinstance(node, IRTextNode):
            self._emit_text(slide, node)
        elif isinstance(node, IRShapeNode):
            self._emit_shape(slide, node)
        elif isinstance(node, IRPictureNode):
            self._emit_picture(slide, node)
        elif isinstance(node, IRRasterNode):
            self._emit_raster(slide, node)
        elif isinstance(node, IRPathShapeNode):
            self._emit_path(slide, node)
        elif isinstance(node, IRGroupNode):
            for child in node.children:
                self._emit_node(slide, child)

    # ---- Per-node emit -------------------------------------------------------

    def _emit_text(self, slide, node: IRTextNode) -> None:
        bbox = node.bbox or IRBbox(x=0, y=0, w=SLIDE_W_PX, h=80)
        x, y, w, h = _emu_rect(bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        if node.fill is not None:
            _apply_fill(tb, node.fill, slide=slide, bbox=bbox)
        # Multi-shadow path (§3.3) — handles both legacy `shadow` and new `shadows`.
        shadows = normalize_shadows(node)
        if shadows:
            apply_shadow_stack(tb, shadows, slide=slide, parent_bbox=bbox)

        # Text-on-path stretch goal (§9.1): map to prstTxWarp when possible.
        if node.onPath is not None:
            _apply_text_warp(tb, node.onPath)

        first = True
        for para in node.paragraphs:
            self._emit_paragraph(tf, para, first)
            first = False
        _stamp_recipe_id(tb, node.recipeId)

    def _emit_paragraph(self, tf, para: IRParagraph, first: bool) -> None:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = _TEXT_ALIGN_MAP.get(para.align, PP_ALIGN.LEFT)
        for run_spec in para.runs:
            run = p.add_run()
            run.text = run_spec.text
            font = run.font
            if run_spec.fontFamily:
                font.name = resolve_font(run_spec.fontFamily)
            else:
                font.name = "Inter"
            if run_spec.fontSizePx:
                font.size = Pt(max(8.0, run_spec.fontSizePx * 72.0 / 96.0))
            if run_spec.fontWeight is not None:
                font.bold = run_spec.fontWeight >= 600
            font.italic = run_spec.italic
            font.underline = run_spec.underline
            if run_spec.color is not None:
                hex_, _alpha = _color_to_hex_alpha(run_spec.color)
                col = parse_color(hex_)
                if col is not None:
                    try:
                        font.color.rgb = col[0]
                    except Exception:
                        pass

    def _emit_shape(self, slide, node: IRShapeNode) -> None:
        bbox = node.bbox or IRBbox(x=0, y=0, w=100, h=100)
        x, y, w, h = _emu_rect(bbox)

        # Resolve preset; rounded-rect with radius=0 collapses to plain rect.
        kind = node.shape
        preset = _SHAPE_KIND_MAP.get(kind, MSO_SHAPE.RECTANGLE)
        if kind == "rounded-rect" and node.borderRadiusPx <= 0:
            preset = MSO_SHAPE.RECTANGLE

        shape = slide.shapes.add_shape(preset, x, y, w, h)
        shape.line.fill.background()

        # Special-case post-creation transforms.
        if kind == "chevron-left":
            _apply_flip_h(shape)
        elif kind == "brace-top":
            shape.rotation = 90.0
        elif kind == "brace-bottom":
            shape.rotation = -90.0

        # callout-bubble pointer geometry (§3.5) — adj1/adj2 normalize the
        # pointer tip from (offset, length) into OOXML adjustment values.
        if (
            kind == "callout-bubble"
            and node.calloutPointerSide is not None
        ):
            _apply_callout_adjustments(
                shape,
                node.calloutPointerSide,
                node.calloutPointerOffset,
                node.calloutPointerLengthPx,
                bbox,
            )

        # Fill: pattern fills require slide+bbox for the dot tiler.
        if isinstance(node.fill, IRPatternFill):
            apply_pattern_fill(
                shape, node.fill, slide=slide, bbox=bbox
            )
        else:
            _apply_fill(shape, node.fill, slide=slide, bbox=bbox)

        if node.border is not None:
            self._apply_border(shape, node.border)

        # Multi-shadow split (§3.3).
        shadows = normalize_shadows(node)
        if shadows:
            apply_shadow_stack(shape, shadows, slide=slide, parent_bbox=bbox)

        # ClipPath routing (§3.4) — rounded-rect on a non-picture node logs
        # and ignores; path-clip falls back to raster (stub for Wave-2).
        if node.clipPath is not None:
            self._route_clip_path(shape, node, slide)

        _stamp_recipe_id(shape, node.recipeId)

    def _apply_border(self, shape, border) -> None:
        hex_, _alpha = _color_to_hex_alpha(border.color)
        col = parse_color(hex_)
        try:
            shape.line.width = Emu(px_to_emu(border.width))
            if col is not None:
                shape.line.color.rgb = col[0]
        except Exception:
            pass

    def _emit_picture(self, slide, node: IRPictureNode) -> None:
        bbox = node.bbox or IRBbox(x=0, y=0, w=SLIDE_W_PX, h=SLIDE_H_PX)
        x, y, w, h = _emu_rect(bbox)

        # Rounded-rect clip on a picture: emit as a ROUNDED_RECTANGLE
        # filled with the picture (§3.4).
        if isinstance(node.clipPath, IRClipPathRoundedRect):
            return self._emit_picture_in_rounded_rect(
                slide, node, bbox, node.clipPath
            )

        # Path-shaped clipPath on a picture: emit a freeform shape
        # (custGeom from the clip's path) and picture-fill it. This is
        # the proper §3.4 path-clip implementation for pictures, the
        # raster fallback only kicks in for ShapeNode (where the host
        # shape's geometry would also have to be intersected against
        # the clip path — left for Wave-3).
        if isinstance(node.clipPath, IRClipPathPath):
            return self._emit_picture_in_path_clip(
                slide, node, bbox, node.clipPath
            )

        try:
            data = _fetch_picture(node.src)
        except Exception as e:
            log.warning("compile_ir.picture_fetch_failed", src=node.src[:80], error=str(e))
            return
        pic = slide.shapes.add_picture(io.BytesIO(data), x, y, w, h)

        # Image mask (§3.4 spirit). Wave-2 ships a soft fallback: emit a
        # full-bbox shape with a matching gradient fill on top of the
        # picture. This is not a true alpha-mask but is good enough for
        # gradient-fade edges; a real `<a:alphaModFix>` chain requires
        # patching the embedded blip relationship and is deferred.
        if node.mask is not None:
            self._emit_picture_mask_overlay(slide, bbox, node.mask)

        _stamp_recipe_id(pic, node.recipeId)

    def _emit_picture_in_rounded_rect(
        self,
        slide,
        node: IRPictureNode,
        bbox: IRBbox,
        clip: IRClipPathRoundedRect,
    ) -> None:
        """Render picture inside a ROUNDED_RECTANGLE shape via picture-fill."""
        # Fetch BEFORE adding the auto-shape: if the source is missing or
        # unreachable we want the slide to skip the picture entirely (matching
        # `_emit_picture`'s non-clipped path), not be left with an orphan
        # rounded rectangle that has no fill.
        try:
            data = _fetch_picture(node.src)
        except Exception as e:
            log.warning(
                "compile_ir.picture_fetch_failed",
                src=node.src[:80],
                error=str(e),
            )
            return

        # Apply inset to the bbox first (clip is relative to the node bbox).
        inset = max(0.0, clip.insetPx)
        inner = IRBbox(
            x=bbox.x + inset,
            y=bbox.y + inset,
            w=max(1.0, bbox.w - 2 * inset),
            h=max(1.0, bbox.h - 2 * inset),
        )
        x, y, w, h = _emu_rect(inner)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h
        )
        shape.line.fill.background()

        # Adjust the corner radius. Rounded-rect adj1 is in 1/100,000 of the
        # min(w,h)/2 — mapping radiusPx → adj1 keeps the visual radius.
        try:
            ref = max(1.0, min(inner.w, inner.h) / 2.0)
            adj_norm = max(0.0, min(0.5, clip.radiusPx / ref / 2.0))
            shape.adjustments[0] = adj_norm
        except Exception:
            pass
        # Replace solid fill with picture fill via lxml.
        try:
            sp_pr = shape._element.spPr
            for tag in ("solidFill", "gradFill", "pattFill", "noFill", "blipFill"):
                for existing in sp_pr.findall(f"{{{NS_A}}}{tag}"):
                    sp_pr.remove(existing)
            # Embed the image via the slide's part so the relationship is
            # created and python-pptx tracks the embedded blob correctly.
            _image_part, rel_id = slide.part.get_or_add_image_part(io.BytesIO(data))
            blip_fill = etree.SubElement(sp_pr, f"{{{NS_A}}}blipFill")
            sp_pr.remove(blip_fill)
            sp_pr.insert(0, blip_fill)
            r_ns = (
                "http://schemas.openxmlformats.org/officeDocument/"
                "2006/relationships"
            )
            etree.SubElement(
                blip_fill,
                f"{{{NS_A}}}blip",
                attrib={f"{{{r_ns}}}embed": rel_id},
            )
            etree.SubElement(blip_fill, f"{{{NS_A}}}stretch").append(
                etree.Element(f"{{{NS_A}}}fillRect")
            )
        except Exception as e:
            log.warning("compile_ir.picture_fill_failed", error=str(e))

        if node.mask is not None:
            self._emit_picture_mask_overlay(slide, inner, node.mask)
        _stamp_recipe_id(shape, node.recipeId)

    def _emit_picture_in_path_clip(
        self,
        slide,
        node: IRPictureNode,
        bbox: IRBbox,
        clip: IRClipPathPath,
    ) -> None:
        """Render picture inside an arbitrary-path clip via custGeom + blipFill.

        Strategy mirrors `_emit_picture_in_rounded_rect`: emit a freeform
        shape whose geometry is the clip path, then patch its fill with
        a `<a:blipFill>` carrying the picture. PowerPoint clips the
        picture to the shape's path geometry natively — no raster needed.

        Fetches the picture before adding any shape so a missing source
        leaves no orphan auto-shape (matches `_emit_picture_in_rounded_rect`).
        """
        try:
            data = _fetch_picture(node.src)
        except Exception as e:
            log.warning(
                "compile_ir.picture_fetch_failed",
                src=node.src[:80],
                error=str(e),
            )
            return

        # Anchor the freeform shape at the picture's bbox (`node.bbox`) so
        # the picture is placed/scaled exactly as it would be without any
        # clip. Project the path's slide-pixel commands into that same
        # local-space reference — coordinates outside `node.bbox` simply
        # extend past the shape and get clipped by PowerPoint.
        #
        # The earlier draft sized the shape from `path_bbox(clip.commands)`,
        # which reprojected the picture into the clip's tight bounds. That
        # squashed the image whenever the clip path didn't fill `node.bbox`
        # (e.g. inset clips, off-center keyholes, paths that extend beyond
        # the picture). Anchoring to node.bbox keeps placement stable
        # regardless of clip geometry.
        node_box = (bbox.x, bbox.y, bbox.x + bbox.w, bbox.y + bbox.h)
        x, y, w, h = _emu_rect(bbox)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.line.fill.background()

        # Replace prstGeom with custGeom built from the clip path,
        # projected against node.bbox so the path renders inside the
        # picture's natural placement.
        try:
            sp_pr = shape._element.spPr
            for prst in sp_pr.findall(f"{{{NS_A}}}prstGeom"):
                sp_pr.remove(prst)
            for cg in sp_pr.findall(f"{{{NS_A}}}custGeom"):
                sp_pr.remove(cg)
            cust_xml = make_custgeom_xml(clip.commands, node_box)
            wrapper = f'<root xmlns:a="{NS_A}">{cust_xml}</root>'
            parsed = etree.fromstring(wrapper)
            cust = parsed[0]
            xfrm_idx = 0
            for i, child in enumerate(sp_pr):
                if child.tag == f"{{{NS_A}}}xfrm":
                    xfrm_idx = i + 1
                    break
            sp_pr.insert(xfrm_idx, cust)
        except Exception as e:
            log.warning("compile_ir.path_clip_custgeom_failed", error=str(e))

        # Replace solidFill with blipFill carrying the picture.
        try:
            sp_pr = shape._element.spPr
            for tag in (
                "solidFill", "gradFill", "pattFill", "noFill", "blipFill"
            ):
                for existing in sp_pr.findall(f"{{{NS_A}}}{tag}"):
                    sp_pr.remove(existing)
            _image_part, rel_id = slide.part.get_or_add_image_part(io.BytesIO(data))
            blip_fill = etree.SubElement(sp_pr, f"{{{NS_A}}}blipFill")
            sp_pr.remove(blip_fill)
            sp_pr.insert(0, blip_fill)
            r_ns = (
                "http://schemas.openxmlformats.org/officeDocument/"
                "2006/relationships"
            )
            etree.SubElement(
                blip_fill,
                f"{{{NS_A}}}blip",
                attrib={f"{{{r_ns}}}embed": rel_id},
            )
            etree.SubElement(blip_fill, f"{{{NS_A}}}stretch").append(
                etree.Element(f"{{{NS_A}}}fillRect")
            )
        except Exception as e:
            log.warning("compile_ir.path_clip_blipfill_failed", error=str(e))

        if node.mask is not None:
            self._emit_picture_mask_overlay(slide, bbox, node.mask)
        _stamp_recipe_id(shape, node.recipeId)

    def _emit_picture_mask_overlay(self, slide, bbox: IRBbox, mask) -> None:
        """Wave-2 fallback for picture masks: paint a gradient overlay rect.

        Per CONTRACT §3.4 spirit, a true `<a:alphaModFix>` chain is the
        target; the gradient overlay is good enough for fade-to-black
        edges (the most common use). The overlay sits *above* the picture
        so the visible result blends correctly with the slide background.
        Documented deviation: this is NOT a true alpha mask.
        """
        x, y, w, h = _emu_rect(bbox)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        try:
            rect.line.fill.background()
        except Exception:
            pass
        # Convert mask stops into gradient stops, with the 'alpha' mapped
        # to a soft black overlay (alpha=stop.alpha → black at full alpha).
        stops = [
            GGradientStop(
                color_hex="000000",
                alpha=max(0.0, min(1.0, 1.0 - s.alpha)),
                position=max(0.0, min(1.0, s.position)),
            )
            for s in mask.stops
        ]
        if isinstance(mask, IRMaskLinearGradient):
            grad = GLinearGradient(angle_deg=mask.angleDeg, stops=stops)
        elif isinstance(mask, IRMaskRadialGradient):
            grad = GRadialGradient(stops=stops, cx=mask.cx, cy=mask.cy)
        else:
            return
        try:
            apply_gradient_fill(rect, grad, densify=False)
        except Exception as e:
            log.warning("compile_ir.mask_overlay_failed", error=str(e))

    def _emit_raster(self, slide, node: IRRasterNode) -> None:
        bbox = node.bbox or IRBbox(x=0, y=0, w=SLIDE_W_PX, h=SLIDE_H_PX)
        # Escape-hatch detour: detect, render-if-empty, embed, stamp metadata,
        # and bookkeep for `report.escapeRate`. The detection key is
        # `metadata.role == 'escape-hatch'` (set by EscapeHatch's TSX
        # emitter); we also accept the recipeId as a fallback so handcrafted
        # IR fixtures don't have to remember the metadata convention.
        is_escape = (
            node.metadata.get("role") == "escape-hatch"
            or node.recipeId == "chrome.escape-hatch"
        )
        if is_escape:
            self._emit_escape_hatch(slide, node, bbox)
            return

        x, y, w, h = _emu_rect(bbox)
        try:
            png_bytes = base64.b64decode(node.pngBase64)
        except Exception as e:
            log.warning("compile_ir.raster_decode_failed", error=str(e))
            return
        pic = slide.shapes.add_picture(io.BytesIO(png_bytes), x, y, w, h)
        _stamp_recipe_id(pic, node.recipeId)

    def _emit_path(self, slide, node: IRPathShapeNode) -> None:
        """Emit an IRPathShapeNode as a freeform shape with `<a:custGeom>`.

        Strategy (§3.1):
          1. Compute path bbox (from IR or `path_bbox(commands)`).
          2. Add a freeform shape, then replace its <a:prstGeom> with
             <a:custGeom> built from `commands`.
          3. Apply fill (solid/gradient/pattern) and stroke.
          4. Map dash pattern → <a:prstDash> when standard, else custDash.
          5. Attach arrowheads on <a:ln>.
          6. Run multi-shadow stack.
          7. Stamp recipeId.
        """
        # 1. Resolve bbox.
        if node.bbox is not None:
            bbox = node.bbox
            box = (bbox.x, bbox.y, bbox.x + bbox.w, bbox.y + bbox.h)
        else:
            box = path_bbox(node.commands)
            bbox = IRBbox(
                x=box[0], y=box[1], w=max(1.0, box[2] - box[0]), h=max(1.0, box[3] - box[1])
            )
            box = (bbox.x, bbox.y, bbox.x + bbox.w, bbox.y + bbox.h)
        x, y, w, h = _emu_rect(bbox)

        # 2. Add a placeholder rectangle, then swap its geometry for custGeom.
        # We build the custGeom XML as a string (path.make_custgeom_xml uses
        # the `a:` prefix) and parse it inside an inline-namespaced wrapper
        # so lxml binds NS_A correctly.
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        try:
            sp_pr = shape._element.spPr
            for prst in sp_pr.findall(f"{{{NS_A}}}prstGeom"):
                sp_pr.remove(prst)
            for cg in sp_pr.findall(f"{{{NS_A}}}custGeom"):
                sp_pr.remove(cg)
            cust_xml = make_custgeom_xml(node.commands, box)
            wrapper = f'<root xmlns:a="{NS_A}">{cust_xml}</root>'
            parsed = etree.fromstring(wrapper)
            cust = parsed[0]
            # OOXML schema: <a:xfrm> precedes geometry inside spPr; insert
            # custGeom right after xfrm so reopen-tools see the expected order.
            xfrm_idx = 0
            for i, child in enumerate(sp_pr):
                if child.tag == f"{{{NS_A}}}xfrm":
                    xfrm_idx = i + 1
                    break
            sp_pr.insert(xfrm_idx, cust)
        except Exception as e:
            log.warning("compile_ir.path_custgeom_failed", error=str(e))

        # 3. Fill (or no-fill for stroke-only paths).
        if node.fill is not None:
            if isinstance(node.fill, IRPatternFill):
                apply_pattern_fill(
                    shape, node.fill, slide=slide, bbox=bbox
                )
            else:
                _apply_fill(shape, node.fill, slide=slide, bbox=bbox)
        else:
            try:
                shape.fill.background()
            except Exception:
                pass

        # 4. Stroke (width, color, dash, cap, join) — stroke-only paths skip fill.
        if node.strokeWidthPx > 0:
            try:
                shape.line.width = Emu(px_to_emu(node.strokeWidthPx))
                if node.strokeColor is not None:
                    hex_, _ = _color_to_hex_alpha(node.strokeColor)
                    col = parse_color(hex_)
                    if col is not None:
                        shape.line.color.rgb = col[0]
            except Exception:
                pass
            if node.strokeDasharray:
                _apply_dash(shape, node.strokeDasharray)
            # Cap/join were defined in IR but never written until now;
            # skipping them silently rounded square corners on connectors.
            _apply_line_cap_join(shape, node.strokeLinecap, node.strokeLinejoin)

        # 5. Arrowheads — markerStart → headEnd, markerEnd → tailEnd.
        if node.markerStart is not None or node.markerEnd is not None:
            _apply_arrowheads(shape, node.markerStart, node.markerEnd)

        # 6. Multi-shadow.
        shadows = normalize_shadows(node)
        if shadows:
            apply_shadow_stack(shape, shadows, slide=slide, parent_bbox=bbox)

        # 7. Recipe id.
        _stamp_recipe_id(shape, node.recipeId)

    def _route_clip_path(self, shape, node: IRShapeNode, slide) -> None:
        """Per §3.4: clipPath on a non-picture node.

        Rounded-rect clips are ignored on shapes/text (the shape's own
        radius is the canonical clip). Path clips fall back to raster.
        """
        if isinstance(node.clipPath, IRClipPathRoundedRect):
            log.info(
                "compile_ir.clip_path_rounded_rect_ignored_on_shape",
                recipe_id=node.recipeId,
            )
        elif isinstance(node.clipPath, IRClipPathPath):
            # Path clip on a ShapeNode would require intersecting the
            # shape's preset geometry with the clip path — non-trivial
            # polygon math left for a Wave-3 follow-up. Surface the
            # gap loudly (warning, not info) and stamp a metadata
            # extension so the harvester sees the unimplemented case
            # rather than the deck silently rendering un-clipped.
            log.warning(
                "compile_ir.clip_path_path_unimplemented_on_shape",
                recipe_id=node.recipeId,
                note=(
                    "ShapeNode + path-clip is not yet supported; the "
                    "shape renders without the path clip. PictureNode + "
                    "path-clip works (custGeom + blipFill)."
                ),
            )
            try:
                sp_pr = shape._element.spPr
                ext_lst = sp_pr.find(f"{{{NS_A}}}extLst")
                if ext_lst is None:
                    ext_lst = etree.SubElement(sp_pr, f"{{{NS_A}}}extLst")
                ext = etree.SubElement(
                    ext_lst,
                    f"{{{NS_A}}}ext",
                    attrib={"uri": "https://slidify.dev/2026/clip-fallback"},
                )
                ext.text = "path-clip-not-implemented-on-shape"
            except Exception:
                pass

    def _emit_escape_hatch(
        self, slide, node: IRRasterNode, bbox: IRBbox
    ) -> None:
        """Embed an escape-hatch raster and stamp its metering payload.

        Pipeline:
          1. Decode the PNG (or render via Chromium if `pngBase64` is empty).
          2. Embed as a picture at the requested bbox.
          3. Stamp the recipe-id extension (so PPTX→IR round-trip recognizes it).
          4. Stamp a separate `escape-hatch` extension carrying the cssPayload,
             intent, attempted, etc., so the harvester can find clusters.
          5. Update `self.metering` for `report.escapeRate`.
        """
        x, y, w, h = _emu_rect(bbox)
        css_payload = str(node.metadata.get("cssPayload", ""))
        body = str(node.metadata.get("body", ""))
        intent = str(node.metadata.get("intent", "unspecified"))
        attempted = node.metadata.get("attempted")

        # 1. Source the PNG. Prefer the IR-supplied pngBase64; fall back to a
        # Chromium screenshot of the cssPayload; final fallback is a 1x1
        # transparent placeholder so the slide still embeds an editable
        # picture frame for the harvester to find later.
        png_bytes: bytes
        if node.pngBase64:
            try:
                png_bytes = base64.b64decode(node.pngBase64)
            except Exception as e:
                log.warning(
                    "compile_ir.escape_hatch_decode_failed", error=str(e)
                )
                png_bytes = base64.b64decode(_TRANSPARENT_1PX_PNG_B64)
        else:
            png_bytes = _render_escape_payload(
                css_payload=css_payload,
                body=body,
                width_px=int(max(1, bbox.w)),
                height_px=int(max(1, bbox.h)),
            )

        try:
            pic = slide.shapes.add_picture(io.BytesIO(png_bytes), x, y, w, h)
        except Exception as e:
            log.warning("compile_ir.escape_hatch_picture_failed", error=str(e))
            return

        # 2. Stamps for round-trip + harvester discovery.
        _stamp_recipe_id(pic, node.recipeId)
        _stamp_escape_metadata(
            pic,
            css_payload=css_payload,
            body=body,
            intent=intent,
            attempted=attempted if isinstance(attempted, str) else None,
        )

        # 3. Metering. Area is in slide-coordinate pixels (matches the
        # denominator added in `_emit_slide`).
        area = max(0.0, bbox.w) * max(0.0, bbox.h)
        self.metering.add_escape(intent, area)
        log.info(
            "compile_ir.escape_hatch_emitted",
            intent=intent,
            attempted=attempted,
            bbox_area=area,
            had_payload=bool(css_payload),
        )

    def _emit_background_shape(self, slide, fill: Fill) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0),
            Emu(0),
            Emu(SLIDE_W_EMU),
            Emu(SLIDE_H_EMU),
        )
        shape.line.fill.background()
        bg_bbox = IRBbox(x=0, y=0, w=SLIDE_W_PX, h=SLIDE_H_PX)
        if isinstance(fill, IRPatternFill):
            apply_pattern_fill(shape, fill, slide=slide, bbox=bg_bbox)
        else:
            _apply_fill(shape, fill, slide=slide, bbox=bg_bbox)
        _stamp_recipe_id(shape, "slide-background")
        # Send to back so children render on top. python-pptx doesn't expose
        # z-order helpers cleanly; the background is added first, so it's
        # naturally at the bottom.


# ---- Helpers ----------------------------------------------------------------


def _is_white_solid(fill: Fill) -> bool:
    if not isinstance(fill, FillSolid):
        return False
    hex_, _ = _color_to_hex_alpha(fill.color)
    return hex_.lower() in ("#ffffff", "#fff")


def _emu_rect(b: IRBbox) -> tuple[Emu, Emu, Emu, Emu]:
    return (
        Emu(px_to_emu(b.x)),
        Emu(px_to_emu(b.y)),
        Emu(px_to_emu(b.w)),
        Emu(px_to_emu(b.h)),
    )


def _apply_fill(shape, fill: Fill, *, slide=None, bbox: IRBbox | None = None) -> None:
    if isinstance(fill, FillSolid):
        hex_, alpha = _color_to_hex_alpha(fill.color)
        col = parse_color(hex_)
        if col is None:
            return
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = col[0]
        except Exception:
            pass
        # python-pptx's RGBColor doesn't carry alpha. If the IR specified a
        # translucent solid (e.g. `{hex: "#ffffff", alpha: 0.5}`), the
        # opacity was previously discarded silently. Patch the solidFill
        # XML directly so the rendered shape matches the IR.
        if alpha < 0.999:
            try:
                sp_pr = shape._element.spPr
                solid = sp_pr.find(f"{{{NS_A}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{NS_A}}}srgbClr")
                    if srgb is not None:
                        for existing in srgb.findall(f"{{{NS_A}}}alpha"):
                            srgb.remove(existing)
                        etree.SubElement(
                            srgb,
                            f"{{{NS_A}}}alpha",
                            attrib={"val": str(int(round(alpha * 100_000)))},
                        )
            except Exception as e:
                log.debug("compile_ir.solid_alpha_failed", error=str(e))
        return
    if isinstance(fill, FillLinearGradient):
        grad = GLinearGradient(
            angle_deg=fill.angleDeg,
            stops=[_to_grad_stop(s) for s in fill.stops],
        )
        apply_gradient_fill(shape, grad)
        return
    if isinstance(fill, FillRadialGradient):
        grad = GRadialGradient(
            stops=[_to_grad_stop(s) for s in fill.stops],
            cx=fill.cx,
            cy=fill.cy,
            shape=fill.shape,
        )
        apply_gradient_fill(shape, grad)
        return
    if isinstance(fill, IRPatternFill):
        # Pattern fills route via apply_pattern_fill but need slide+bbox for
        # the dot tiler. _apply_fill is called from contexts that may not
        # have those handy (e.g. text fills); silently fall through.
        if slide is not None and bbox is not None:
            apply_pattern_fill(shape, fill, slide=slide, bbox=bbox)
        return
    # FillNone
    try:
        shape.fill.background()
    except Exception:
        pass


def _to_grad_stop(s) -> GGradientStop:
    hex_, alpha = _color_to_hex_alpha(s.color)
    # GradientStop wants 6-hex (no leading #).
    bare = hex_[1:] if hex_.startswith("#") else hex_
    return GGradientStop(color_hex=bare.upper(), alpha=alpha, position=s.position)


def _apply_shadow_ir(shape, sh) -> None:
    """Legacy helper for single shadow (kept for back-compat)."""
    hex_, alpha = _color_to_hex_alpha(sh.color)
    bare = hex_[1:] if hex_.startswith("#") else hex_
    apply_shadow(
        shape,
        ShBoxShadow(
            inset=sh.inset,
            offset_x_px=sh.offsetX,
            offset_y_px=sh.offsetY,
            blur_px=max(0.0, sh.blur),
            spread_px=sh.spread,
            color_hex=bare.upper(),
            alpha=alpha,
        ),
    )


def _stamp_escape_metadata(
    shape,
    *,
    css_payload: str,
    body: str,
    intent: str,
    attempted: str | None,
) -> None:
    """Embed the escape-hatch payload as an OOXML extension entry.

    Distinct URI (`SLIDIFY_ESCAPE_NS`) so the harvester's grep can be a
    single substring match; the body is JSON so it stays compact and
    schema-stable across compiler versions.
    """
    try:
        sp_pr = shape._element.spPr
    except AttributeError:
        return
    if sp_pr is None:
        return
    ext_lst = sp_pr.find(f"{{{NS_A}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(sp_pr, f"{{{NS_A}}}extLst")
    # Replace any prior escape-hatch ext for idempotent re-emission.
    for existing in ext_lst.findall(
        f"{{{NS_A}}}ext[@uri='{SLIDIFY_ESCAPE_NS}']"
    ):
        ext_lst.remove(existing)
    ext = etree.SubElement(
        ext_lst, f"{{{NS_A}}}ext", attrib={"uri": SLIDIFY_ESCAPE_NS}
    )
    payload = etree.SubElement(
        ext, f"{{{SLIDIFY_ESCAPE_NS}}}escapePayload"
    )
    payload.text = json.dumps(
        {
            "intent": intent,
            "attempted": attempted,
            "cssPayload": css_payload,
            "body": body,
        },
        sort_keys=True,
    )


def _render_escape_payload(
    *, css_payload: str, body: str, width_px: int, height_px: int
) -> bytes:
    """Render the escape-hatch CSS payload to PNG via Chromium.

    Best-effort: if Playwright isn't available (CI without browser deps,
    unit tests), returns a 1×1 transparent PNG placeholder so the picture
    frame still embeds and the harvester can later re-render from the
    metadata extension. The metering is unaffected — area accounting uses
    the bbox, not the PNG dimensions.
    """
    html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>html,body{margin:0;padding:0;background:transparent;}"
        f".__esc{{width:{width_px}px;height:{height_px}px;{css_payload}}}"
        "</style></head><body>"
        f"<div class='__esc'>{body}</div>"
        "</body></html>"
    )
    try:
        return asyncio.run(
            _screenshot_html(html, width_px=width_px, height_px=height_px)
        )
    except RuntimeError:
        # We're already inside an event loop (e.g. compile_ir called from an
        # async context). Fall through to placeholder; an async-aware caller
        # should pre-render the PNG and ship it in pngBase64.
        log.info("compile_ir.escape_payload_in_loop_fallback")
        return base64.b64decode(_TRANSPARENT_1PX_PNG_B64)
    except Exception as e:
        log.warning("compile_ir.escape_payload_render_failed", error=str(e))
        return base64.b64decode(_TRANSPARENT_1PX_PNG_B64)


async def _screenshot_html(html: str, *, width_px: int, height_px: int) -> bytes:
    """One-shot Chromium screenshot of an HTML string. Used for escape-hatch
    rasterization. Imported lazily so unit tests without Playwright don't
    fail at import time."""
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        log.info("compile_ir.playwright_unavailable")
        return base64.b64decode(_TRANSPARENT_1PX_PNG_B64)

    async with async_playwright() as pw:
        browser = await pw.chromium.launch()
        try:
            ctx = await browser.new_context(
                viewport={"width": max(1, width_px), "height": max(1, height_px)},
                device_scale_factor=1,
            )
            page = await ctx.new_page()
            await page.set_content(html, wait_until="load", timeout=10_000)
            try:
                await page.evaluate("document.fonts && document.fonts.ready")
            except Exception:
                pass
            png = await page.screenshot(
                clip={
                    "x": 0,
                    "y": 0,
                    "width": max(1, width_px),
                    "height": max(1, height_px),
                },
                full_page=False,
                type="png",
                omit_background=True,
            )
            return png
        finally:
            await browser.close()


def _stamp_recipe_id(shape, recipe_id: str) -> None:
    """Tag an emitted shape with its IR recipe id via an extension entry.

    The reverse path (PPTX → IR) reads this back to recognize what
    component produced the shape, enabling round-trip JSX reconstitution.
    """
    try:
        sp_pr = shape._element.spPr
    except AttributeError:
        return
    if sp_pr is None:
        return
    # Find or create extLst
    ext_lst = sp_pr.find(f"{{{NS_A}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(sp_pr, f"{{{NS_A}}}extLst")
    # One ext per recipe-id (overwrite if already present)
    for existing in ext_lst.findall(f"{{{NS_A}}}ext[@uri='{SLIDIFY_RECIPE_NS}']"):
        ext_lst.remove(existing)
    ext = etree.SubElement(
        ext_lst, f"{{{NS_A}}}ext", attrib={"uri": SLIDIFY_RECIPE_NS}
    )
    rid = etree.SubElement(ext, f"{{{SLIDIFY_RECIPE_NS}}}recipeId")
    rid.text = recipe_id


def _fetch_picture(src: str) -> bytes:
    """Fetch picture bytes from data: URI or http(s) URL."""
    if src.startswith("data:"):
        head, _, payload = src.partition(",")
        if ";base64" in head:
            return base64.b64decode(payload)
        # RFC 2397: a data URI without ";base64" is URL-encoded ASCII.
        # Encoding the literal string would keep `%20` etc. as bytes and
        # produce broken image payloads. Decode percent-escapes to the
        # original bytes before returning.
        from urllib.parse import unquote_to_bytes

        return unquote_to_bytes(payload)
    if src.startswith(("http://", "https://")):
        import httpx

        with httpx.Client(timeout=10.0, follow_redirects=True) as c:
            r = c.get(src)
            r.raise_for_status()
            return r.content
    return Path(src).read_bytes()


def _apply_flip_h(shape) -> None:
    """Set flipH=1 on a shape's <a:xfrm> (used for chevron-left)."""
    try:
        sp_pr = shape._element.spPr
        xfrm = sp_pr.find(f"{{{NS_A}}}xfrm")
        if xfrm is None:
            return
        xfrm.set("flipH", "1")
    except Exception:
        pass


def _apply_callout_adjustments(
    shape,
    side: str,
    offset: float | None,
    length_px: float | None,
    bbox: IRBbox,
) -> None:
    """Set the wedgeRectCallout adjustment values from §1.7 fields.

    OOXML's `wedgeRectCallout` exposes two adjustments (adj1, adj2) that
    place the pointer tip relative to the shape's origin (proportions of
    -50%..+50% of width/height, expressed in 1/1000-fraction space:
    -50000..+50000). Here we interpret IR fields as: the pointer base
    sits at `(offset along side)` and the tip extends `length_px` outward.
    """
    if offset is None and length_px is None:
        return
    o = 0.5 if offset is None else max(0.0, min(1.0, float(offset)))
    length = max(0.0, float(length_px or 24.0))

    # Compute tip in the shape's local 0..1 box, with center=(0.5,0.5).
    if side == "top":
        tip_x = o
        tip_y = -length / max(1.0, bbox.h)
    elif side == "bottom":
        tip_x = o
        tip_y = 1.0 + length / max(1.0, bbox.h)
    elif side == "left":
        tip_x = -length / max(1.0, bbox.w)
        tip_y = o
    else:  # right
        tip_x = 1.0 + length / max(1.0, bbox.w)
        tip_y = o

    # Translate to OOXML's center-relative system (adjustments are signed
    # fractions of the bbox where 0 = center; positive = right/down).
    adj1 = (tip_x - 0.5)
    adj2 = (tip_y - 0.5)
    try:
        shape.adjustments[0] = adj1
        shape.adjustments[1] = adj2
    except Exception:
        pass


def _apply_dash(shape, pattern: list[float]) -> None:
    """Translate IR strokeDasharray → `<a:prstDash>` or `<a:custDash>`."""
    try:
        sp_pr = shape._element.spPr
        ln = sp_pr.find(f"{{{NS_A}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, f"{{{NS_A}}}ln")
        # Strip any existing prstDash/custDash
        for tag in ("prstDash", "custDash"):
            for existing in ln.findall(f"{{{NS_A}}}{tag}"):
                ln.remove(existing)
        key = tuple(float(x) for x in pattern)
        # Per SVG/CSS spec: an odd-length dash array repeats once to
        # become even, so `[3, 2, 1]` is equivalent to `[3, 2, 1, 3, 2, 1]`.
        # The previous loop used `range(0, len(key) - 1, 2)` and silently
        # dropped the trailing value, changing the stroke rhythm.
        if len(key) % 2 == 1:
            key = key + key
        prst = _DASH_PRESET_MAP.get(key)
        if prst is not None:
            etree.SubElement(ln, f"{{{NS_A}}}prstDash", attrib={"val": prst})
        else:
            cust = etree.SubElement(ln, f"{{{NS_A}}}custDash")
            # Pairs: (dash, gap) → `<a:ds d="…" sp="…"/>`
            for i in range(0, len(key), 2):
                etree.SubElement(
                    cust,
                    f"{{{NS_A}}}ds",
                    attrib={
                        "d": str(int(round(key[i] * 100_000))),
                        "sp": str(int(round(key[i + 1] * 100_000))),
                    },
                )
    except Exception as e:
        log.debug("compile_ir.dash_apply_failed", error=str(e))


_LINECAP_MAP = {"butt": "flat", "round": "rnd", "square": "sq"}
_LINEJOIN_TAGS = {"miter": "miter", "round": "round", "bevel": "bevel"}


def _apply_line_cap_join(shape, cap: str | None, join: str | None) -> None:
    """Write <a:ln cap="…"> attribute and <a:miter|round|bevel/> child.

    PathShapeNode declared strokeLinecap/strokeLinejoin from day one, but
    `_emit_path` only wrote width/color/dash — so authored 'round' or
    'square' caps and 'round'/'bevel' joins were silently lost.
    """
    if cap is None and join is None:
        return
    try:
        sp_pr = shape._element.spPr
        ln = sp_pr.find(f"{{{NS_A}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, f"{{{NS_A}}}ln")
        if cap is not None and cap in _LINECAP_MAP:
            ln.set("cap", _LINECAP_MAP[cap])
        if join is not None and join in _LINEJOIN_TAGS:
            # Drop any existing join element so we own this slot.
            for tag in _LINEJOIN_TAGS.values():
                for existing in ln.findall(f"{{{NS_A}}}{tag}"):
                    ln.remove(existing)
            etree.SubElement(ln, f"{{{NS_A}}}{_LINEJOIN_TAGS[join]}")
    except Exception as e:
        log.debug("compile_ir.line_cap_join_failed", error=str(e))


def _apply_arrowheads(shape, marker_start, marker_end) -> None:
    """Attach <a:headEnd>/<a:tailEnd> children to shape's <a:ln>."""
    try:
        sp_pr = shape._element.spPr
        ln = sp_pr.find(f"{{{NS_A}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, f"{{{NS_A}}}ln")

        def _emit(tag: str, marker) -> None:
            if marker is None or marker.kind == "none":
                return
            for existing in ln.findall(f"{{{NS_A}}}{tag}"):
                ln.remove(existing)
            attrs = {"type": _ARROWHEAD_TYPE_MAP.get(marker.kind, "triangle")}
            sz = _ARROWHEAD_SIZE_MAP.get(marker.size, "med")
            attrs["w"] = sz
            attrs["len"] = sz
            etree.SubElement(ln, f"{{{NS_A}}}{tag}", attrib=attrs)

        _emit("headEnd", marker_start)
        _emit("tailEnd", marker_end)
    except Exception as e:
        log.debug("compile_ir.arrowhead_apply_failed", error=str(e))


def _apply_text_warp(shape, on_path) -> None:
    """Attach an `<a:prstTxWarp prst=…>` to a textbox's bodyPr.

    Falls back to straight text when the path doesn't map to a known
    preset (raster fallback per §9.1 is a stub for Wave-2). Logs at
    `warning` level so the lossy fallback is visible by default — users
    were previously seeing "straight text" without knowing they
    requested a curve.
    """
    preset = detect_prst_txwarp(on_path.commands)
    if preset is None:
        log.warning(
            "compile_ir.text_on_path_warp_fallback",
            n_commands=len(on_path.commands),
            note="arbitrary onPath paths fall back to straight text "
                 "until §9.1 raster fallback ships",
        )
        return
    try:
        body_pr = shape._element.txBody.bodyPr
        for existing in body_pr.findall(f"{{{NS_A}}}prstTxWarp"):
            body_pr.remove(existing)
        warp = etree.SubElement(
            body_pr, f"{{{NS_A}}}prstTxWarp", attrib={"prst": preset}
        )
        # avLst is required; default-empty is fine.
        etree.SubElement(warp, f"{{{NS_A}}}avLst")
    except Exception as e:
        log.debug("compile_ir.text_warp_failed", error=str(e))
