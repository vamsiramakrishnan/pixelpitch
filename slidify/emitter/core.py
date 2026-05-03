"""PPTX emitter: turns EmitOps into a python-pptx Presentation."""

from __future__ import annotations

import io
import re
from pathlib import Path
from urllib.parse import urlparse

import httpx
import structlog
from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from slidify.colors import parse_color
from slidify.decorations import derive_decorations
from slidify.emitter.geometry import (
    _apply_rotation,
    _clamp_bbox,
    _emu_rect,
    _emu_rect_offcanvas,
    _grow_bbox_for_fallback_wrap,
    _rotation_degrees,
)
from slidify.emitter.metrics import native_area_ratio
from slidify.emitter.shapes import _shape_kind_for_anchor
from slidify.emitter.text import (
    _are_spatially_distinct,
    _has_title_sized_text,
    _is_bg_clip_text,
    _resolve_run_color,
    _try_apply_gradient_text_fill,
    _union_line_box,
)
from slidify.exceptions import EmitError
from slidify.fonts import is_bold
from slidify.fonts import resolve as resolve_font
from slidify.geom import (
    SLIDE_H_EMU,
    SLIDE_H_PX,
    SLIDE_W_EMU,
    SLIDE_W_PX,
    parse_pt,
    parse_px,
    px_to_emu,
)
from slidify.gradients import apply_gradient_fill, parse_gradient
from slidify.models import (
    BoundingBox,
    DecisionKind,
    DomElement,
    EmitOp,
    RenderedSlide,
    VisualUnit,
)
from slidify.shadows import apply_shadows, parse_box_shadows
from slidify.svg.shapes import emit_svg_shapes
from slidify.text_metrics import (
    compute_font_scale_for_textbox,
    genre_for_family,
    shrink_pill_bbox_to_fit_text,
)

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _apply_explicit_autofit(
    text_frame, font_scale_pct: int, line_space_reduction_pct: int
) -> None:
    """Bake an explicit ``<a:normAutofit fontScale=... lnSpcReduction=.../>``
    into the textframe, replacing any existing autofit element. Pre-computing
    the shrink at emit time makes the deck render identically across
    LibreOffice/PowerPoint regardless of which font gets substituted."""
    from lxml import etree as _et

    bodyPr = text_frame._txBody.bodyPr
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for existing in bodyPr.findall(f"{{{_A_NS}}}{tag}"):
            bodyPr.remove(existing)
    norm = _et.SubElement(bodyPr, f"{{{_A_NS}}}normAutofit")
    norm.set("fontScale", str(font_scale_pct))
    norm.set("lnSpcReduction", str(line_space_reduction_pct))


def _ensure_default_normautofit(text_frame) -> None:
    """Default-on safety net: emit a parameter-free ``<a:normAutofit/>`` on
    every text frame that doesn't already have an explicit autofit element.

    With this, a textbox whose body text exceeds its frame at runtime
    (because the slide author over-budgeted, the font fell back to a wider
    family, or both) shrinks instead of clipping. The compile-time
    overflow detector still surfaces the design bug; this is the runtime
    rescue so the PPTX never displays clipped text in slideshow mode.

    Skipped if any autofit child element is already present — callers
    that have *measured* the right shrink (see ``_apply_explicit_autofit``)
    win, because their value is more accurate than PowerPoint's heuristic.
    """
    bodyPr = text_frame._txBody.bodyPr
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        if bodyPr.find(f"{{{_A_NS}}}{tag}") is not None:
            return
    from lxml import etree as _et

    _et.SubElement(bodyPr, f"{{{_A_NS}}}normAutofit")

log = structlog.get_logger(__name__)

__all__ = [
    "Emitter",
    "_apply_explicit_autofit",
    "_grow_bbox_for_fallback_wrap",
    "_resolve_run_color",
    "_rotation_degrees",
    "_try_apply_gradient_text_fill",
    "_union_line_box",
    "native_area_ratio",
]


_TEXT_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "right": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER,
    "justify": PP_ALIGN.JUSTIFY,
    "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT,
}


class Emitter:
    """Builds a PPTX from rendered slides + per-slide EmitOps."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W_EMU)
        self.prs.slide_height = Emu(SLIDE_H_EMU)
        self._image_cache: dict[str, bytes] = {}
        self._http = httpx.Client(timeout=10.0, follow_redirects=True)
        # Brand palette used by decoration layers (mesh glow, etc.) when the
        # source HTML opts in via `data-slidify-decorate`. The API populates
        # this from the deck-wide color harvest right before save.
        self._brand_palette: list[str] = []

    def set_brand_palette(self, palette: list[str]) -> None:
        """Set the deck-wide brand palette used by decoration layers."""
        self._brand_palette = list(palette)

    def close(self) -> None:
        self._http.close()

    # ----- public --------------------------------------------------------

    async def emit_slide(
        self,
        slide_index: int,
        rendered: RenderedSlide,
        units_by_id: dict[str, VisualUnit],
        ops: list[EmitOp],
        renderer,
        notes: str = "",
    ) -> None:
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)

        if rendered.degraded:
            log.warning("emitter.degraded_slide", slide_index=slide_index, reason=rendered.reason)
            # Fall back to ground-truth raster as a safety net.
            self._emit_full_raster(slide, rendered.ground_truth_png)
            if notes:
                self._set_notes(slide, notes)
            return

        # Sort ops by z_order (already pre-order; but EmitOps may include explicit z)
        ops_sorted = sorted(ops, key=lambda o: o.z_order)
        for op in ops_sorted:
            try:
                await self._emit_op(slide, op, units_by_id, rendered, renderer)
            except Exception as e:
                log.warning(
                    "emitter.op_failed",
                    unit_id=op.unit_id,
                    kind=op.decision.kind,
                    error=str(e),
                )
                # Last-resort: fall back to ground-truth crop for this region.
                try:
                    self._emit_region_raster(
                        slide, rendered.ground_truth_png, op.bbox
                    )
                except Exception as e2:
                    log.error("emitter.fallback_failed", error=str(e2))

        if notes:
            self._set_notes(slide, notes)

    def save(self, path: Path | str) -> None:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._sanitize_for_repair_dialog()
        self.prs.save(str(path))

    def _sanitize_for_repair_dialog(self) -> None:
        """Walk every slide's XML and fix the two patterns that most often
        trigger PowerPoint's "found a problem with some content" repair
        dialog at open time:

          1. Empty paragraph `<a:p/>` with no child runs and no
             `<a:endParaRPr>` — PowerPoint logs an error then quietly
             repairs by inserting one. We insert it ourselves so the
             open is silent.
          2. Run text `<a:t>foo bar</a:t>` containing leading/trailing
             whitespace without `xml:space="preserve"` — XML's default
             whitespace handling strips it, and PowerPoint repairs by
             re-asserting `preserve`.

        Both fixes are idempotent and additive — they only add elements
        / attributes that are missing.
        """
        from lxml import etree as _et

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"
        for slide in self.prs.slides:
            tree = slide._element
            for p in tree.iter(f"{{{ns_a}}}p"):
                # 1. Insert <a:endParaRPr> on empty paragraphs.
                has_run = any(
                    child.tag in (f"{{{ns_a}}}r", f"{{{ns_a}}}fld", f"{{{ns_a}}}br")
                    for child in p
                )
                has_end_pr = p.find(f"{{{ns_a}}}endParaRPr") is not None
                if not has_run and not has_end_pr:
                    _et.SubElement(
                        p,
                        f"{{{ns_a}}}endParaRPr",
                        attrib={"lang": "en-US"},
                    )
                # 2. xml:space="preserve" on whitespace-bearing <a:t>.
                for t in p.iter(f"{{{ns_a}}}t"):
                    if t.text and (
                        t.text != t.text.strip() or "  " in t.text
                    ):
                        if t.get(XML_SPACE) != "preserve":
                            t.set(XML_SPACE, "preserve")

    # ----- per-op dispatch ----------------------------------------------

    async def _emit_op(
        self,
        slide,
        op: EmitOp,
        units_by_id: dict[str, VisualUnit],
        rendered: RenderedSlide,
        renderer,
    ) -> None:
        unit = units_by_id.get(op.unit_id)
        if unit is None:
            return
        kind = op.decision.kind

        if kind == DecisionKind.Skip:
            return

        # Skip shapes whose bbox is entirely off the slide canvas. Without
        # this, source HTML that overflows (CSS overflow:hidden hides it
        # in the browser; the bboxes still get captured) emits dozens of
        # h=1 placeholder shapes clamped to y=719. They aren't visible
        # but bloat shape count and confuse the oracle's region attribution.
        b = op.bbox
        if (
            b.y >= SLIDE_H_PX
            or b.x >= SLIDE_W_PX
            or b.y + b.h <= 0
            or b.x + b.w <= 0
        ):
            return

        if kind == DecisionKind.NativeText:
            self._emit_native_text(slide, unit, op)
            return

        if kind == DecisionKind.NativeBullet:
            self._emit_native_bullet(slide, unit, op)
            return

        if kind == DecisionKind.NativeShape:
            self._emit_native_shape(slide, unit, op)
            return

        if kind == DecisionKind.NativePicture:
            self._emit_native_picture(slide, unit, op)
            return

        if kind == DecisionKind.NativeSvg:
            self._emit_native_svg(slide, unit, op)
            return

        if kind == DecisionKind.NativeTable:
            self._emit_native_table(slide, unit, op)
            return

        if kind == DecisionKind.Raster:
            await self._emit_raster(slide, unit, op, rendered, renderer)
            return

        if kind == DecisionKind.Hybrid:
            # Surgical hybrid emission, in order of preference:
            #   1) Try emitting a native gradient/shape (works when the unit's
            #      decoration is fully translatable — gradient + solid + shadow).
            #   2) Else if the renderer captured a no-text decoration-only pass,
            #      crop from THAT (pixel-exact decoration with no text bleed).
            #   3) Else fall back to live region screenshot or ground-truth crop.
            # Children emit independently on top in all three branches.
            if self._try_emit_native_decoration(slide, unit, op):
                return
            if rendered.no_text_png:
                self._emit_region_raster(slide, rendered.no_text_png, op.bbox)
                return
            await self._emit_raster(slide, unit, op, rendered, renderer)
            return

    # ----- native text ---------------------------------------------------

    def _pick_text_anchor(self, unit: VisualUnit) -> DomElement | None:
        """Return the deepest text-bearing element (leaf or text container)."""
        text_elems = [
            e
            for e in unit.all_elements()
            if (e.text and e.text.strip())
            or (e.runs and any(r.text.strip() for r in e.runs if not r.is_break))
            or (getattr(e, "mixed_content_text", None) or "").strip()
        ]
        if not text_elems:
            return None
        for e in text_elems:
            if e.pptx_text:
                return e
        return max(text_elems, key=lambda e: parse_px(e.font_size))

    def _emit_native_text(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        anchor = self._pick_text_anchor(unit)
        if anchor is None:
            return
        # Browser-derived line boxes (`Range.getClientRects()`) give the
        # exact pixel rectangles each text run occupied in the source
        # render. When we have them, prefer the UNION bbox of all line
        # boxes over the unit's bbox — the unit bbox is the CSS box-model
        # container, which is often padded or oversized vs. the actual
        # text. The union is the truth: where the type really sat.
        #
        # Exception: when the anchor carries `mixed_content_text` (parent
        # text alongside block children), the line-box union pulls in the
        # children's positions and would render the parent's text frame
        # ON TOP of the child unit. Use the anchor's own bbox so the
        # parent's text sits in its intended slot above the children.
        if (getattr(anchor, "mixed_content_text", None) or "").strip():
            emit_bbox = anchor.bbox
        else:
            emit_bbox = _union_line_box(unit) or op.bbox

        # Give multi-run inline text containers some horizontal slack so font
        # substitution doesn't wrap them. CSS sizes flex-item widths by their
        # content; PPTX uses fixed bboxes, so a "slidify · v1.0" inline-flex
        # container that fit in 102px in Chromium overflows in Calibri. Grow
        # the bbox by 30% horizontally — only for short multi-run frames where
        # the wrap risk is real. SKIP this when we have line boxes — those
        # are already the truth.
        skip_slack = _union_line_box(unit) is not None
        # Naked inline text only — if the anchor paints its own background
        # (a pill, a chip, a trend tag) the slack would stretch that fill
        # past the text and leave a visible "tail" rectangle.
        anchor_paints_bg = (
            (anchor.background_color and anchor.background_color != "rgba(0, 0, 0, 0)")
            or (anchor.background_image and anchor.background_image != "none")
        )
        if (
            anchor.is_text_container
            and anchor.runs
            and len(anchor.runs) >= 2
            and op.bbox.h <= 48
            and op.bbox.w <= 320
            and not anchor_paints_bg
            and not skip_slack
        ):
            extra = op.bbox.w * 0.30
            emit_bbox = BoundingBox(
                x=op.bbox.x,
                y=op.bbox.y,
                w=op.bbox.w + extra,
                h=op.bbox.h,
            )

        # Pill shrink — kills the "green tail" past the substituted-font
        # text in single-line bg-painted containers. Fires for any
        # narrow single-line element with its own text and a painted
        # background, whether the bg is a solid color OR a gradient
        # (the slide-12 CEO badge has a 90deg gradient bg, which the
        # is_text_container path missed because leaf text elements
        # aren't marked as text containers).
        if (
            anchor_paints_bg
            and (not anchor.runs or len(anchor.runs) <= 1)
            and emit_bbox.h <= 48
            and (anchor.text or "").strip()
        ):
            try:
                pt_size = max(8.0, parse_pt(anchor.font_size))
                emit_bbox = shrink_pill_bbox_to_fit_text(
                    anchor.text.strip(), emit_bbox, pt_size
                )
            except Exception:
                pass

        x, y, w, h = _emu_rect(emit_bbox)
        # Decoration layer (BELOW): mesh glow / drop shadow attached to
        # the text container's anchor. Opt-in via `data-slidify-decorate`.
        deco_stack = derive_decorations(anchor, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, emit_bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        _apply_rotation(tb, anchor)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        _ensure_default_normautofit(tf)
        # PowerPoint substitutes Inter/etc with Calibri when the source font
        # isn't installed — Calibri is wider, so big headlines that fit in 3
        # lines in Chromium can overflow into the lede below in PowerPoint.
        # TEXT_TO_FIT_SHAPE shrinks the font just enough to keep everything
        # in-frame. We only apply it for *title-sized* text where the wrap
        # risk is real; on small frames (pills, badges, kickers) LibreOffice
        # over-shrinks the text to invisibility.
        # Heuristic: apply when the largest run is >= 20px (subtitle-ish or
        # bigger). Small inline labels are unaffected.
        if _has_title_sized_text(unit):
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

        # If unit has a background that we can render, paint it on the textbox.
        # Special case: when this looks like a `background-clip: text` recipe
        # (transparent text on a gradient bg), DON'T paint the gradient as
        # the shape's fill — that produces a gradient rectangle with the
        # text rendered as a transparent silhouette inside it. Instead,
        # leave the textbox transparent; the run-color resolver will
        # substitute the gradient's first stop as a solid text color.
        unit_anchor_el = unit.elements[0] if unit.elements else None
        if unit_anchor_el is not None and not _is_bg_clip_text(unit_anchor_el):
            self._apply_fill(tb, unit_anchor_el)
            self._apply_shadow(tb, unit_anchor_el)

        text_elems = [
            e
            for e in unit.all_elements()
            if (e.text and e.text.strip())
            or (e.runs and any(r.text.strip() for r in e.runs if not r.is_break))
            or (getattr(e, "mixed_content_text", None) or "").strip()
        ]
        # Leaf-text-wins: when a text-bearing parent has text-bearing
        # descendants in the same list, the descendants will paint the
        # words at their precise child bboxes — so emitting the parent
        # too produces a stacked duplicate (slide 8's `.author` div case
        # contained `.name` + `.title`, both with their own text).
        text_elem_ids = {e.id for e in text_elems}
        all_in_unit = unit.all_elements()
        parent_of = {e.id: e.parent_id for e in all_in_unit}

        def _has_text_descendant(eid: int) -> bool:
            for other in all_in_unit:
                if other.id == eid or other.id not in text_elem_ids:
                    continue
                cur = parent_of.get(other.id)
                while cur is not None:
                    if cur == eid:
                        return True
                    cur = parent_of.get(cur)
            return False

        text_elems = [e for e in text_elems if not _has_text_descendant(e.id)]
        para_targets = text_elems if len(text_elems) > 1 else [anchor]

        # Spatially-distinct text emit: when para_targets are arranged
        # horizontally (e.g. process steps with `position:absolute; left:N`),
        # their text MUST land at each element's own x — not stacked into
        # a single tall textbox at the unit's bbox. Detect by checking if
        # any two elements have non-overlapping x-ranges. When detected,
        # bypass the single-textbox path and emit each element as its own
        # native_text shape.
        if (
            len(para_targets) >= 2
            and _are_spatially_distinct(para_targets)
        ):
            for e in para_targets:
                self._emit_text_element_as_own_shape(slide, e, op)
            # Decorations on the unit-level still fire ABOVE.
            if not deco_stack.is_empty():
                deco_stack.emit_above(slide, emit_bbox)
            return

        # Pre-compute fontScale + wrap policy. Single-line textboxes get
        # `wrap=none` so the wider fallback font doesn't reflow text into
        # a phantom second line; multi-paragraph textboxes keep wrapping.
        font_scale = 1.0
        try:
            scale_runs: list[tuple[str, float]] = []
            for e in para_targets:
                t = (e.pptx_text or e.text or "").strip()
                if not t:
                    continue
                scale_runs.append((t, max(8.0, parse_pt(e.font_size))))
            if scale_runs:
                font_scale = self._apply_textframe_wrap_and_scale(
                    tf, scale_runs, emit_bbox.w, emit_bbox.h,
                    genre=genre_for_family(anchor.font_family),
                )
        except Exception:
            pass
        first = True
        for e in para_targets:
            # Build paragraph(s) from runs if available, else from text.
            paragraphs = self._element_to_paragraphs(e)
            for para_runs in paragraphs:
                if not para_runs:
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = _TEXT_ALIGN_MAP.get(e.text_align, PP_ALIGN.LEFT)
                for run_spec in para_runs:
                    self._add_styled_run(
                        p, run_spec, fallback_el=e, font_scale=font_scale
                    )

        # Decoration layer (ABOVE): rim highlight / hairline / inset glow.
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, emit_bbox)

    def _emit_text_element_as_own_shape(
        self, slide, e: DomElement, op: EmitOp
    ) -> None:
        """Emit a single text element as a self-contained textbox at its
        own bbox. Used for spatially-distinct siblings inside one unit
        (e.g. absolutely-positioned process-step labels)."""
        x, y, w, h = _emu_rect(e.bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        _apply_rotation(tb, e)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        _ensure_default_normautofit(tf)
        # Per-element font-metrics + wrap policy.
        font_scale = 1.0
        try:
            text = (e.pptx_text or e.text or "").strip()
            if text:
                pt_size = max(8.0, parse_pt(e.font_size))
                font_scale = self._apply_textframe_wrap_and_scale(
                    tf, [(text, pt_size)], e.bbox.w, e.bbox.h,
                    genre=genre_for_family(e.font_family),
                )
        except Exception:
            pass
        paragraphs = self._element_to_paragraphs(e)
        first = True
        for para_runs in paragraphs:
            if not para_runs:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = _TEXT_ALIGN_MAP.get(e.text_align, PP_ALIGN.LEFT)
            for run_spec in para_runs:
                self._add_styled_run(
                    p, run_spec, fallback_el=e, font_scale=font_scale
                )

    def _apply_textframe_wrap_and_scale(
        self,
        tf,
        runs: list[tuple[str, float]],
        bbox_w_px: float,
        bbox_h_px: float,
        genre: str | None = None,
    ) -> float:
        """Set wrap policy + explicit fontScale for a textframe.

        Single-line textboxes (height <= 1.5 line-heights of the largest
        run) get ``wrap=none``: even if the fallback font happens to be
        slightly wider than fonttools predicts, the text won't reflow
        into a phantom second line that overlaps the next vertical row.

        Returns the font_scale (1.0 = no shrink, 0.91 = 9% shrink) the
        caller should multiply each per-run pt size by. We return a bake-
        scale (instead of always emitting ``<a:normAutofit fontScale=>``)
        because LibreOffice **silently ignores** ``normAutofit``'s
        ``fontScale`` when ``bodyPr.wrap="none"`` — the headline overflows
        the slide and the renderer clips at the canvas edge instead of
        shrinking to fit (slide-13 "Pro, in the smallest possible frame.").
        Baking the scale into ``<a:rPr sz=...>`` works in every renderer.

        ``genre`` ("serif" / "mono" / "sans") routes the autofit
        measurement to a genre-matched fallback font — without this a
        Source Serif Pro headline measured against Liberation Sans
        under-counts by ~15% and skips the shrink that would have kept
        it in-bbox.

        For multi-line wrapping textboxes we keep emitting ``normAutofit`` so
        PowerPoint's runtime shrink-to-fit still applies in the wrapped case.
        """
        if not runs:
            return 1.0
        max_pt = max(pt for _, pt in runs) if runs else 12.0
        # Approximate line-height factor; CSS default is 1.2.
        line_height_px = max_pt * 1.333 * 1.2
        is_single_line = bbox_h_px <= line_height_px * 1.5

        if is_single_line:
            try:
                bodyPr = tf._txBody.bodyPr
                bodyPr.set("wrap", "none")
            except Exception:
                pass

        try:
            scale = compute_font_scale_for_textbox(
                runs, bbox_w_px=bbox_w_px, bbox_h_px=bbox_h_px, genre=genre
            )
            if not scale.needs_autofit:
                return 1.0
            if is_single_line:
                # Bake into rPr — normAutofit is dropped by LibreOffice when
                # wrap=none, and would conflict with the per-run sz anyway.
                return scale.font_scale_pct / 100_000.0
            # Multi-line: normAutofit is honored; let the renderer shrink.
            _apply_explicit_autofit(
                tf, scale.font_scale_pct, scale.line_space_reduction_pct
            )
            return 1.0
        except Exception:
            return 1.0

    def _element_to_paragraphs(
        self, e: DomElement
    ) -> list[list[dict]]:
        """Turn an element's runs (or text) into a list of paragraphs, where
        each paragraph is a list of run-spec dicts. <br> splits paragraphs.
        """
        if e.runs:
            paragraphs: list[list[dict]] = [[]]
            for r in e.runs:
                if r.is_break:
                    paragraphs.append([])
                    continue
                spec = {
                    "text": r.text,
                    "font_family": r.font_family,
                    "font_size": r.font_size,
                    "font_weight": r.font_weight,
                    "color": r.color,
                    # Forward the run's parent background-image so the emitter
                    # can substitute a solid color when the run is gradient-
                    # clipped (color: transparent + background-clip: text).
                    "background_image": r.background_image,
                    "italic": r.italic,
                    "underline": r.underline,
                }
                if not r.text.strip():
                    # Preserve inter-run whitespace; skip empties.
                    if r.text and paragraphs[-1]:
                        paragraphs[-1].append(spec)
                    continue
                paragraphs[-1].append(spec)
            # Drop trailing empty paragraphs
            while paragraphs and not paragraphs[-1]:
                paragraphs.pop()
            return paragraphs
        # `mixed_content_text` is the parent's own direct-text fragment when
        # an element ALSO has block descendants (`<div>parent text<div
        # class="sub">child</div></div>`). The walker captures it on the
        # parent and the children classify themselves; without this fallback
        # the parent's text would silently disappear because `text` is only
        # set on leaf-text and text-container elements.
        text = (
            e.pptx_text or e.text or getattr(e, "mixed_content_text", None) or ""
        ).strip()
        if not text:
            return []
        return [
            [
                {
                    "text": text,
                    "font_family": e.font_family,
                    "font_size": e.font_size,
                    "font_weight": e.font_weight,
                    "color": e.color,
                    "italic": False,
                    "underline": False,
                }
            ]
        ]

    def _add_styled_run(
        self,
        paragraph,
        spec: dict,
        fallback_el: DomElement,
        font_scale: float = 1.0,
    ) -> None:
        run = paragraph.add_run()
        run.text = spec["text"]
        font = run.font
        font.name = resolve_font(spec.get("font_family") or fallback_el.font_family)
        size_px = spec.get("font_size") or fallback_el.font_size
        # Bake font_scale into the per-run pt size when the textframe asked
        # us to (single-line `wrap=none` boxes — see
        # `_apply_textframe_wrap_and_scale` for the rationale). We never
        # shrink below 8pt to keep small footers / kickers readable; the
        # scale is meant for headline-class text where 5–10% shrink avoids
        # canvas-edge clipping under the wider fallback font.
        pt = parse_pt(size_px) * max(0.1, font_scale)
        font.size = Pt(max(8.0, pt))
        font.bold = is_bold(spec.get("font_weight") or fallback_el.font_weight)
        font.italic = bool(spec.get("italic"))
        font.underline = bool(spec.get("underline"))
        # If this is a `background-clip: text` run with a translatable
        # gradient on the source span, emit a true OOXML gradient text fill
        # via python-pptx's Font.fill API. PPTX renders the gradient through
        # the glyph silhouettes — the visual the source CSS actually wants.
        if _try_apply_gradient_text_fill(font, spec, fallback_el):
            return
        color = _resolve_run_color(spec, fallback_el)
        if color is not None:
            try:
                font.color.rgb = color
            except Exception:
                pass

    # ----- native bullets ------------------------------------------------

    def _emit_native_bullet(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        # Treat a single ListItem unit as a one-paragraph bullet.
        x, y, w, h = _emu_rect(op.bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        _ensure_default_normautofit(tf)
        # PowerPoint substitutes Inter/etc with Calibri when the source font
        # isn't installed — Calibri is wider, so big headlines that fit in 3
        # lines in Chromium can overflow into the lede below in PowerPoint.
        # TEXT_TO_FIT_SHAPE shrinks the font just enough to keep everything
        # in-frame. We only apply it for *title-sized* text where the wrap
        # risk is real; on small frames (pills, badges, kickers) LibreOffice
        # over-shrinks the text to invisibility.
        # Heuristic: apply when the largest run is >= 20px (subtitle-ish or
        # bigger). Small inline labels are unaffected.
        if _has_title_sized_text(unit):
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        text_elems = [e for e in unit.all_elements() if e.text and e.text.strip()]
        if not text_elems:
            return
        first = True
        for e in text_elems:
            text = (e.pptx_text or e.text or "").strip()
            if not text:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ""
            run = p.add_run()
            run.text = "• " + text
            font = run.font
            font.name = resolve_font(e.font_family)
            font.size = Pt(max(8.0, parse_pt(e.font_size)))
            font.bold = is_bold(e.font_weight)
            color = parse_color(e.color)
            if color is not None:
                font.color.rgb = color[0]

    # ----- native shape --------------------------------------------------

    def _try_emit_native_decoration(
        self, slide, unit: VisualUnit, op: EmitOp
    ) -> bool:
        """If the unit's anchor decoration (gradient/shadow/solid) is all
        translatable, emit a native shape covering the bbox. Returns True if
        a native shape was emitted; False to fall back to raster.
        """
        anchor = unit.elements[0] if unit.elements else None
        if anchor is None:
            return False
        # Gradients we can translate — go native.
        bg_img = anchor.background_image
        if bg_img and bg_img != "none":
            if "url(" in bg_img:
                return False
            grad = parse_gradient(bg_img)
            if grad is None:
                return False
        # Decorative overlays may sit off-canvas (top:-220px etc); preserve
        # the original origin so the gradient center isn't displaced.
        radius = parse_px(anchor.border_radius)
        shape_kind, bbox_override = _shape_kind_for_anchor(anchor, radius)
        shape_bbox = bbox_override if bbox_override is not None else op.bbox
        x, y, w, h = _emu_rect_offcanvas(shape_bbox)
        deco_stack = derive_decorations(anchor, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, op.bbox)
        shape = slide.shapes.add_shape(shape_kind, x, y, w, h)
        _apply_rotation(shape, anchor)
        shape.line.fill.background()
        self._apply_fill(shape, anchor)
        self._apply_border(shape, anchor)
        self._apply_shadow(shape, anchor)
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, op.bbox)
        return True

    def _emit_native_shape(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        anchor_el = unit.elements[0] if unit.elements else None
        if anchor_el is None:
            return
        # Decorative overlays may sit off-canvas — see _emit_native_decoration.
        radius = parse_px(anchor_el.border_radius)
        shape_kind, bbox_override = _shape_kind_for_anchor(anchor_el, radius)
        shape_bbox = bbox_override if bbox_override is not None else op.bbox
        x, y, w, h = _emu_rect_offcanvas(shape_bbox)
        deco_stack = derive_decorations(anchor_el, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, op.bbox)
        shape = slide.shapes.add_shape(shape_kind, x, y, w, h)
        _apply_rotation(shape, anchor_el)
        shape.line.fill.background()  # default to no border
        self._apply_fill(shape, anchor_el)
        self._apply_border(shape, anchor_el)
        self._apply_shadow(shape, anchor_el)
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, op.bbox)

    def _apply_fill(self, shape, el: DomElement) -> None:
        # Try native gradient first (background-image: linear-gradient(...) / radial-gradient(...)).
        if el.background_image and el.background_image != "none":
            grad = parse_gradient(el.background_image)
            if grad is not None and apply_gradient_fill(shape, grad):
                return
            # Fall through to background_color / no-fill if gradient unparseable.
        bg = parse_color(el.background_color)
        if bg is None:
            try:
                shape.fill.background()
            except Exception:
                pass
            return
        rgb_color, alpha = bg
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb_color
        except Exception:
            return
        # Honor alpha when < 1: cards with `rgba(255,255,255,0.04)` were
        # rendering as solid WHITE and burning their light-gray text into
        # invisibility. python-pptx doesn't expose alpha on solid fills, so
        # drop down to OOXML and add an `<a:alpha val="..."/>` child.
        if alpha < 0.999:
            self._set_solid_fill_alpha(shape, alpha)

    def _set_solid_fill_alpha(self, shape, alpha: float) -> None:

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        try:
            sp_pr = shape._element.spPr
        except AttributeError:
            return
        if sp_pr is None:
            return
        solid = sp_pr.find(f"{{{ns_a}}}solidFill")
        if solid is None:
            return
        srgb = solid.find(f"{{{ns_a}}}srgbClr")
        if srgb is None:
            return
        # Clear any pre-existing alpha child.
        for existing in srgb.findall(f"{{{ns_a}}}alpha"):
            srgb.remove(existing)
        etree.SubElement(
            srgb,
            f"{{{ns_a}}}alpha",
            attrib={"val": str(int(round(max(0.0, min(1.0, alpha)) * 100_000)))},
        )

    def _apply_shadow(self, shape, el: DomElement) -> None:
        if not el.box_shadow or el.box_shadow == "none":
            return
        layers = parse_box_shadows(el.box_shadow)
        if layers:
            apply_shadows(shape, layers)

    def _apply_border(self, shape, el: DomElement) -> None:
        if not el.border or el.border == "none":
            try:
                shape.line.fill.background()
            except Exception:
                pass
            return
        # "1px solid rgb(0, 0, 0)"
        parts = el.border.split()
        if not parts:
            return
        width = parse_px(parts[0])
        if width <= 0:
            try:
                shape.line.fill.background()
            except Exception:
                pass
            return
        # Try to find a color in the border string
        m = re.search(r"(rgba?\([^)]+\)|#[0-9a-fA-F]{3,8})", el.border)
        col = parse_color(m.group(1)) if m else None
        try:
            shape.line.width = Emu(px_to_emu(width))
            if col is not None:
                shape.line.color.rgb = col[0]
        except Exception:
            pass

    # ----- native picture (img) ------------------------------------------

    def _emit_native_svg(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        svg_el = next((e for e in unit.all_elements() if e.is_svg and e.svg_shapes), None)
        if svg_el is None or not svg_el.svg_shapes:
            return
        emit_svg_shapes(slide, svg_el.svg_shapes)

    def _emit_native_table(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        """Emit an HTML <table> as a real PPTX table primitive.

        We size the table at the source bbox, then post-walk the shape's
        XML to apply per-cell colspan/rowspan merges (python-pptx exposes
        a `merge` API, but it works on cell objects after creation).

        Cells are populated with one paragraph each. Per-cell font, color,
        background, alignment, italic and underline are honored. Header
        rows (rows owned by <thead> or rows containing only <th> cells)
        force bold and inherit any background captured on the cell.
        """
        table_el = next(
            (e for e in unit.all_elements() if e.is_table and e.table_data),
            None,
        )
        if table_el is None or not table_el.table_data:
            return
        td = table_el.table_data
        rows: list[list[dict]] = td.get("rows") or []
        n_cols: int = int(td.get("n_cols") or 0)
        if not rows or n_cols <= 0:
            return

        # Build a logical grid (n_rows × n_cols) with each visible cell's
        # source row/col index. We emit every grid slot first, then merge.
        n_rows = len(rows)
        # Pre-walk to establish the grid → (source_row, source_col_in_row)
        # for every covered slot (including those swallowed by spans).
        grid: list[list[tuple[int, int] | None]] = [
            [None] * n_cols for _ in range(n_rows)
        ]
        for r_idx, row in enumerate(rows):
            c_idx = 0
            for cell_idx, cell in enumerate(row):
                while c_idx < n_cols and grid[r_idx][c_idx] is not None:
                    c_idx += 1
                if c_idx >= n_cols:
                    break
                cs = max(1, int(cell.get("colspan", 1)))
                rs = max(1, int(cell.get("rowspan", 1)))
                for dr in range(rs):
                    for dc in range(cs):
                        rr, cc = r_idx + dr, c_idx + dc
                        if 0 <= rr < n_rows and 0 <= cc < n_cols:
                            grid[rr][cc] = (r_idx, cell_idx)
                c_idx += cs

        x, y, w, h = _emu_rect(op.bbox)
        try:
            shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
        except Exception as e:
            log.warning("emitter.table_create_failed", error=str(e))
            return
        table = shape.table

        # Populate every visible (top-left of its span) cell.
        for r_idx, row in enumerate(rows):
            c_idx = 0
            for cell_idx, cell in enumerate(row):
                while c_idx < n_cols and (
                    grid[r_idx][c_idx] is None
                    or grid[r_idx][c_idx] != (r_idx, cell_idx)
                ):
                    c_idx += 1
                if c_idx >= n_cols:
                    break
                self._populate_table_cell(
                    table.cell(r_idx, c_idx), cell
                )
                c_idx += max(1, int(cell.get("colspan", 1)))

        # Apply spans by merging the originating cell with the bottom-right
        # of its span. python-pptx supports `origin.merge(other)` where
        # both must be top-left and bottom-right corners of the rect.
        for r_idx, row in enumerate(rows):
            c_idx = 0
            for cell in row:
                while c_idx < n_cols and (
                    grid[r_idx][c_idx] is None
                    or grid[r_idx][c_idx][0] != r_idx
                ):
                    c_idx += 1
                if c_idx >= n_cols:
                    break
                cs = max(1, int(cell.get("colspan", 1)))
                rs = max(1, int(cell.get("rowspan", 1)))
                if cs > 1 or rs > 1:
                    r2 = min(n_rows - 1, r_idx + rs - 1)
                    c2 = min(n_cols - 1, c_idx + cs - 1)
                    try:
                        table.cell(r_idx, c_idx).merge(table.cell(r2, c2))
                    except Exception as e:
                        log.debug(
                            "emitter.table_merge_failed",
                            r=r_idx, c=c_idx, rs=rs, cs=cs, error=str(e),
                        )
                c_idx += cs

    def _populate_table_cell(self, pptx_cell, cell: dict) -> None:
        """Set text + per-cell styling on a python-pptx table cell."""

        text = cell.get("text", "") or ""
        tf = pptx_cell.text_frame
        tf.word_wrap = True
        # Reset paragraphs (python-pptx pre-creates one).
        p = tf.paragraphs[0]
        # Clear any existing runs.
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        align = _TEXT_ALIGN_MAP.get(
            (cell.get("text_align") or "start").strip().lower(),
            PP_ALIGN.LEFT,
        )
        p.alignment = align
        run = p.add_run()
        run.text = text

        font = run.font
        size_px = parse_px(cell.get("font_size") or "14px")
        if size_px > 0:
            font.size = Pt(size_px * 0.75)  # px → pt at 96dpi
        weight = (cell.get("font_weight") or "400").strip()
        font.bold = bool(cell.get("is_header")) or is_bold(weight)
        if cell.get("italic"):
            font.italic = True
        if cell.get("underline"):
            font.underline = True
        family = (cell.get("font_family") or "").split(",")[0].strip().strip("\"'")
        if family:
            font.name = resolve_font(family)

        col = parse_color(cell.get("color") or "rgb(0,0,0)")
        if col is not None:
            try:
                font.color.rgb = col[0]
            except Exception:
                pass

        bg = parse_color(cell.get("background_color") or "rgba(0,0,0,0)")
        if bg is not None and bg[1] > 0.001:
            try:
                pptx_cell.fill.solid()
                pptx_cell.fill.fore_color.rgb = bg[0]
            except Exception:
                pass
        else:
            # Override the default theme fill so cells without an explicit
            # background render transparent rather than picking up the
            # accent-color row banding python-pptx applies by default.
            try:
                pptx_cell.fill.background()
            except Exception:
                pass

    def _emit_native_picture(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        src = op.decision.metadata.get("src") or self._first_img_src(unit)
        if not src:
            return
        try:
            data = self._fetch_image(src)
        except Exception as e:
            log.warning("emitter.image_fetch_failed", src=src[:120], error=str(e))
            return
        # Honor `<img style="opacity: 0.X">` by baking the alpha into the
        # image bytes against the slide's underlying bg color. We do this
        # at emit time (rather than via OOXML `<a:alphaModFix>`) because
        # LibreOffice and many older PowerPoint versions silently ignore
        # blip-level alpha — pre-multiplying produces the same visual on
        # every renderer at the cost of needing the bg color. Without this,
        # slides that use a faded photo as a backdrop under big type render
        # the photo at full strength, blowing out contrast and forcing the
        # overlay text into illegibility.
        opacity = self._picture_opacity(unit)
        if opacity is not None and opacity < 0.999:
            bg_rgb = self._underlying_bg_rgb(unit)
            try:
                data = self._bake_image_opacity(data, opacity, bg_rgb)
            except Exception as e:
                log.warning("emitter.bake_alpha_failed", error=str(e))
        x, y, w, h = _emu_rect(op.bbox)
        slide.shapes.add_picture(io.BytesIO(data), x, y, w, h)

    def _picture_opacity(self, unit: VisualUnit) -> float | None:
        """The IMG's CSS opacity (None ⇔ no IMG / opacity == 1)."""
        for e in unit.all_elements():
            if e.is_img and e.img_src:
                try:
                    op = float(e.opacity)
                except (TypeError, ValueError):
                    return None
                return op if 0.0 <= op < 1.0 else None
        return None

    def _underlying_bg_rgb(self, unit: VisualUnit) -> tuple[int, int, int]:
        """Best-effort guess at the color sitting behind a translucent IMG.

        We walk the IMG's ancestor chain looking for the first opaque
        `background-color` and use it as the alpha-composite base. Defaults
        to black, which is what most photo-backdrop slides use.
        """
        # Anchor parent chain via parent_id stored on every DomElement.
        elems = {e.id: e for e in unit.all_elements()}
        # Pick the IMG anchor.
        img = next(
            (e for e in unit.all_elements() if e.is_img and e.img_src), None
        )
        if img is None:
            return (0, 0, 0)
        cur = img
        while cur is not None:
            bgc = cur.background_color or ""
            parsed = parse_color(bgc)
            if parsed is not None:
                rgb_color, alpha = parsed
                if alpha >= 0.5:
                    return (
                        int((rgb_color >> 16) & 0xFF),
                        int((rgb_color >> 8) & 0xFF),
                        int(rgb_color & 0xFF),
                    )
            parent_id = cur.parent_id
            cur = elems.get(parent_id) if parent_id is not None else None
        return (0, 0, 0)

    def _bake_image_opacity(
        self, data: bytes, opacity: float, bg_rgb: tuple[int, int, int]
    ) -> bytes:
        """Alpha-composite an image against `bg_rgb` and return PNG bytes.

        `out = src * opacity + bg * (1 - opacity)` per pixel. Fast, exact,
        and renderer-portable.
        """
        from PIL import Image

        img = Image.open(io.BytesIO(data)).convert("RGB")
        # Use NumPy if available; fall back to a Python loop for portability.
        try:
            import numpy as np

            arr = np.asarray(img, dtype=np.float32)
            bg = np.asarray(bg_rgb, dtype=np.float32)
            out = arr * opacity + bg * (1.0 - opacity)
            out = out.clip(0, 255).astype("uint8")
            faded = Image.fromarray(out, mode="RGB")
        except ImportError:  # pragma: no cover
            blend = Image.new("RGB", img.size, bg_rgb)
            faded = Image.blend(blend, img, opacity)
        buf = io.BytesIO()
        faded.save(buf, format="PNG")
        return buf.getvalue()

    def _first_img_src(self, unit: VisualUnit) -> str | None:
        for e in unit.all_elements():
            if e.is_img and e.img_src:
                return e.img_src
        return None

    def _fetch_image(self, src: str) -> bytes:
        if src in self._image_cache:
            return self._image_cache[src]
        if src.startswith("data:"):
            import base64

            head, _, b64 = src.partition(",")
            if ";base64" in head:
                data = base64.b64decode(b64)
            else:
                data = b64.encode("utf-8")
            self._image_cache[src] = data
            return data
        parsed = urlparse(src)
        if parsed.scheme in ("http", "https"):
            r = self._http.get(src)
            r.raise_for_status()
            data = r.content
        elif parsed.scheme in ("file", ""):
            data = Path(parsed.path or src).read_bytes()
        else:
            raise EmitError(f"unsupported image scheme: {parsed.scheme}")
        self._image_cache[src] = data
        return data

    # ----- raster --------------------------------------------------------

    async def _emit_raster(
        self,
        slide,
        unit: VisualUnit,
        op: EmitOp,
        rendered: RenderedSlide,
        renderer,
    ) -> None:
        # Try a region screenshot from the live page; fall back to ground-truth crop.
        png: bytes | None = None
        try:
            png = await renderer.screenshot_region(
                rendered.html,
                "",
                (op.bbox.x, op.bbox.y, op.bbox.w, op.bbox.h),
            )
        except Exception as e:
            log.warning("emitter.region_screenshot_failed", error=str(e))
            png = None
        if not png:
            self._emit_region_raster(slide, rendered.ground_truth_png, op.bbox)
            return
        x, y, w, h = _emu_rect(op.bbox)
        slide.shapes.add_picture(io.BytesIO(png), x, y, w, h)

    def _emit_region_raster(
        self, slide, ground_truth_png: bytes, bbox: BoundingBox
    ) -> None:
        from PIL import Image

        gt = Image.open(io.BytesIO(ground_truth_png))
        bb = _clamp_bbox(bbox)
        crop = gt.crop(
            (
                int(bb.x),
                int(bb.y),
                int(bb.x + bb.w),
                int(bb.y + bb.h),
            )
        )
        out = io.BytesIO()
        crop.save(out, format="PNG")
        out.seek(0)
        x, y, w, h = _emu_rect(bb)
        slide.shapes.add_picture(out, x, y, w, h)

    def _emit_full_raster(self, slide, png_bytes: bytes) -> None:
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            Emu(0),
            Emu(0),
            Emu(SLIDE_W_EMU),
            Emu(SLIDE_H_EMU),
        )

    def _set_notes(self, slide, text: str) -> None:
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception as e:
            log.warning("emitter.notes_failed", error=str(e))
