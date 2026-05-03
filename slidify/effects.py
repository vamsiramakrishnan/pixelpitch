"""Wave-2 multi-shadow stack — native split + sibling-rect fallback.

PPTX shape effects support exactly **one** outer shadow + one inner shadow
per shape (via `<a:effectLst>`). When the IR carries more than that on a
node, F3 splits:

  * The first outer shadow + first inner shadow attach natively.
  * Each additional shadow becomes a sibling, transparent, same-bbox
    shape with its own native shadow effect attached.

This is the `StatCardWithDepth` strategy generalized — see CONTRACT §3.3.
The compiler routes any node with `len(normalize_shadows(node)) >= 1` here;
the function returns the list of sibling shapes the caller adds in the
correct z-order.
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from slidify.geom import px_to_emu
from slidify.ir import _color_to_hex_alpha
from slidify.shadows import BoxShadow as ShBoxShadow
from slidify.shadows import apply_shadows

if TYPE_CHECKING:
    from slidify.ir import IRBbox, IRBoxShadow


def _to_internal(sh: IRBoxShadow) -> ShBoxShadow:
    hex_, alpha = _color_to_hex_alpha(sh.color)
    bare = hex_.lstrip("#").upper()[:6]
    return ShBoxShadow(
        inset=sh.inset,
        offset_x_px=sh.offsetX,
        offset_y_px=sh.offsetY,
        blur_px=max(0.0, sh.blur),
        spread_px=sh.spread,
        color_hex=bare,
        alpha=alpha,
    )


def _split_native_pair(
    shadows: list[IRBoxShadow],
) -> tuple[list[IRBoxShadow], list[IRBoxShadow]]:
    """Pick the first outer + first inner; return (native_pair, leftover).

    The relative order of the native pair preserves the IR list's order
    (CSS-paint semantics: first-listed = topmost). Leftover keeps source
    order so the sibling-rect emission z-orders the same way users
    authored.
    """
    native: list[IRBoxShadow] = []
    leftover: list[IRBoxShadow] = []
    seen_outer = False
    seen_inner = False
    for sh in shadows:
        if sh.inset:
            if not seen_inner:
                seen_inner = True
                native.append(sh)
            else:
                leftover.append(sh)
        else:
            if not seen_outer:
                seen_outer = True
                native.append(sh)
            else:
                leftover.append(sh)
    return native, leftover


def _add_sibling_rect_with_shadow(
    slide,
    bbox: IRBbox,
    sh: IRBoxShadow,
):
    """Add a transparent rect at `bbox` carrying one native shadow."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(px_to_emu(bbox.x)),
        Emu(px_to_emu(bbox.y)),
        Emu(px_to_emu(bbox.w)),
        Emu(px_to_emu(bbox.h)),
    )
    try:
        rect.fill.background()
        rect.line.fill.background()
    except Exception:
        pass
    apply_shadows(rect, [_to_internal(sh)])
    return rect


def apply_shadow_stack(
    shape,
    shadows: list[IRBoxShadow],
    slide=None,
    parent_bbox: IRBbox | None = None,
) -> list:
    """Attach a multi-shadow stack to `shape`, spilling extras as siblings.

    Args:
        shape: The primary shape carrying the first outer + first inner.
        shadows: IR shadow list (already normalized — see
            `slidify.ir.normalize_shadows`). Empty list = no-op.
        slide: Required when `len(shadows) > 2` and the leftover needs
            sibling rects. May be None when only 1–2 shadows are present.
        parent_bbox: The IRBbox of `shape`. Required when slide is.

    Returns:
        List of sibling shapes added to `slide.shapes` (empty for ≤2
        compatible shadows).
    """
    if not shadows:
        return []
    # Cap at the §1.2 hard limit: 4. Crews exceeding that should redesign.
    capped = list(shadows)[:4]
    native, leftover = _split_native_pair(capped)
    apply_shadows(shape, [_to_internal(sh) for sh in native])
    if not leftover:
        return []
    if slide is None or parent_bbox is None:
        # Caller didn't provide enough context to spill; silently drop.
        return []
    out: list = []
    for sh in leftover:
        out.append(_add_sibling_rect_with_shadow(slide, parent_bbox, sh))
    # Z-order fix: leftover shadow rects were appended AFTER `shape`, so
    # their soft halos render ON TOP of the primary shape — producing
    # visible darkening around the perimeter for outer shadows. Move
    # `shape` to the end of the spTree so the halos sit beneath it,
    # matching the manual sibling-rect-before-base ordering used by
    # StatCardWithDepth.tsx (the original spec exemplar).
    try:
        sp_tree = shape._element.getparent()
        if sp_tree is not None:
            sp_tree.remove(shape._element)
            sp_tree.append(shape._element)
    except Exception:
        pass
    return out


__all__ = [
    "apply_shadow_stack",
]
