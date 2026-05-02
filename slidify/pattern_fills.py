"""Wave-2 PatternFill — `<a:pattFill>` / tiled-shape emission.

PPTX has 50+ built-in `prstFill` patterns; we map IR's six pattern variants
to the closest-looking presets where possible. The `dots` pattern can't be
satisfied by any preset (PPTX's `dotGrid`/`dotDmnd` are spaced too widely
and offer no spacing knob), so it emits as a tiled group of small ovals.

Module is named `pattern_fills` (not `patterns`) because `slidify.patterns`
already exists as a sub-package. The CONTRACT §3.2 spec lists `patterns.py`;
the rename is a deviation documented in the F3 PR body.
"""

from __future__ import annotations

import math
from dataclasses import dataclass
from typing import TYPE_CHECKING

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from slidify.colors import parse_color
from slidify.geom import px_to_emu

if TYPE_CHECKING:
    from slidify.ir import IRBbox, IRPatternFill

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# IR pattern → OOXML `<a:pattFill prst="…">` preset. `'dots'` is absent
# because it tiles small ovals instead.
_PRST_MAP_LIGHT = {
    "lines-h": "ltHorz",
    "lines-v": "ltVert",
    "lines-grid": "smGrid",
    "diagonal": "ltUpDiag",
    "crosshatch": "diagCross",
}
# Heavier-stroke variant when featureSizePx >= 2.
_PRST_MAP_DARK = {
    "lines-h": "dkHorz",
    "lines-v": "dkVert",
    "lines-grid": "lgGrid",
    "diagonal": "dkUpDiag",
    "crosshatch": "diagCross",  # no "dark" variant; keep the light one
}


@dataclass(frozen=True)
class DotShapeSpec:
    """One small oval emitted by the dot-grid tiler.

    Coordinates are absolute slide-pixel space (top-left origin),
    matching IRBbox conventions.
    """

    x: float
    y: float
    diameter: float
    color_hex: str  # 6-char hex, no leading #
    alpha: float


def _hex_no_hash(color_hex: str) -> str:
    return color_hex.lstrip("#").upper()[:6]


def _set_pattfill(
    shape,
    prst: str,
    fg_hex: str,
    fg_alpha: float,
    bg_hex: str,
    bg_alpha: float,
) -> bool:
    """Replace `shape`'s fill with `<a:pattFill>` carrying fg+bg colors."""
    try:
        sp_pr = shape._element.spPr
    except AttributeError:
        return False
    if sp_pr is None:
        return False

    # Strip existing fills.
    for tag_local in (
        "solidFill",
        "gradFill",
        "blipFill",
        "pattFill",
        "noFill",
        "grpFill",
    ):
        for existing in sp_pr.findall(f"{{{NS_A}}}{tag_local}"):
            sp_pr.remove(existing)

    patt = etree.SubElement(
        sp_pr, f"{{{NS_A}}}pattFill", attrib={"prst": prst}
    )
    sp_pr.remove(patt)
    sp_pr.insert(0, patt)

    fg = etree.SubElement(patt, f"{{{NS_A}}}fgClr")
    fg_srgb = etree.SubElement(
        fg, f"{{{NS_A}}}srgbClr", attrib={"val": _hex_no_hash(fg_hex)}
    )
    if fg_alpha < 0.999:
        etree.SubElement(
            fg_srgb,
            f"{{{NS_A}}}alpha",
            attrib={"val": str(int(round(fg_alpha * 100_000)))},
        )

    bg = etree.SubElement(patt, f"{{{NS_A}}}bgClr")
    bg_srgb = etree.SubElement(
        bg, f"{{{NS_A}}}srgbClr", attrib={"val": _hex_no_hash(bg_hex)}
    )
    if bg_alpha < 0.999:
        etree.SubElement(
            bg_srgb,
            f"{{{NS_A}}}alpha",
            attrib={"val": str(int(round(bg_alpha * 100_000)))},
        )
    return True


def _resolve_color(c) -> tuple[str, float]:
    """Coerce IR Color (str | ColorWithAlpha) → ('RRGGBB', alpha)."""
    if c is None:
        return ("FFFFFF", 0.0)
    if isinstance(c, str):
        col = parse_color(c)
        if col is None:
            return ("000000", 1.0)
        rgb, a = col
        return (f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}", a)
    # ColorWithAlpha
    hex_ = getattr(c, "hex", "#000000")
    alpha = float(getattr(c, "alpha", 1.0))
    return (hex_.lstrip("#").upper()[:6], alpha)


def _dot_specs(
    bbox: IRBbox,
    pattern_fill: IRPatternFill,
) -> list[DotShapeSpec]:
    """Tile small ovals across `bbox` per `pattern_fill` spacing.

    Spacing is `tileWidthPx` × `tileHeightPx`; dot diameter is
    `2 * featureSizePx`. Each dot is centered on a tile.
    """
    fg_hex, fg_alpha = _resolve_color(pattern_fill.fgColor)
    diameter = max(0.5, 2.0 * pattern_fill.featureSizePx)
    tile_w = max(2.0, pattern_fill.tileWidthPx)
    tile_h = max(2.0, pattern_fill.tileHeightPx)
    cols = max(1, int(math.floor(bbox.w / tile_w)))
    rows = max(1, int(math.floor(bbox.h / tile_h)))
    out: list[DotShapeSpec] = []
    for r in range(rows):
        cy = bbox.y + (r + 0.5) * tile_h
        for c in range(cols):
            cx = bbox.x + (c + 0.5) * tile_w
            out.append(
                DotShapeSpec(
                    x=cx - diameter / 2.0,
                    y=cy - diameter / 2.0,
                    diameter=diameter,
                    color_hex=fg_hex,
                    alpha=fg_alpha,
                )
            )
    return out


def _emit_dots_into(
    slide,
    bbox: IRBbox,
    pattern_fill: IRPatternFill,
) -> list:
    """Add an oval per tile to `slide.shapes`. Returns the list."""
    specs = _dot_specs(bbox, pattern_fill)
    out: list = []
    for spec in specs:
        d_emu = Emu(px_to_emu(spec.diameter))
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(px_to_emu(spec.x)),
            Emu(px_to_emu(spec.y)),
            d_emu,
            d_emu,
        )
        try:
            oval.line.fill.background()
        except Exception:
            pass
        # Fill via lxml so we can attach an alpha if needed.
        try:
            sp_pr = oval._element.spPr
            for tag in ("solidFill", "noFill", "gradFill", "pattFill"):
                for existing in sp_pr.findall(f"{{{NS_A}}}{tag}"):
                    sp_pr.remove(existing)
            sf = etree.SubElement(sp_pr, f"{{{NS_A}}}solidFill")
            sp_pr.remove(sf)
            sp_pr.insert(0, sf)
            srgb = etree.SubElement(
                sf, f"{{{NS_A}}}srgbClr", attrib={"val": spec.color_hex}
            )
            if spec.alpha < 0.999:
                etree.SubElement(
                    srgb,
                    f"{{{NS_A}}}alpha",
                    attrib={"val": str(int(round(spec.alpha * 100_000)))},
                )
        except Exception:
            pass
        out.append(oval)
    return out


def apply_pattern_fill(
    shape,
    pattern_fill: IRPatternFill,
    *,
    slide=None,
    bbox: IRBbox | None = None,
) -> list:
    """Apply an IR `IRPatternFill` to `shape`. Returns sibling shapes the
    caller should treat as part of the pattern (only non-empty for `dots`).

    Args:
        shape: The primary shape (rect/oval) whose surface receives the
            pattern. For `dots`, this shape's fill is set to `none` and the
            ovals are added to `slide.shapes` as siblings.
        pattern_fill: IR pattern fill spec.
        slide: Required iff `pattern_fill.pattern == 'dots'`.
        bbox: Required iff `pattern_fill.pattern == 'dots'` — the area to tile.

    Returns:
        List of additional shapes added to `slide.shapes` (empty for the
        prst-pattern variants).
    """
    if pattern_fill.pattern == "dots":
        if slide is None or bbox is None:
            raise ValueError(
                "apply_pattern_fill('dots'): slide and bbox are required"
            )
        # Make the host shape transparent so only the dots show.
        try:
            sp_pr = shape._element.spPr
            for tag in ("solidFill", "gradFill", "pattFill", "blipFill"):
                for existing in sp_pr.findall(f"{{{NS_A}}}{tag}"):
                    sp_pr.remove(existing)
            etree.SubElement(sp_pr, f"{{{NS_A}}}noFill")
        except Exception:
            pass
        return _emit_dots_into(slide, bbox, pattern_fill)

    # Pattern preset path.
    use_dark = pattern_fill.featureSizePx >= 2.0
    table = _PRST_MAP_DARK if use_dark else _PRST_MAP_LIGHT
    prst = table.get(pattern_fill.pattern, "ltHorz")
    fg_hex, fg_alpha = _resolve_color(pattern_fill.fgColor)
    if pattern_fill.bgColor is not None:
        bg_hex, bg_alpha = _resolve_color(pattern_fill.bgColor)
    else:
        bg_hex, bg_alpha = ("FFFFFF", 0.0)
    _set_pattfill(shape, prst, fg_hex, fg_alpha, bg_hex, bg_alpha)
    # angleDeg → shape rotation. PPTX has no per-pattern rotation.
    if pattern_fill.angleDeg:
        try:
            shape.rotation = float(pattern_fill.angleDeg)
        except Exception:
            pass
    return []


__all__ = [
    "DotShapeSpec",
    "apply_pattern_fill",
]
