"""Post-emit editability round-trip check.

After `Emitter.save()` we re-open the produced .pptx and count, per slide,
how many editable primitives (text frames, shapes, tables, pictures, group
shapes) actually made it into the file. We compare that against the
*intended* counts derived from the slide's emit ops. A drift means
python-pptx (or the sanitization pass, or the OOXML repair logic) silently
dropped a shape — exactly the class of bug that's invisible in the SSIM
oracle (the emitter "succeeded", the file just lost a shape).

This is a small, post-hoc structural diff. It does NOT verify pixel
fidelity — the SSIM/OCR oracle covers that. It DOES verify that what we
asked PowerPoint to render is what's actually in the file.
"""

from __future__ import annotations

from collections import Counter
from pathlib import Path

from pydantic import BaseModel, Field

from slidify.models import DecisionKind, EmitOp

# Decision kinds that should produce at least one editable shape in the
# .pptx. (Skip / Raster live separately — Raster IS a shape, but a picture,
# not an editable primitive; we count it explicitly below.)
_EDITABLE_KINDS = {
    DecisionKind.NativeText,
    DecisionKind.NativeShape,
    DecisionKind.NativeBullet,
    DecisionKind.NativePicture,
    DecisionKind.NativeSvg,
    DecisionKind.NativeTable,
    DecisionKind.Hybrid,  # decoration raster + native children — counts as native
}


class SlideEditabilityReport(BaseModel):
    """Per-slide structural diff between intent and what the .pptx contains."""

    slide_index: int
    intended_editable: int
    actual_editable: int
    intended_pictures: int
    actual_pictures: int
    intended_tables: int
    actual_tables: int
    intended_text_frames: int
    actual_text_frames: int
    notes: str = ""

    @property
    def passed(self) -> bool:
        # We treat ≥ as a pass: emitters routinely add MORE shapes than the
        # decision count (decoration stacks, gradient fills emitted as
        # multiple layers). The failure mode is shapes silently DROPPED.
        return (
            self.actual_editable >= self.intended_editable
            and self.actual_tables >= self.intended_tables
            and self.actual_pictures >= self.intended_pictures
        )


class EditabilityReport(BaseModel):
    """Deck-level rollup of per-slide editability checks."""

    n_slides: int
    n_passed: int
    per_slide: list[SlideEditabilityReport] = Field(default_factory=list)

    @property
    def passed(self) -> bool:
        return self.n_passed == self.n_slides


def _intended_counts(ops: list[EmitOp]) -> dict[str, int]:
    """Roll a slide's emit ops into per-primitive intended counts."""
    by_kind: Counter[DecisionKind] = Counter(o.decision.kind for o in ops)
    editable = sum(by_kind.get(k, 0) for k in _EDITABLE_KINDS)
    return {
        "editable": editable,
        "pictures": by_kind.get(DecisionKind.NativePicture, 0),
        "tables": by_kind.get(DecisionKind.NativeTable, 0),
        "text_frames": (
            by_kind.get(DecisionKind.NativeText, 0)
            + by_kind.get(DecisionKind.NativeBullet, 0)
        ),
    }


def _actual_counts(slide) -> dict[str, int]:
    """Walk a python-pptx slide and tally what's actually in it."""
    n_text = 0
    n_pic = 0
    n_table = 0
    n_editable = 0
    for shp in slide.shapes:
        # has_text_frame / has_table / has_chart are exposed by python-pptx.
        try:
            if getattr(shp, "has_text_frame", False) and shp.text_frame.text.strip():
                n_text += 1
        except Exception:
            pass
        if getattr(shp, "has_table", False):
            n_table += 1
            n_editable += 1
            continue
        # Pictures (shape_type == 13 / PICTURE in python-pptx).
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_pic += 1
                # Pictures count as editable in the "didn't disappear" sense
                # (the file holds the bytes), but not as native primitives —
                # we count them in their own bucket above and exclude from
                # n_editable to avoid double-counting against intent.
                continue
            # Every native vector primitive counts as editable: AUTO_SHAPE
            # (rect/oval/preset), TEXT_BOX, GROUP, plus LINE / FREEFORM /
            # CONNECTOR — these last three are emitted by NativeSvg for
            # `<line>`, `<path>`, and `<polyline>` and remain individually
            # selectable & editable in PowerPoint. Omitting them caused
            # decks dense in SVG (echo trails, blueprints, sparkbars) to
            # spuriously fail the round-trip check.
            if shp.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.TEXT_BOX,
                MSO_SHAPE_TYPE.GROUP,
                MSO_SHAPE_TYPE.LINE,
                MSO_SHAPE_TYPE.FREEFORM,
            ):
                n_editable += 1
                continue
        except Exception:
            pass
    return {
        "editable": n_editable,
        "pictures": n_pic,
        "tables": n_table,
        "text_frames": n_text,
    }


def check_pptx_editability(
    pptx_path: str | Path, ops_per_slide: list[list[EmitOp]]
) -> EditabilityReport:
    """Re-open the produced PPTX and verify intended editable shapes survived.

    Args:
        pptx_path: path to the .pptx the emitter produced.
        ops_per_slide: the per-slide list of EmitOps the emitter was asked to
            render. Order must match slide order in the .pptx.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    per_slide: list[SlideEditabilityReport] = []
    for idx, (slide, ops) in enumerate(zip(prs.slides, ops_per_slide, strict=False)):
        intent = _intended_counts(ops)
        actual = _actual_counts(slide)
        notes_parts: list[str] = []
        if actual["editable"] < intent["editable"]:
            notes_parts.append(
                f"editable shapes {actual['editable']} < intended {intent['editable']}"
            )
        if actual["tables"] < intent["tables"]:
            notes_parts.append(
                f"tables {actual['tables']} < intended {intent['tables']}"
            )
        if actual["pictures"] < intent["pictures"]:
            notes_parts.append(
                f"pictures {actual['pictures']} < intended {intent['pictures']}"
            )
        per_slide.append(
            SlideEditabilityReport(
                slide_index=idx,
                intended_editable=intent["editable"],
                actual_editable=actual["editable"],
                intended_pictures=intent["pictures"],
                actual_pictures=actual["pictures"],
                intended_tables=intent["tables"],
                actual_tables=actual["tables"],
                intended_text_frames=intent["text_frames"],
                actual_text_frames=actual["text_frames"],
                notes="; ".join(notes_parts),
            )
        )
    n_passed = sum(1 for r in per_slide if r.passed)
    return EditabilityReport(
        n_slides=len(per_slide), n_passed=n_passed, per_slide=per_slide
    )
