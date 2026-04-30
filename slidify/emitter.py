"""PPTX emitter: turns EmitOps into a python-pptx Presentation."""

from __future__ import annotations

import io
import re
from pathlib import Path
from urllib.parse import urlparse

import httpx
import structlog
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from slidify.colors import parse_color
from slidify.decorations import derive_decorations
from slidify.exceptions import EmitError
from slidify.fonts import is_bold
from slidify.fonts import resolve as resolve_font
from slidify.geom import (
    SLIDE_H_EMU,
    SLIDE_H_PX,
    SLIDE_W_EMU,
    SLIDE_W_PX,
    parse_pt,
    parse_px,
    px_to_emu,
)
from slidify.gradients import apply_gradient_fill, parse_gradient
from slidify.models import (
    BoundingBox,
    DecisionKind,
    DomElement,
    EmitOp,
    RenderedSlide,
    VisualUnit,
)
from slidify.preset_shapes import detect_preset_shape, emit_preset
from slidify.shadows import apply_shadows, parse_box_shadows
from slidify.svg_shapes import emit_svg_shapes
from slidify.text_metrics import (
    compute_font_scale_for_textbox,
    shrink_pill_bbox_to_fit_text,
)

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _apply_explicit_autofit(
    text_frame, font_scale_pct: int, line_space_reduction_pct: int
) -> None:
    """Bake an explicit ``<a:normAutofit fontScale=... lnSpcReduction=.../>``
    into the textframe, replacing any existing autofit element. Pre-computing
    the shrink at emit time makes the deck render identically across
    LibreOffice/PowerPoint regardless of which font gets substituted."""
    from lxml import etree as _et

    bodyPr = text_frame._txBody.bodyPr
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for existing in bodyPr.findall(f"{{{_A_NS}}}{tag}"):
            bodyPr.remove(existing)
    norm = _et.SubElement(bodyPr, f"{{{_A_NS}}}normAutofit")
    norm.set("fontScale", str(font_scale_pct))
    norm.set("lnSpcReduction", str(line_space_reduction_pct))

log = structlog.get_logger(__name__)


_TEXT_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "right": PP_ALIGN.RIGHT,
    "center": PP_ALIGN.CENTER,
    "justify": PP_ALIGN.JUSTIFY,
    "start": PP_ALIGN.LEFT,
    "end": PP_ALIGN.RIGHT,
}


def _clamp_bbox(b: BoundingBox) -> BoundingBox:
    """Tight clamp to slide bounds — used only when we need a contained
    region (e.g., screenshot crops). DO NOT use for placing decorative
    shapes; clamping a top:-220px overlay to top:0 *moves* its visual
    center. Use `_emu_rect_offcanvas` instead for shapes whose origin
    can legitimately sit outside the slide.
    """
    x = max(0.0, min(b.x, SLIDE_W_PX - 1))
    y = max(0.0, min(b.y, SLIDE_H_PX - 1))
    w = max(1.0, min(b.w, SLIDE_W_PX - x))
    h = max(1.0, min(b.h, SLIDE_H_PX - y))
    return BoundingBox(x=x, y=y, w=w, h=h)


def _emu_rect(b: BoundingBox) -> tuple[Emu, Emu, Emu, Emu]:
    b = _clamp_bbox(b)
    return (
        Emu(px_to_emu(b.x)),
        Emu(px_to_emu(b.y)),
        Emu(px_to_emu(b.w)),
        Emu(px_to_emu(b.h)),
    )


def _emu_rect_offcanvas(b: BoundingBox) -> tuple[Emu, Emu, Emu, Emu]:
    """Convert px bbox → EMU rect WITHOUT clamping origin to >= 0. PPTX
    accepts negative shape positions and clips at the slide boundary, which
    is exactly what we want for decorative overlays anchored off-canvas
    (auroras, glow ovals positioned with top: -220px etc).

    Width/height are clamped to >= 1 px to avoid python-pptx errors on
    zero-size shapes; otherwise the shape is emitted as-is.
    """
    w = max(1.0, b.w)
    h = max(1.0, b.h)
    return (
        Emu(px_to_emu(b.x)),
        Emu(px_to_emu(b.y)),
        Emu(px_to_emu(w)),
        Emu(px_to_emu(h)),
    )


class Emitter:
    """Builds a PPTX from rendered slides + per-slide EmitOps."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Emu(SLIDE_W_EMU)
        self.prs.slide_height = Emu(SLIDE_H_EMU)
        self._image_cache: dict[str, bytes] = {}
        self._http = httpx.Client(timeout=10.0, follow_redirects=True)
        # Brand palette used by decoration layers (mesh glow, etc.) when the
        # source HTML opts in via `data-slidify-decorate`. The API populates
        # this from the deck-wide color harvest right before save.
        self._brand_palette: list[str] = []

    def set_brand_palette(self, palette: list[str]) -> None:
        """Set the deck-wide brand palette used by decoration layers."""
        self._brand_palette = list(palette)

    def close(self) -> None:
        self._http.close()

    # ----- public --------------------------------------------------------

    async def emit_slide(
        self,
        slide_index: int,
        rendered: RenderedSlide,
        units_by_id: dict[str, VisualUnit],
        ops: list[EmitOp],
        renderer,
        notes: str = "",
    ) -> None:
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)

        if rendered.degraded:
            log.warning("emitter.degraded_slide", slide_index=slide_index, reason=rendered.reason)
            # Fall back to ground-truth raster as a safety net.
            self._emit_full_raster(slide, rendered.ground_truth_png)
            if notes:
                self._set_notes(slide, notes)
            return

        # Sort ops by z_order (already pre-order; but EmitOps may include explicit z)
        ops_sorted = sorted(ops, key=lambda o: o.z_order)
        for op in ops_sorted:
            try:
                await self._emit_op(slide, op, units_by_id, rendered, renderer)
            except Exception as e:
                log.warning(
                    "emitter.op_failed",
                    unit_id=op.unit_id,
                    kind=op.decision.kind,
                    error=str(e),
                )
                # Last-resort: fall back to ground-truth crop for this region.
                try:
                    self._emit_region_raster(
                        slide, rendered.ground_truth_png, op.bbox
                    )
                except Exception as e2:
                    log.error("emitter.fallback_failed", error=str(e2))

        if notes:
            self._set_notes(slide, notes)

    def save(self, path: Path | str) -> None:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._sanitize_for_repair_dialog()
        self.prs.save(str(path))

    def _sanitize_for_repair_dialog(self) -> None:
        """Walk every slide's XML and fix the two patterns that most often
        trigger PowerPoint's "found a problem with some content" repair
        dialog at open time:

          1. Empty paragraph `<a:p/>` with no child runs and no
             `<a:endParaRPr>` — PowerPoint logs an error then quietly
             repairs by inserting one. We insert it ourselves so the
             open is silent.
          2. Run text `<a:t>foo bar</a:t>` containing leading/trailing
             whitespace without `xml:space="preserve"` — XML's default
             whitespace handling strips it, and PowerPoint repairs by
             re-asserting `preserve`.

        Both fixes are idempotent and additive — they only add elements
        / attributes that are missing.
        """
        from lxml import etree as _et

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"
        for slide in self.prs.slides:
            tree = slide._element  # noqa: SLF001
            for p in tree.iter(f"{{{ns_a}}}p"):
                # 1. Insert <a:endParaRPr> on empty paragraphs.
                has_run = any(
                    child.tag in (f"{{{ns_a}}}r", f"{{{ns_a}}}fld", f"{{{ns_a}}}br")
                    for child in p
                )
                has_end_pr = p.find(f"{{{ns_a}}}endParaRPr") is not None
                if not has_run and not has_end_pr:
                    _et.SubElement(
                        p,
                        f"{{{ns_a}}}endParaRPr",
                        attrib={"lang": "en-US"},
                    )
                # 2. xml:space="preserve" on whitespace-bearing <a:t>.
                for t in p.iter(f"{{{ns_a}}}t"):
                    if t.text and (
                        t.text != t.text.strip() or "  " in t.text
                    ):
                        if t.get(XML_SPACE) != "preserve":
                            t.set(XML_SPACE, "preserve")

    # ----- per-op dispatch ----------------------------------------------

    async def _emit_op(
        self,
        slide,
        op: EmitOp,
        units_by_id: dict[str, VisualUnit],
        rendered: RenderedSlide,
        renderer,
    ) -> None:
        unit = units_by_id.get(op.unit_id)
        if unit is None:
            return
        kind = op.decision.kind

        if kind == DecisionKind.Skip:
            return

        if kind == DecisionKind.NativeText:
            self._emit_native_text(slide, unit, op)
            return

        if kind == DecisionKind.NativeBullet:
            self._emit_native_bullet(slide, unit, op)
            return

        if kind == DecisionKind.NativeShape:
            self._emit_native_shape(slide, unit, op)
            return

        if kind == DecisionKind.NativePicture:
            self._emit_native_picture(slide, unit, op)
            return

        if kind == DecisionKind.NativeSvg:
            self._emit_native_svg(slide, unit, op)
            return

        if kind == DecisionKind.Raster:
            await self._emit_raster(slide, unit, op, rendered, renderer)
            return

        if kind == DecisionKind.Hybrid:
            # Surgical hybrid emission, in order of preference:
            #   1) Try emitting a native gradient/shape (works when the unit's
            #      decoration is fully translatable — gradient + solid + shadow).
            #   2) Else if the renderer captured a no-text decoration-only pass,
            #      crop from THAT (pixel-exact decoration with no text bleed).
            #   3) Else fall back to live region screenshot or ground-truth crop.
            # Children emit independently on top in all three branches.
            if self._try_emit_native_decoration(slide, unit, op):
                return
            if rendered.no_text_png:
                self._emit_region_raster(slide, rendered.no_text_png, op.bbox)
                return
            await self._emit_raster(slide, unit, op, rendered, renderer)
            return

    # ----- native text ---------------------------------------------------

    def _pick_text_anchor(self, unit: VisualUnit) -> DomElement | None:
        """Return the deepest text-bearing element (leaf or text container)."""
        text_elems = [
            e
            for e in unit.all_elements()
            if (e.text and e.text.strip())
            or (e.runs and any(r.text.strip() for r in e.runs if not r.is_break))
        ]
        if not text_elems:
            return None
        for e in text_elems:
            if e.pptx_text:
                return e
        return max(text_elems, key=lambda e: parse_px(e.font_size))

    def _emit_native_text(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        anchor = self._pick_text_anchor(unit)
        if anchor is None:
            return
        # Give multi-run inline text containers some horizontal slack so font
        # substitution doesn't wrap them. CSS sizes flex-item widths by their
        # content; PPTX uses fixed bboxes, so a "slidify · v1.0" inline-flex
        # container that fit in 102px in Chromium overflows in Calibri. Grow
        # the bbox by 30% horizontally — only for short multi-run frames where
        # the wrap risk is real.
        emit_bbox = op.bbox
        # Naked inline text only — if the anchor paints its own background
        # (a pill, a chip, a trend tag) the slack would stretch that fill
        # past the text and leave a visible "tail" rectangle.
        anchor_paints_bg = (
            (anchor.background_color and anchor.background_color != "rgba(0, 0, 0, 0)")
            or (anchor.background_image and anchor.background_image != "none")
        )
        if (
            anchor.is_text_container
            and anchor.runs
            and len(anchor.runs) >= 2
            and op.bbox.h <= 48
            and op.bbox.w <= 320
            and not anchor_paints_bg
        ):
            extra = op.bbox.w * 0.30
            emit_bbox = BoundingBox(
                x=op.bbox.x,
                y=op.bbox.y,
                w=op.bbox.w + extra,
                h=op.bbox.h,
            )

        # Pill shrink — kills the "green tail" past the substituted-font
        # text in single-line bg-painted containers. Only applied to
        # truly pill-sized boxes (single-line, narrow).
        if (
            anchor_paints_bg
            and anchor.is_text_container
            and (not anchor.runs or len(anchor.runs) <= 1)
            and emit_bbox.h <= 48
            and (anchor.text or "").strip()
        ):
            try:
                pt_size = max(8.0, parse_pt(anchor.font_size))
                emit_bbox = shrink_pill_bbox_to_fit_text(
                    anchor.text.strip(), emit_bbox, pt_size
                )
            except Exception:
                pass

        x, y, w, h = _emu_rect(emit_bbox)
        # Decoration layer (BELOW): mesh glow / drop shadow attached to
        # the text container's anchor. Opt-in via `data-slidify-decorate`.
        deco_stack = derive_decorations(anchor, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, emit_bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        # PowerPoint substitutes Inter/etc with Calibri when the source font
        # isn't installed — Calibri is wider, so big headlines that fit in 3
        # lines in Chromium can overflow into the lede below in PowerPoint.
        # TEXT_TO_FIT_SHAPE shrinks the font just enough to keep everything
        # in-frame. We only apply it for *title-sized* text where the wrap
        # risk is real; on small frames (pills, badges, kickers) LibreOffice
        # over-shrinks the text to invisibility.
        # Heuristic: apply when the largest run is >= 20px (subtitle-ish or
        # bigger). Small inline labels are unaffected.
        if _has_title_sized_text(unit):
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass

        # If unit has a background that we can render, paint it on the textbox.
        # Special case: when this looks like a `background-clip: text` recipe
        # (transparent text on a gradient bg), DON'T paint the gradient as
        # the shape's fill — that produces a gradient rectangle with the
        # text rendered as a transparent silhouette inside it. Instead,
        # leave the textbox transparent; the run-color resolver will
        # substitute the gradient's first stop as a solid text color.
        unit_anchor_el = unit.elements[0] if unit.elements else None
        if unit_anchor_el is not None and not _is_bg_clip_text(unit_anchor_el):
            self._apply_fill(tb, unit_anchor_el)
            self._apply_shadow(tb, unit_anchor_el)

        text_elems = [
            e
            for e in unit.all_elements()
            if (e.text and e.text.strip()) or (e.runs and any(r.text.strip() for r in e.runs if not r.is_break))
        ]
        # Leaf-text-wins: when a text-bearing parent has text-bearing
        # descendants in the same list, the descendants will paint the
        # words at their precise child bboxes — so emitting the parent
        # too produces a stacked duplicate (slide 8's `.author` div case
        # contained `.name` + `.title`, both with their own text).
        text_elem_ids = {e.id for e in text_elems}
        all_in_unit = unit.all_elements()
        parent_of = {e.id: e.parent_id for e in all_in_unit}

        def _has_text_descendant(eid: int) -> bool:
            for other in all_in_unit:
                if other.id == eid or other.id not in text_elem_ids:
                    continue
                cur = parent_of.get(other.id)
                while cur is not None:
                    if cur == eid:
                        return True
                    cur = parent_of.get(cur)
            return False

        text_elems = [e for e in text_elems if not _has_text_descendant(e.id)]
        para_targets = text_elems if len(text_elems) > 1 else [anchor]

        # Pre-compute fontScale against the fallback (Liberation Sans /
        # Calibri) rather than trusting MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        # at runtime. Bake the result into <a:normAutofit> so the deck
        # renders identically across PowerPoint and LibreOffice.
        try:
            scale_runs: list[tuple[str, float]] = []
            for e in para_targets:
                t = (e.pptx_text or e.text or "").strip()
                if not t:
                    continue
                scale_runs.append((t, max(8.0, parse_pt(e.font_size))))
            if scale_runs:
                scale = compute_font_scale_for_textbox(
                    scale_runs, bbox_w_px=emit_bbox.w, bbox_h_px=emit_bbox.h
                )
                if scale.needs_autofit:
                    _apply_explicit_autofit(
                        tf, scale.font_scale_pct, scale.line_space_reduction_pct
                    )
        except Exception:
            pass
        first = True
        for e in para_targets:
            # Build paragraph(s) from runs if available, else from text.
            paragraphs = self._element_to_paragraphs(e)
            for para_runs in paragraphs:
                if not para_runs:
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = _TEXT_ALIGN_MAP.get(e.text_align, PP_ALIGN.LEFT)
                for run_spec in para_runs:
                    self._add_styled_run(p, run_spec, fallback_el=e)

        # Decoration layer (ABOVE): rim highlight / hairline / inset glow.
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, emit_bbox)

    def _element_to_paragraphs(
        self, e: DomElement
    ) -> list[list[dict]]:
        """Turn an element's runs (or text) into a list of paragraphs, where
        each paragraph is a list of run-spec dicts. <br> splits paragraphs.
        """
        if e.runs:
            paragraphs: list[list[dict]] = [[]]
            for r in e.runs:
                if r.is_break:
                    paragraphs.append([])
                    continue
                spec = {
                    "text": r.text,
                    "font_family": r.font_family,
                    "font_size": r.font_size,
                    "font_weight": r.font_weight,
                    "color": r.color,
                    # Forward the run's parent background-image so the emitter
                    # can substitute a solid color when the run is gradient-
                    # clipped (color: transparent + background-clip: text).
                    "background_image": r.background_image,
                    "italic": r.italic,
                    "underline": r.underline,
                }
                if not r.text.strip():
                    # Preserve inter-run whitespace; skip empties.
                    if r.text and paragraphs[-1]:
                        paragraphs[-1].append(spec)
                    continue
                paragraphs[-1].append(spec)
            # Drop trailing empty paragraphs
            while paragraphs and not paragraphs[-1]:
                paragraphs.pop()
            return paragraphs
        text = (e.pptx_text or e.text or "").strip()
        if not text:
            return []
        return [
            [
                {
                    "text": text,
                    "font_family": e.font_family,
                    "font_size": e.font_size,
                    "font_weight": e.font_weight,
                    "color": e.color,
                    "italic": False,
                    "underline": False,
                }
            ]
        ]

    def _add_styled_run(self, paragraph, spec: dict, fallback_el: DomElement) -> None:
        run = paragraph.add_run()
        run.text = spec["text"]
        font = run.font
        font.name = resolve_font(spec.get("font_family") or fallback_el.font_family)
        size_px = spec.get("font_size") or fallback_el.font_size
        font.size = Pt(max(8.0, parse_pt(size_px)))
        font.bold = is_bold(spec.get("font_weight") or fallback_el.font_weight)
        font.italic = bool(spec.get("italic"))
        font.underline = bool(spec.get("underline"))
        # If this is a `background-clip: text` run with a translatable
        # gradient on the source span, emit a true OOXML gradient text fill
        # via python-pptx's Font.fill API. PPTX renders the gradient through
        # the glyph silhouettes — the visual the source CSS actually wants.
        if _try_apply_gradient_text_fill(font, spec):
            return
        color = _resolve_run_color(spec, fallback_el)
        if color is not None:
            try:
                font.color.rgb = color
            except Exception:
                pass

    # ----- native bullets ------------------------------------------------

    def _emit_native_bullet(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        # Treat a single ListItem unit as a one-paragraph bullet.
        x, y, w, h = _emu_rect(op.bbox)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        # PowerPoint substitutes Inter/etc with Calibri when the source font
        # isn't installed — Calibri is wider, so big headlines that fit in 3
        # lines in Chromium can overflow into the lede below in PowerPoint.
        # TEXT_TO_FIT_SHAPE shrinks the font just enough to keep everything
        # in-frame. We only apply it for *title-sized* text where the wrap
        # risk is real; on small frames (pills, badges, kickers) LibreOffice
        # over-shrinks the text to invisibility.
        # Heuristic: apply when the largest run is >= 20px (subtitle-ish or
        # bigger). Small inline labels are unaffected.
        if _has_title_sized_text(unit):
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        text_elems = [e for e in unit.all_elements() if e.text and e.text.strip()]
        if not text_elems:
            return
        first = True
        for e in text_elems:
            text = (e.pptx_text or e.text or "").strip()
            if not text:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ""
            run = p.add_run()
            run.text = "• " + text
            font = run.font
            font.name = resolve_font(e.font_family)
            font.size = Pt(max(8.0, parse_pt(e.font_size)))
            font.bold = is_bold(e.font_weight)
            color = parse_color(e.color)
            if color is not None:
                font.color.rgb = color[0]

    # ----- native shape --------------------------------------------------

    def _try_emit_native_decoration(
        self, slide, unit: VisualUnit, op: EmitOp
    ) -> bool:
        """If the unit's anchor decoration (gradient/shadow/solid) is all
        translatable, emit a native shape covering the bbox. Returns True if
        a native shape was emitted; False to fall back to raster.
        """
        anchor = unit.elements[0] if unit.elements else None
        if anchor is None:
            return False
        # Gradients we can translate — go native.
        bg_img = anchor.background_image
        if bg_img and bg_img != "none":
            if "url(" in bg_img:
                return False
            grad = parse_gradient(bg_img)
            if grad is None:
                return False
        # Decorative overlays may sit off-canvas (top:-220px etc); preserve
        # the original origin so the gradient center isn't displaced.
        x, y, w, h = _emu_rect_offcanvas(op.bbox)
        radius = parse_px(anchor.border_radius)
        shape_kind = _shape_kind_for_anchor(anchor, radius)
        deco_stack = derive_decorations(anchor, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, op.bbox)
        shape = slide.shapes.add_shape(shape_kind, x, y, w, h)
        shape.line.fill.background()
        self._apply_fill(shape, anchor)
        self._apply_border(shape, anchor)
        self._apply_shadow(shape, anchor)
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, op.bbox)
        return True

    def _emit_native_shape(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        anchor_el = unit.elements[0] if unit.elements else None
        if anchor_el is None:
            return
        # Decorative overlays may sit off-canvas — see _emit_native_decoration.
        x, y, w, h = _emu_rect_offcanvas(op.bbox)
        radius = parse_px(anchor_el.border_radius)
        shape_kind = _shape_kind_for_anchor(anchor_el, radius)
        deco_stack = derive_decorations(anchor_el, palette=self._brand_palette)
        if not deco_stack.is_empty():
            deco_stack.emit_below(slide, op.bbox)
        shape = slide.shapes.add_shape(shape_kind, x, y, w, h)
        shape.line.fill.background()  # default to no border
        self._apply_fill(shape, anchor_el)
        self._apply_border(shape, anchor_el)
        self._apply_shadow(shape, anchor_el)
        if not deco_stack.is_empty():
            deco_stack.emit_above(slide, op.bbox)

    def _apply_fill(self, shape, el: DomElement) -> None:
        # Try native gradient first (background-image: linear-gradient(...) / radial-gradient(...)).
        if el.background_image and el.background_image != "none":
            grad = parse_gradient(el.background_image)
            if grad is not None and apply_gradient_fill(shape, grad):
                return
            # Fall through to background_color / no-fill if gradient unparseable.
        bg = parse_color(el.background_color)
        if bg is None:
            try:
                shape.fill.background()
            except Exception:
                pass
            return
        rgb_color, alpha = bg
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb_color
        except Exception:
            return
        # Honor alpha when < 1: cards with `rgba(255,255,255,0.04)` were
        # rendering as solid WHITE and burning their light-gray text into
        # invisibility. python-pptx doesn't expose alpha on solid fills, so
        # drop down to OOXML and add an `<a:alpha val="..."/>` child.
        if alpha < 0.999:
            self._set_solid_fill_alpha(shape, alpha)

    def _set_solid_fill_alpha(self, shape, alpha: float) -> None:
        from lxml import etree

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        try:
            sp_pr = shape._element.spPr
        except AttributeError:
            return
        if sp_pr is None:
            return
        solid = sp_pr.find(f"{{{ns_a}}}solidFill")
        if solid is None:
            return
        srgb = solid.find(f"{{{ns_a}}}srgbClr")
        if srgb is None:
            return
        # Clear any pre-existing alpha child.
        for existing in srgb.findall(f"{{{ns_a}}}alpha"):
            srgb.remove(existing)
        etree.SubElement(
            srgb,
            f"{{{ns_a}}}alpha",
            attrib={"val": str(int(round(max(0.0, min(1.0, alpha)) * 100_000)))},
        )

    def _apply_shadow(self, shape, el: DomElement) -> None:
        if not el.box_shadow or el.box_shadow == "none":
            return
        layers = parse_box_shadows(el.box_shadow)
        if layers:
            apply_shadows(shape, layers)

    def _apply_border(self, shape, el: DomElement) -> None:
        if not el.border or el.border == "none":
            try:
                shape.line.fill.background()
            except Exception:
                pass
            return
        # "1px solid rgb(0, 0, 0)"
        parts = el.border.split()
        if not parts:
            return
        width = parse_px(parts[0])
        if width <= 0:
            try:
                shape.line.fill.background()
            except Exception:
                pass
            return
        # Try to find a color in the border string
        m = re.search(r"(rgba?\([^)]+\)|#[0-9a-fA-F]{3,8})", el.border)
        col = parse_color(m.group(1)) if m else None
        try:
            shape.line.width = Emu(px_to_emu(width))
            if col is not None:
                shape.line.color.rgb = col[0]
        except Exception:
            pass

    # ----- native picture (img) ------------------------------------------

    def _emit_native_svg(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        svg_el = next((e for e in unit.all_elements() if e.is_svg and e.svg_shapes), None)
        if svg_el is None or not svg_el.svg_shapes:
            return
        emit_svg_shapes(slide, svg_el.svg_shapes)

    def _emit_native_picture(self, slide, unit: VisualUnit, op: EmitOp) -> None:
        src = op.decision.metadata.get("src") or self._first_img_src(unit)
        if not src:
            return
        try:
            data = self._fetch_image(src)
        except Exception as e:
            log.warning("emitter.image_fetch_failed", src=src[:120], error=str(e))
            return
        x, y, w, h = _emu_rect(op.bbox)
        slide.shapes.add_picture(io.BytesIO(data), x, y, w, h)

    def _first_img_src(self, unit: VisualUnit) -> str | None:
        for e in unit.all_elements():
            if e.is_img and e.img_src:
                return e.img_src
        return None

    def _fetch_image(self, src: str) -> bytes:
        if src in self._image_cache:
            return self._image_cache[src]
        if src.startswith("data:"):
            import base64

            head, _, b64 = src.partition(",")
            if ";base64" in head:
                data = base64.b64decode(b64)
            else:
                data = b64.encode("utf-8")
            self._image_cache[src] = data
            return data
        parsed = urlparse(src)
        if parsed.scheme in ("http", "https"):
            r = self._http.get(src)
            r.raise_for_status()
            data = r.content
        elif parsed.scheme in ("file", ""):
            data = Path(parsed.path or src).read_bytes()
        else:
            raise EmitError(f"unsupported image scheme: {parsed.scheme}")
        self._image_cache[src] = data
        return data

    # ----- raster --------------------------------------------------------

    async def _emit_raster(
        self,
        slide,
        unit: VisualUnit,
        op: EmitOp,
        rendered: RenderedSlide,
        renderer,
    ) -> None:
        # Try a region screenshot from the live page; fall back to ground-truth crop.
        png: bytes | None = None
        try:
            png = await renderer.screenshot_region(
                rendered.html,
                "",
                (op.bbox.x, op.bbox.y, op.bbox.w, op.bbox.h),
            )
        except Exception as e:
            log.warning("emitter.region_screenshot_failed", error=str(e))
            png = None
        if not png:
            self._emit_region_raster(slide, rendered.ground_truth_png, op.bbox)
            return
        x, y, w, h = _emu_rect(op.bbox)
        slide.shapes.add_picture(io.BytesIO(png), x, y, w, h)

    def _emit_region_raster(
        self, slide, ground_truth_png: bytes, bbox: BoundingBox
    ) -> None:
        from PIL import Image

        gt = Image.open(io.BytesIO(ground_truth_png))
        bb = _clamp_bbox(bbox)
        crop = gt.crop(
            (
                int(bb.x),
                int(bb.y),
                int(bb.x + bb.w),
                int(bb.y + bb.h),
            )
        )
        out = io.BytesIO()
        crop.save(out, format="PNG")
        out.seek(0)
        x, y, w, h = _emu_rect(bb)
        slide.shapes.add_picture(out, x, y, w, h)

    def _emit_full_raster(self, slide, png_bytes: bytes) -> None:
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            Emu(0),
            Emu(0),
            Emu(SLIDE_W_EMU),
            Emu(SLIDE_H_EMU),
        )

    def _set_notes(self, slide, text: str) -> None:
        try:
            slide.notes_slide.notes_text_frame.text = text
        except Exception as e:
            log.warning("emitter.notes_failed", error=str(e))


def _shape_kind_for_anchor(anchor: DomElement, radius_px: float):
    """Pick the right MSO_SHAPE for an anchor.

    Priority: explicit clip-path / class hint → preset (chevron, arrow,
    star, hexagon, callout, etc.) via `slidify.preset_shapes`. Then fall
    back to OVAL for true circles, ROUNDED_RECTANGLE for radius > 0,
    RECTANGLE otherwise.
    """
    match = detect_preset_shape(anchor)
    if match is not None:
        return match.preset
    w = anchor.bbox.w
    h = anchor.bbox.h
    if w > 0 and h > 0:
        if radius_px >= min(w, h) * 0.45:
            return MSO_SHAPE.OVAL
    return MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE


def _has_title_sized_text(unit: VisualUnit) -> bool:
    """True iff the unit has at least one text element/run rendered at >= 20px.

    The cutoff matches CSS `text-xl` / `text-2xl` (display headings). Below it
    the text is body / labels / captions where wrap-on-overflow is acceptable
    and auto-fit causes more harm than good (LibreOffice's TEXT_TO_FIT_SHAPE
    over-shrinks small frames).
    """
    for e in unit.all_elements():
        if e.runs:
            for r in e.runs:
                if parse_px(r.font_size) >= 20.0:
                    return True
        elif e.text and e.text.strip():
            if parse_px(e.font_size) >= 20.0:
                return True
    return False


def _is_bg_clip_text(anchor: DomElement) -> bool:
    """Heuristic detection of the CSS `background-clip: text` recipe.

    Pattern:
      .x { background: linear-gradient(...); -webkit-background-clip: text;
           color: transparent; }

    The browser computes:
      anchor.color           = "rgba(0, 0, 0, 0)"
      anchor.background_image = "linear-gradient(...)"

    We can't see the `background-clip` property in the computed-style snapshot
    we collect, but the (transparent color + gradient bg-image) combination is
    a sufficiently distinctive marker. PPTX cannot natively reproduce
    gradient-clipped text on a block-level shape, so we paint the run with the
    gradient's first stop color and leave the textbox unfilled.
    """
    if not anchor.background_image or anchor.background_image == "none":
        return False
    if "gradient(" not in anchor.background_image:
        return False
    # color: transparent? In computed form it's `rgba(*, *, *, 0)`.
    col_str = (anchor.color or "").strip().lower()
    import re as _re

    return _re.match(r"rgba?\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*0(?:\.0+)?\s*\)", col_str) is not None


def _try_apply_gradient_text_fill(font, spec: dict) -> bool:
    """Native gradient fill on a text run via OOXML's `<a:gradFill>` inside
    `<a:rPr>`. Returns True iff a gradient was applied.

    Triggers only when the spec has a transparent text color AND a
    parseable gradient bg-image — i.e. the CSS `background-clip: text`
    pattern. python-pptx exposes `font.fill.gradient()` for this; we then
    populate the stops directly. The result is what the CSS author actually
    wanted: glyph-shaped gradient, not a gradient rectangle.
    """
    raw_color = spec.get("color") or ""
    if parse_color(raw_color) is not None:
        return False  # color is a real RGB, not transparent
    bg_image = spec.get("background_image") or "none"
    if not bg_image or bg_image == "none":
        return False
    grad = parse_gradient(bg_image)
    if grad is None or len(grad.stops) < 2:
        return False
    try:
        font.fill.gradient()
        gradient_stops = font.fill.gradient_stops
    except Exception:
        return False
    # python-pptx pre-creates 2 stops; we may have more or fewer.
    # Easiest robust path: drop down to OOXML and rewrite the gsLst.
    try:
        from lxml import etree as _et

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        # font.fill._xPr is the rPr element where the fill lives.
        rpr = font._rPr  # noqa: SLF001
        if rpr is None:
            return False
        grad_fill = rpr.find(f"{{{ns_a}}}gradFill")
        if grad_fill is None:
            return False
        gs_lst = grad_fill.find(f"{{{ns_a}}}gsLst")
        if gs_lst is None:
            return False
        # Wipe existing stops and rebuild.
        for child in list(gs_lst):
            gs_lst.remove(child)
        for stop in grad.stops:
            gs = _et.SubElement(
                gs_lst,
                f"{{{ns_a}}}gs",
                attrib={"pos": str(int(round(max(0.0, min(1.0, stop.position)) * 100_000)))},
            )
            srgb = _et.SubElement(
                gs, f"{{{ns_a}}}srgbClr", attrib={"val": stop.color_hex}
            )
            if stop.alpha < 0.999:
                _et.SubElement(
                    srgb,
                    f"{{{ns_a}}}alpha",
                    attrib={"val": str(int(round(stop.alpha * 100_000)))},
                )
        # Set direction: linear → <a:lin ang="..."/>, radial → <a:path path="circle">.
        # python-pptx may have already inserted a <a:lin> default; replace it.
        for tag in ("lin", "path"):
            for el in grad_fill.findall(f"{{{ns_a}}}{tag}"):
                grad_fill.remove(el)
        from slidify.gradients import LinearGradient, RadialGradient

        if isinstance(grad, LinearGradient):
            pptx_deg = (grad.angle_deg - 90.0) % 360.0
            _et.SubElement(
                grad_fill,
                f"{{{ns_a}}}lin",
                attrib={"ang": str(int(round(pptx_deg * 60_000))), "scaled": "0"},
            )
        elif isinstance(grad, RadialGradient):
            path = _et.SubElement(
                grad_fill, f"{{{ns_a}}}path", attrib={"path": "circle"}
            )
            cx = max(0.0, min(1.0, grad.cx))
            cy = max(0.0, min(1.0, grad.cy))
            _et.SubElement(
                path,
                f"{{{ns_a}}}fillToRect",
                attrib={
                    "l": str(int(round(cx * 100_000))),
                    "t": str(int(round(cy * 100_000))),
                    "r": str(int(round((1.0 - cx) * 100_000))),
                    "b": str(int(round((1.0 - cy) * 100_000))),
                },
            )
        return True
    except Exception:
        return False


def _resolve_run_color(spec: dict, fallback_el: DomElement):
    """Pick the RGB color for a styled text run.

    Special case: when the run's color is fully transparent (rgba(*,*,*,0))
    AND the run's parent has a background-image gradient, this is the
    `background-clip: text` pattern — text is filled with the gradient. PPTX
    can't reproduce gradient-clipped text natively, so we substitute the
    gradient's first (focal) stop color as a solid fallback. This matches
    most users' expectation: the dominant accent color shows through.
    """
    raw_color = spec.get("color") or fallback_el.color
    parsed = parse_color(raw_color)
    if parsed is not None:
        return parsed[0]
    # parse_color returned None — color is transparent or unrecognized.
    # Try the run's own bg-image first, then fall back to the anchor's
    # bg-image (covers leaf-text elements where the gradient is on the
    # element itself, not on a wrapping span).
    bg_image = (
        spec.get("background_image")
        or fallback_el.background_image
        or "none"
    )
    if bg_image and bg_image != "none":
        from slidify.gradients import parse_gradient

        grad = parse_gradient(bg_image)
        if grad is not None and grad.stops:
            # Use the FIRST stop (focal/dominant color in radial; start of
            # the linear). For a 3-stop gradient like `from-indigo-500 via-
            # purple-500 to-pink-500`, we get the indigo — visually the
            # closest single-color match for "this text should look colorful".
            stop = grad.stops[0]
            from pptx.dml.color import RGBColor

            try:
                return RGBColor(
                    int(stop.color_hex[0:2], 16),
                    int(stop.color_hex[2:4], 16),
                    int(stop.color_hex[4:6], 16),
                )
            except (ValueError, IndexError):
                pass
    return None


def native_area_ratio(
    ops: list[EmitOp], slide_w: int = SLIDE_W_PX, slide_h: int = SLIDE_H_PX
) -> float:
    """Compute fraction of slide area covered by native (non-raster) ops.

    Naively sums non-overlapping native bbox area. Approximate but bounded.
    """
    total = float(slide_w * slide_h)
    if total <= 0:
        return 0.0
    native_area = 0.0
    for op in ops:
        if op.decision.kind in (
            DecisionKind.NativeText,
            DecisionKind.NativeShape,
            DecisionKind.NativeBullet,
            DecisionKind.NativePicture,
            DecisionKind.NativeSvg,
            DecisionKind.Hybrid,  # hybrid keeps text editable; fill counts as native shape
        ):
            bb = _clamp_bbox(op.bbox)
            native_area += bb.w * bb.h
    return min(1.0, native_area / total)
