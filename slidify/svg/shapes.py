"""Translate simple SVG primitives to native PPTX shapes.

Supports `<rect>`, `<circle>`, `<ellipse>`, `<line>`, `<polygon>`, `<polyline>`,
`<path>` (via slidify.svg_path), and `<text>`.

Coordinates are in SVG userspace; we map them into the SVG element's viewport
bounding box on the slide, honoring `viewBox` and `preserveAspectRatio`.
"""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

from slidify.colors import parse_color
from slidify.geom import px_to_emu


@dataclass
class _Mapping:
    """Linear map from SVG userspace into a sub-rect of the slide."""

    svg_origin_x: float
    svg_origin_y: float
    scale_x: float
    scale_y: float

    def x(self, ux: float) -> float:
        return self.svg_origin_x + ux * self.scale_x

    def y(self, uy: float) -> float:
        return self.svg_origin_y + uy * self.scale_y


def _build_mapping(svg_rect: dict) -> _Mapping:
    """Map SVG userspace into a sub-rect of the slide.

    If the SVG has no `viewBox`, userspace == viewport pixels (1:1).
    Otherwise the viewBox `(min-x, min-y, vb-w, vb-h)` is mapped into the
    rendered viewport `(x, y, w, h)`, honoring `preserveAspectRatio` (the
    default `xMidYMid meet` keeps aspect and centers; `none` stretches).
    """
    rect_x = svg_rect.get("x", 0.0) or 0.0
    rect_y = svg_rect.get("y", 0.0) or 0.0
    rect_w = svg_rect.get("w", 0.0) or 0.0
    rect_h = svg_rect.get("h", 0.0) or 0.0

    vb = svg_rect.get("viewBox")
    if not vb or rect_w <= 0 or rect_h <= 0:
        return _Mapping(rect_x, rect_y, 1.0, 1.0)

    try:
        min_x, min_y, vb_w, vb_h = (float(v) for v in vb)
    except (TypeError, ValueError):
        return _Mapping(rect_x, rect_y, 1.0, 1.0)
    if vb_w <= 0 or vb_h <= 0:
        return _Mapping(rect_x, rect_y, 1.0, 1.0)

    par = (svg_rect.get("preserveAspectRatio") or "xMidYMid meet").strip().lower()
    parts = par.split()
    align = parts[0] if parts else "xmidymid"
    meet_or_slice = parts[1] if len(parts) > 1 else "meet"

    if align == "none":
        scale_x = rect_w / vb_w
        scale_y = rect_h / vb_h
        return _Mapping(
            rect_x - min_x * scale_x,
            rect_y - min_y * scale_y,
            scale_x,
            scale_y,
        )

    sx = rect_w / vb_w
    sy = rect_h / vb_h
    scale = min(sx, sy) if meet_or_slice != "slice" else max(sx, sy)
    extra_x = rect_w - vb_w * scale
    extra_y = rect_h - vb_h * scale

    if "xmin" in align:
        ox = 0.0
    elif "xmax" in align:
        ox = extra_x
    else:
        ox = extra_x / 2.0
    if "ymin" in align:
        oy = 0.0
    elif "ymax" in align:
        oy = extra_y
    else:
        oy = extra_y / 2.0

    return _Mapping(
        rect_x + ox - min_x * scale,
        rect_y + oy - min_y * scale,
        scale,
        scale,
    )


def is_translatable_svg(svg_shapes: list[dict] | None) -> bool:
    if not svg_shapes:
        return False
    for s in svg_shapes:
        tag = s.get("tag")
        if tag in ("rect", "circle", "ellipse", "line", "polygon", "polyline", "text"):
            continue
        if tag == "path":
            # Try a real parse to be sure svg_path can handle it. A bad
            # path string falls through to raster, which is the safer
            # default than emitting a broken custGeom.
            d = (s.get("d") or "").strip()
            if not d:
                return False
            try:
                from slidify.svg.path_parser import parse_path

                cmds = parse_path(d)
                if cmds:
                    continue
            except Exception:
                pass
            return False
        return False
    return True


def emit_svg_shapes(slide, svg_shapes: list[dict]) -> int:
    """Emit each translatable SVG primitive as a native PPTX shape on `slide`.

    Returns the number of shapes successfully emitted.
    """
    if not svg_shapes:
        return 0
    rect_meta = getattr(svg_shapes, "_svgRect", None)
    if rect_meta is None and isinstance(svg_shapes, dict):
        rect_meta = svg_shapes.get("_svgRect")
    # When deserialized through pydantic, the special `_svgRect` attribute
    # doesn't survive on the list. Look for it on the first dict instead.
    if rect_meta is None and svg_shapes and isinstance(svg_shapes[0], dict):
        rect_meta = svg_shapes[0].get("_svgRect")
    if rect_meta is None:
        return 0
    mapping = _build_mapping(rect_meta)

    n_emitted = 0
    for s in svg_shapes:
        # `_svgRect` used to ride as its own pseudo-entry in the list; now
        # the DOM walker embeds it as a key on every shape dict. Either way,
        # only entries with a real `tag` are emitted.
        if not isinstance(s, dict) or not s.get("tag"):
            continue
        try:
            if _emit_one(slide, s, mapping):
                n_emitted += 1
        except Exception:
            pass
    return n_emitted


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_solid_alpha(srgb_el, alpha: float) -> None:
    """Append (or replace) an `<a:alpha val="N"/>` child on an `<a:srgbClr>`.

    OOXML alpha is in 1/1000 of a percent (val=100000 → fully opaque).
    Without this, SVG `fill-opacity="0.55"` blobs render as fully-opaque
    walls of color, drowning the slide content underneath (slide-43 /
    slide-44 organic blobs).
    """
    from lxml import etree as _et

    if srgb_el is None:
        return
    for existing in srgb_el.findall(f"{{{_NS_A}}}alpha"):
        srgb_el.remove(existing)
    val = int(round(max(0.0, min(1.0, alpha)) * 100_000))
    _et.SubElement(srgb_el, f"{{{_NS_A}}}alpha", attrib={"val": str(val)})


def _apply_svg_fill(shape, fill, fill_opacity: float) -> None:
    """Apply a fill to an emitted PPTX shape.

    `fill` is one of:
      * ``None`` / ``"none"`` / ``"transparent"`` — clear fill.
      * ``str`` — a CSS colour (`"#1F2924"`, `"rgba(...)"`, …).  Emits
        `<a:solidFill>`.
      * ``dict`` — a structured gradient def the walker resolved from
        `fill="url(#id)"`.  Emits `<a:gradFill>` with the same stops,
        positions, and angle/radial geometry.

    Round-2 visual-QA traced the duotone-drift regression to the
    silent failure of `parse_color("url(#duo-bg)")` — every
    gradient-filled rect was rendering with no fill.  The walker now
    resolves url() refs inline, and this dispatch routes them through
    the native gradFill emit instead.
    """
    if fill is None or (isinstance(fill, str) and fill.lower() in ("none", "transparent")):
        try:
            shape.fill.background()
        except Exception:
            pass
        return
    if isinstance(fill, dict):
        _apply_svg_gradient_fill(shape, fill, fill_opacity)
        return
    parsed = parse_color(fill)
    if parsed is None:
        return
    rgb, color_alpha = parsed
    final_alpha = float(color_alpha) * float(max(0.0, min(1.0, fill_opacity)))
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    except Exception:
        return
    if final_alpha < 0.999:
        try:
            sp_pr = shape._element.spPr
            solid = sp_pr.find(f"{{{_NS_A}}}solidFill")
            if solid is not None:
                srgb = solid.find(f"{{{_NS_A}}}srgbClr")
                _set_solid_alpha(srgb, final_alpha)
        except Exception:
            pass


def _apply_svg_gradient_fill(shape, gradient: dict, fill_opacity: float) -> None:
    """Translate a resolved SVG gradient def into a native `<a:gradFill>`.

    The walker captures `<linearGradient>` / `<radialGradient>` defs
    from the SVG <defs> block, normalises stops to (offset 0..1, color,
    opacity), and resolves `fill="url(#id)"` references inline.  This
    function takes the resolved dict and rewrites the shape's spPr
    fill block.

    Linear gradients map directly: SVG (x1,y1)→(x2,y2) becomes an
    `<a:lin ang=...>` angle.  Radial gradients map to OOXML's
    "circle" path with `<a:tileRect>` controlling focal placement.
    Stops carry through with `<a:alpha>` per stop when the source
    `stop-opacity` is partial — the OOXML attribute is identical in
    semantics (1/100000 percent) so the round-trip is lossless.
    """
    import math

    from lxml import etree as _et

    stops = gradient.get("stops") or []
    if not stops:
        try:
            shape.fill.background()
        except Exception:
            pass
        return

    sp_pr = shape._element.spPr
    # Strip any existing fill children so we're authoring fresh.
    for tag in ("solidFill", "gradFill", "blipFill", "pattFill", "noFill"):
        for existing in sp_pr.findall(f"{{{_NS_A}}}{tag}"):
            sp_pr.remove(existing)

    # Determine where to insert the fill.  spPr's child ordering rules
    # (per OOXML schema) put fills after xfrm and prstGeom but before
    # ln, effectLst, etc.  Insert immediately after the geometry
    # element if present, else as the first child.
    insert_after = None
    for tag in ("custGeom", "prstGeom"):
        m = sp_pr.find(f"{{{_NS_A}}}{tag}")
        if m is not None:
            insert_after = m
            break

    grad_fill = _et.Element(f"{{{_NS_A}}}gradFill",
                             attrib={"flip": "none", "rotWithShape": "1"})
    if insert_after is not None:
        insert_after.addnext(grad_fill)
    else:
        sp_pr.insert(0, grad_fill)

    # <a:gsLst> with one <a:gs pos="..."> per stop
    gs_lst = _et.SubElement(grad_fill, f"{{{_NS_A}}}gsLst")
    for st in stops:
        offset = float(st.get("offset", 0))
        color  = st.get("color") or "#000000"
        op     = float(st.get("opacity", 1)) * float(max(0.0, min(1.0, fill_opacity)))
        parsed = parse_color(color)
        if parsed is None:
            continue
        rgb, color_alpha = parsed
        final_alpha = op * float(color_alpha)
        pos = int(round(max(0.0, min(1.0, offset)) * 100_000))
        gs = _et.SubElement(gs_lst, f"{{{_NS_A}}}gs", attrib={"pos": str(pos)})
        # `parse_color` returns python-pptx's `RGBColor`, whose `str()`
        # is the canonical 6-hex representation (`"0F1311"`).  That's
        # exactly what `<a:srgbClr val="...">` wants.
        srgb = _et.SubElement(gs, f"{{{_NS_A}}}srgbClr",
                              attrib={"val": str(rgb).upper()})
        if final_alpha < 0.999:
            _et.SubElement(srgb, f"{{{_NS_A}}}alpha",
                            attrib={"val": str(int(round(final_alpha * 100_000)))})

    # Geometry — linear angle or radial path
    if gradient.get("kind") == "linear":
        # SVG (x1,y1) → (x2,y2) defines the gradient direction in
        # objectBoundingBox units (0..1).  OOXML <a:lin ang=...> is in
        # 1/60000 of a degree, with 0° pointing right (matching SVG).
        x1 = float(gradient.get("x1", 0))
        y1 = float(gradient.get("y1", 0))
        x2 = float(gradient.get("x2", 1))
        y2 = float(gradient.get("y2", 0))
        dx, dy = x2 - x1, y2 - y1
        ang_deg = math.degrees(math.atan2(dy, dx)) if (dx or dy) else 0.0
        # Normalise into 0..360 and convert to OOXML's 60000ths.
        while ang_deg < 0:
            ang_deg += 360
        ooxml_ang = int(round(ang_deg * 60_000)) % (360 * 60_000)
        _et.SubElement(grad_fill, f"{{{_NS_A}}}lin",
                        attrib={"ang": str(ooxml_ang), "scaled": "1"})
    else:
        # Radial — emit a circle path with a focal tileRect derived
        # from `fx`/`fy` (objectBoundingBox 0..1).
        fx = float(gradient.get("fx", gradient.get("cx", 0.5)))
        fy = float(gradient.get("fy", gradient.get("cy", 0.5)))
        path_el = _et.SubElement(grad_fill, f"{{{_NS_A}}}path",
                                  attrib={"path": "circle"})
        _et.SubElement(path_el, f"{{{_NS_A}}}fillToRect",
                        attrib={
                            "l": str(int(round(fx * 100_000))),
                            "t": str(int(round(fy * 100_000))),
                            "r": str(int(round((1 - fx) * 100_000))),
                            "b": str(int(round((1 - fy) * 100_000))),
                        })


def _apply_svg_stroke(
    shape,
    stroke: str | None,
    stroke_width: float,
    stroke_opacity: float = 1.0,
) -> None:
    if stroke is None or stroke.lower() in ("none", "transparent") or stroke_width <= 0:
        try:
            shape.line.fill.background()
        except Exception:
            pass
        return
    parsed = parse_color(stroke)
    if parsed is None:
        return
    rgb, color_alpha = parsed
    final_alpha = float(color_alpha) * float(max(0.0, min(1.0, stroke_opacity)))
    try:
        shape.line.width = Emu(px_to_emu(stroke_width))
        shape.line.color.rgb = rgb
    except Exception:
        return
    if final_alpha < 0.999:
        try:
            sp_pr = shape._element.spPr
            ln = sp_pr.find(f"{{{_NS_A}}}ln")
            if ln is not None:
                solid = ln.find(f"{{{_NS_A}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{_NS_A}}}srgbClr")
                    _set_solid_alpha(srgb, final_alpha)
        except Exception:
            pass


def _emit_one(slide, s: dict, mapping: _Mapping) -> bool:
    tag = s.get("tag")
    fill = s.get("fill")
    stroke = s.get("stroke")
    stroke_width = float(s.get("stroke_width") or 0.0)
    fill_opacity = float(s.get("fill_opacity") if s.get("fill_opacity") is not None else 1.0)
    stroke_opacity = float(
        s.get("stroke_opacity") if s.get("stroke_opacity") is not None else 1.0
    )

    if tag == "rect":
        x = mapping.x(float(s["x"]))
        y = mapping.y(float(s["y"]))
        w = float(s["w"]) * mapping.scale_x
        h = float(s["h"]) * mapping.scale_y
        if w <= 0 or h <= 0:
            return False
        rx = float(s.get("rx") or 0.0)
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            kind, Emu(px_to_emu(x)), Emu(px_to_emu(y)), Emu(px_to_emu(w)), Emu(px_to_emu(h))
        )
        _apply_svg_fill(shape, fill, fill_opacity)
        _apply_svg_stroke(shape, stroke, stroke_width, stroke_opacity)
        return True

    if tag == "circle":
        cx = mapping.x(float(s["cx"]))
        cy = mapping.y(float(s["cy"]))
        r = float(s["r"])
        if r <= 0:
            return False
        x = cx - r
        y = cy - r
        w = h = 2 * r
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(px_to_emu(x)),
            Emu(px_to_emu(y)),
            Emu(px_to_emu(w)),
            Emu(px_to_emu(h)),
        )
        _apply_svg_fill(shape, fill, fill_opacity)
        _apply_svg_stroke(shape, stroke, stroke_width, stroke_opacity)
        return True

    if tag == "ellipse":
        cx = mapping.x(float(s["cx"]))
        cy = mapping.y(float(s["cy"]))
        rx = float(s["rx"])
        ry = float(s["ry"])
        if rx <= 0 or ry <= 0:
            return False
        x = cx - rx
        y = cy - ry
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(px_to_emu(x)),
            Emu(px_to_emu(y)),
            Emu(px_to_emu(2 * rx)),
            Emu(px_to_emu(2 * ry)),
        )
        _apply_svg_fill(shape, fill, fill_opacity)
        _apply_svg_stroke(shape, stroke, stroke_width, stroke_opacity)
        return True

    if tag == "line":
        x1 = mapping.x(float(s["x1"]))
        y1 = mapping.y(float(s["y1"]))
        x2 = mapping.x(float(s["x2"]))
        y2 = mapping.y(float(s["y2"]))
        connector = slide.shapes.add_connector(
            1,
            Emu(px_to_emu(x1)),
            Emu(px_to_emu(y1)),
            Emu(px_to_emu(x2)),
            Emu(px_to_emu(y2)),
        )
        _apply_svg_stroke(
            connector, stroke or "rgb(0,0,0)", stroke_width or 1.0, stroke_opacity
        )
        return True

    if tag == "path":
        d = (s.get("d") or "").strip()
        if not d:
            return False
        return _emit_path(
            slide, d, mapping, fill, fill_opacity, stroke, stroke_width, stroke_opacity
        )

    if tag == "text":
        return _emit_svg_text(slide, s, mapping)

    if tag in ("polygon", "polyline"):
        points = _parse_points(s.get("points", ""))
        if len(points) < 2:
            return False
        # Build a freeform shape via add_freeform.
        x0_user, y0_user = points[0]
        try:
            ff = slide.shapes.build_freeform(
                Emu(px_to_emu(mapping.x(x0_user))),
                Emu(px_to_emu(mapping.y(y0_user))),
                scale=1.0,
            )
            # IMPORTANT: pass ALL remaining points in ONE call.
            # python-pptx's `add_line_segments(pts, close=False)` writes one
            # contiguous run of `<a:lnTo>`s. Calling it once per point (the
            # previous code's bug) produced N independent `moveTo+lnTo+close`
            # micro-paths — every line chart rendered as a sequence of
            # zigzag triangles instead of a continuous polyline.
            segments = [
                (Emu(px_to_emu(mapping.x(ux))), Emu(px_to_emu(mapping.y(uy))))
                for (ux, uy) in points[1:]
            ]
            # `close=True` only when this is a polygon — polylines stay open.
            ff.add_line_segments(segments, close=(tag == "polygon"))
            shape = ff.convert_to_shape()
        except Exception:
            return False
        _apply_svg_fill(shape, fill, fill_opacity)
        _apply_svg_stroke(shape, stroke, stroke_width, stroke_opacity)
        return True

    return False


_POINT_PAIR = re.compile(r"(-?\d+(?:\.\d+)?)[\s,]+(-?\d+(?:\.\d+)?)")


def _parse_points(s: str) -> list[tuple[float, float]]:
    return [(float(x), float(y)) for x, y in _POINT_PAIR.findall(s)]


def _emit_svg_text(slide, s: dict, mapping: _Mapping) -> bool:
    """Emit an SVG `<text>` element as a native PPTX textbox.

    SVG text positioning is *baseline*-anchored at (x, y) — y is the
    BASELINE, not the top edge. PPTX textboxes are top-anchored, so we
    shift by approximately the font's ascent (≈0.8 × font-size) to
    align baselines.
    """
    text_str = (s.get("text") or "").strip()
    if not text_str:
        return False
    user_x = float(s.get("x") or 0.0)
    user_y = float(s.get("y") or 0.0)
    font_size_px = float(s.get("font_size") or 14.0)
    font_family = s.get("font_family") or "Inter"
    font_weight = str(s.get("font_weight") or "400")
    italic = (s.get("font_style") or "").lower() == "italic"
    fill = s.get("fill") or "#000"
    anchor = (s.get("text_anchor") or "start").lower()

    slide_x = mapping.x(user_x)
    slide_y_baseline = mapping.y(user_y)
    # Shift up so the baseline lands at slide_y_baseline.
    ascent = font_size_px * 0.8
    slide_y_top = slide_y_baseline - ascent

    # Estimate text width — Inter at this weight/size renders ~0.55 ×
    # font_size per character on average. Generous so the textbox isn't
    # too tight (PPT can autofit smaller; can't extend wider than bbox).
    est_width = max(8.0, len(text_str) * font_size_px * 0.55)
    height = font_size_px * 1.25

    if anchor == "middle":
        slide_x -= est_width / 2.0
    elif anchor == "end":
        slide_x -= est_width

    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    from slidify.colors import parse_color
    from slidify.fonts import is_bold
    from slidify.fonts import resolve as resolve_font

    try:
        tb = slide.shapes.add_textbox(
            Emu(px_to_emu(slide_x)),
            Emu(px_to_emu(slide_y_top)),
            Emu(px_to_emu(est_width)),
            Emu(px_to_emu(height)),
        )
    except Exception:
        return False
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    if anchor == "middle":
        p.alignment = PP_ALIGN.CENTER
    elif anchor == "end":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text_str
    font = run.font
    try:
        font.name = resolve_font(font_family)
        # SVG font_size is in user units that we treat as px; PPTX wants pt.
        font.size = Pt(max(6.0, font_size_px / 1.333))
        font.bold = is_bold(font_weight)
        font.italic = italic
        col = parse_color(fill)
        if col is not None:
            font.color.rgb = col[0]
    except Exception:
        pass
    return True


def _path_bounds(commands: list) -> tuple[float, float, float, float]:
    """Return (minx, miny, maxx, maxy) over all endpoint + control points."""
    xs: list[float] = []
    ys: list[float] = []
    for cmd in commands:
        op = cmd[0]
        if op == "moveTo" or op == "lineTo":
            xs.append(cmd[1])
            ys.append(cmd[2])
        elif op == "cubicBezTo":
            for i in (1, 3, 5):
                xs.append(cmd[i])
                ys.append(cmd[i + 1])
        elif op == "quadBezTo":
            for i in (1, 3):
                xs.append(cmd[i])
                ys.append(cmd[i + 1])
    if not xs:
        return (0.0, 0.0, 1.0, 1.0)
    return (min(xs), min(ys), max(xs), max(ys))


def _emit_path(
    slide,
    d: str,
    mapping: _Mapping,
    fill: str | None,
    fill_opacity: float,
    stroke: str | None,
    stroke_width: float,
    stroke_opacity: float = 1.0,
) -> bool:
    """Emit an SVG `<path d="...">` as a native PPTX freeform via custGeom.

    Uses `slidify.svg_path.parse_path` for full curve command coverage
    (cubic / smooth-cubic / quadratic / smooth-quadratic). Builds a
    custGeom XML fragment and grafts it onto a placeholder freeform
    shape, replacing the default prstGeom.
    """
    from lxml import etree

    from slidify.svg.path_parser import commands_to_path_xml, parse_path

    try:
        commands = parse_path(d)
    except Exception:
        return False
    if not commands:
        return False

    # Compute path bounds in SVG-local space, then map to slide-px.
    minx, miny, maxx, maxy = _path_bounds(commands)
    # Horizontal / vertical / dot paths have zero height or width. Pad
    # symmetrically by stroke_width so the stroke is visible *within* the
    # freeform's bbox (otherwise PPT clips half of it).
    pad = max(2.0, stroke_width * 1.5)
    if maxx == minx:
        minx -= pad
        maxx += pad
    elif maxx < minx:
        return False
    if maxy == miny:
        miny -= pad
        maxy += pad
    elif maxy < miny:
        return False

    slide_x = mapping.x(minx)
    slide_y = mapping.y(miny)
    slide_w = (maxx - minx) * mapping.scale_x
    slide_h = (maxy - miny) * mapping.scale_y
    if slide_w <= 0 or slide_h <= 0:
        return False

    # Translate commands so the path's local-origin is (0, 0) at the
    # bbox top-left — custGeom uses path-local coords scaled to the shape.
    PATH_W = 100_000
    PATH_H = 100_000
    norm: list = []
    for cmd in commands:
        op = cmd[0]
        if op == "moveTo" or op == "lineTo":
            x = (cmd[1] - minx) / (maxx - minx) * PATH_W
            y = (cmd[2] - miny) / (maxy - miny) * PATH_H
            norm.append((op, x, y))
        elif op == "cubicBezTo":
            pts = []
            for i in (1, 3, 5):
                x = (cmd[i] - minx) / (maxx - minx) * PATH_W
                y = (cmd[i + 1] - miny) / (maxy - miny) * PATH_H
                pts.extend([x, y])
            norm.append((op, *pts))
        elif op == "quadBezTo":
            pts = []
            for i in (1, 3):
                x = (cmd[i] - minx) / (maxx - minx) * PATH_W
                y = (cmd[i + 1] - miny) / (maxy - miny) * PATH_H
                pts.extend([x, y])
            norm.append((op, *pts))
        elif op == "close":
            norm.append(cmd)

    path_xml = commands_to_path_xml(norm, PATH_W, PATH_H)

    # Bootstrap a freeform so we get a valid <p:sp> with all the
    # plumbing (nvSpPr, spPr, etc.), then swap its <a:prstGeom> for
    # our <a:custGeom>.
    try:
        ff = slide.shapes.build_freeform(
            Emu(px_to_emu(slide_x)),
            Emu(px_to_emu(slide_y)),
            scale=1.0,
        )
        ff.add_line_segments(
            [(Emu(px_to_emu(slide_x + 1)), Emu(px_to_emu(slide_y + 1)))]
        )
        shape = ff.convert_to_shape()
    except Exception:
        return False

    # Resize to the actual path bbox.
    try:
        shape.width = Emu(px_to_emu(slide_w))
        shape.height = Emu(px_to_emu(slide_h))
    except Exception:
        pass

    # Swap prstGeom → custGeom. OOXML's spPr child order matters:
    # xfrm → (prstGeom | custGeom) → fill → ln → effects.
    # LibreOffice silently drops slides whose spPr children violate this.
    sp_pr = shape._element.spPr
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for prst in sp_pr.findall(f"{{{ns_a}}}prstGeom"):
        sp_pr.remove(prst)
    for cust in sp_pr.findall(f"{{{ns_a}}}custGeom"):
        sp_pr.remove(cust)
    custgeom_el = etree.fromstring(
        f'<a:custGeom xmlns:a="{ns_a}"><a:avLst/><a:gdLst/><a:ahLst/>'
        f"<a:cxnLst/><a:rect l=\"0\" t=\"0\" r=\"{PATH_W}\" b=\"{PATH_H}\"/>"
        f"<a:pathLst>{path_xml}</a:pathLst></a:custGeom>"
    )
    # Insert right after <a:xfrm> (or at index 0 if absent).
    xfrm = sp_pr.find(f"{{{ns_a}}}xfrm")
    if xfrm is not None:
        sp_pr.insert(list(sp_pr).index(xfrm) + 1, custgeom_el)
    else:
        sp_pr.insert(0, custgeom_el)

    _apply_svg_fill(shape, fill, fill_opacity)
    _apply_svg_stroke(shape, stroke, stroke_width, stroke_opacity)
    return True
